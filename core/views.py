from rest_framework import viewsets, status, permissions, serializers
from rest_framework.decorators import action
from rest_framework.response import Response
from rest_framework.views import APIView
from django.utils import timezone
from django.db.models import Q, Count, Sum
from django.db import models
from django.views.decorators.clickjacking import xframe_options_exempt
from django.views.decorators.csrf import csrf_exempt
from django.utils.decorators import method_decorator
from .models import (
    User, Municipio, Cliente, Video,
    Playlist, PlaylistItem, DispositivoTV, LogExibicao, Segmento, AppVersion,
    QRCodeClick, AgendamentoExibicao, ConteudoCorporativo, ConfiguracaoAPI,
    HorarioFuncionamento, LogExibicaoWebView,
    Campanha, CampanhaCupomConfig, CampanhaLead,
    CampanhaRoletaConfig, CampanhaRoletaPremio, CampanhaJogada,
    CampanhaCartaConfig,
    CampanhaAlertaConfig, CampanhaAlertaCampo, CampanhaAlertaLead,
    CampanhaSorteioConfig, CampanhaParticipanteSorteio,
    LandingLead,
    AgenteIA, AgenteIAConversa, AgenteIAMensagem,
)
from .serializers import (
    UserSerializer, UserMinimalSerializer, MunicipioSerializer,
    ClienteSerializer, ClienteCreateSerializer, VideoSerializer,
    PlaylistSerializer, PlaylistItemSerializer, DispositivoTVSerializer,
    LogExibicaoSerializer, LogExibicaoWebViewSerializer,
    PlaylistTVSerializer, DispositivoTVAuthSerializer
)
from .permissions import (
    IsOwner, IsFranchiseeOrOwner, IsClientOrAbove,
    IsOwnerOfObject, CanManageClients, CanManagePlaylists, CanManageVideos
)


class UserViewSet(viewsets.ModelViewSet):
    """
    ViewSet para gerenciar usuários
    - OWNER: vê todos
    - FRANCHISEE: vê seus clientes
    - CLIENT: vê apenas a si mesmo
    """
    queryset = User.objects.all()
    serializer_class = UserSerializer
    permission_classes = [permissions.IsAuthenticated]
    
    def get_queryset(self):
        user = self.request.user
        qs = User.objects.all()
        if user.is_owner():
            return qs
        elif user.is_franchisee():
            return qs.filter(
                Q(id=user.id) | Q(created_by=user) | Q(cliente_profile__franqueado=user)
            ).distinct()
        else:
            return qs.filter(id=user.id)
    
    @action(detail=False, methods=['get'])
    def me(self, request):
        """Retorna dados do usuário logado"""
        serializer = self.get_serializer(request.user)
        return Response(serializer.data)
    
    @action(detail=False, methods=['get'], permission_classes=[IsFranchiseeOrOwner])
    def franchisees(self, request):
        """Lista todos os franqueados (apenas para OWNER)"""
        franchisees = User.objects.filter(role='FRANCHISEE')
        serializer = UserMinimalSerializer(franchisees, many=True)
        return Response(serializer.data)


class MunicipioViewSet(viewsets.ModelViewSet):
    """
    ViewSet para gerenciar municípios
    Apenas franqueados e dono podem gerenciar
    """
    queryset = Municipio.objects.select_related('franqueado').all()
    serializer_class = MunicipioSerializer
    permission_classes = [IsFranchiseeOrOwner]
    
    def get_queryset(self):
        user = self.request.user
        qs = Municipio.objects.select_related('franqueado')
        if user.is_owner():
            return qs
        elif user.is_franchisee():
            return qs.filter(franqueado=user)
        return Municipio.objects.none()
    
    def perform_create(self, serializer):
        """Ao criar, define o franqueado automaticamente se não for owner"""
        if self.request.user.is_franchisee():
            serializer.save(franqueado=self.request.user)
        else:
            serializer.save()


class ClienteViewSet(viewsets.ModelViewSet):
    """
    ViewSet para gerenciar clientes
    """
    queryset = Cliente.objects.select_related('user', 'franqueado', 'segmento').all()
    permission_classes = [CanManageClients]
    
    def get_serializer_class(self):
        if self.action == 'create':
            return ClienteCreateSerializer
        return ClienteSerializer
    
    def get_queryset(self):
        user = self.request.user
        qs = Cliente.objects.select_related('user', 'franqueado', 'segmento').prefetch_related('municipios')
        if user.is_owner():
            return qs
        elif user.is_franchisee():
            return qs.filter(franqueado=user)
        return Cliente.objects.none()
    
    def perform_create(self, serializer):
        """Ao criar, define o franqueado automaticamente se não for owner"""
        if self.request.user.is_franchisee():
            serializer.save(franqueado=self.request.user)
        else:
            serializer.save()
    
    @action(detail=True, methods=['get'])
    def videos(self, request, pk=None):
        """Lista vídeos de um cliente específico"""
        cliente = self.get_object()
        videos = Video.objects.filter(cliente=cliente)
        serializer = VideoSerializer(videos, many=True, context={'request': request})
        return Response(serializer.data)


class VideoViewSet(viewsets.ModelViewSet):
    """
    ViewSet para gerenciar vídeos
    """
    queryset = Video.objects.select_related('cliente', 'cliente__user').all()
    serializer_class = VideoSerializer
    permission_classes = [CanManageVideos, IsOwnerOfObject]
    
    def get_queryset(self):
        user = self.request.user
        qs = Video.objects.select_related('cliente', 'cliente__user', 'cliente__segmento')
        if user.is_owner():
            return qs
        elif user.is_franchisee():
            return qs.filter(cliente__franqueado=user)
        elif user.is_client():
            try:
                cliente = user.cliente_profile
                return qs.filter(cliente=cliente)
            except Cliente.DoesNotExist:
                return Video.objects.none()
        return Video.objects.none()
    
    def perform_create(self, serializer):
        """Cliente cria vídeo para si mesmo"""
        user = self.request.user
        if user.is_client():
            try:
                cliente = user.cliente_profile
                serializer.save(cliente=cliente)
            except Cliente.DoesNotExist:
                raise serializers.ValidationError("Usuário não possui perfil de cliente")
        else:
            serializer.save()
    
    @action(detail=True, methods=['post'], permission_classes=[IsFranchiseeOrOwner])
    def approve(self, request, pk=None):
        """Aprova um vídeo"""
        video = self.get_object()
        video.status = 'APPROVED'
        video.save()
        return Response({'status': 'vídeo aprovado'})
    
    @action(detail=True, methods=['post'], permission_classes=[IsFranchiseeOrOwner])
    def reject(self, request, pk=None):
        """Rejeita um vídeo"""
        video = self.get_object()
        video.status = 'REJECTED'
        video.save()
        return Response({'status': 'vídeo rejeitado'})


class PlaylistViewSet(viewsets.ModelViewSet):
    """
    ViewSet para gerenciar playlists
    """
    queryset = Playlist.objects.select_related('municipio', 'franqueado').all()
    serializer_class = PlaylistSerializer
    permission_classes = [CanManagePlaylists, IsOwnerOfObject]
    
    def get_queryset(self):
        user = self.request.user
        qs = Playlist.objects.select_related('municipio', 'franqueado').prefetch_related(
            'items__video'
        )
        if user.is_owner():
            return qs
        elif user.is_franchisee():
            return qs.filter(franqueado=user)
        return Playlist.objects.none()
    
    def perform_create(self, serializer):
        """Ao criar, define o franqueado automaticamente se não for owner"""
        if self.request.user.is_franchisee():
            serializer.save(franqueado=self.request.user)
        else:
            serializer.save()
    
    @action(detail=True, methods=['post'])
    def add_video(self, request, pk=None):
        """Adiciona um vídeo à playlist"""
        playlist = self.get_object()
        video_id = request.data.get('video_id')
        ordem = request.data.get('ordem', 0)
        repeticoes = request.data.get('repeticoes', 1)
        
        try:
            video = Video.objects.get(id=video_id, status='APPROVED')
        except Video.DoesNotExist:
            return Response(
                {'error': 'Vídeo não encontrado ou não aprovado'},
                status=status.HTTP_404_NOT_FOUND
            )
        
        item, created = PlaylistItem.objects.get_or_create(
            playlist=playlist,
            video=video,
            defaults={'ordem': ordem, 'repeticoes': repeticoes}
        )
        
        if not created:
            item.ordem = ordem
            item.repeticoes = repeticoes
            item.save()
        
        playlist.calcular_duracao_total()
        
        serializer = PlaylistItemSerializer(item)
        return Response(serializer.data, status=status.HTTP_201_CREATED)
    
    @action(detail=True, methods=['delete'])
    def remove_video(self, request, pk=None):
        """Remove um vídeo da playlist"""
        playlist = self.get_object()
        video_id = request.data.get('video_id')
        
        try:
            item = PlaylistItem.objects.get(playlist=playlist, video_id=video_id)
            item.delete()
            playlist.calcular_duracao_total()
            return Response({'status': 'vídeo removido da playlist'})
        except PlaylistItem.DoesNotExist:
            return Response(
                {'error': 'Vídeo não encontrado na playlist'},
                status=status.HTTP_404_NOT_FOUND
            )
    
    @action(detail=True, methods=['post'])
    def reorder(self, request, pk=None):
        """Reordena itens da playlist"""
        playlist = self.get_object()
        items_order = request.data.get('items', [])  # Lista de {id, ordem}
        
        for item_data in items_order:
            try:
                item = PlaylistItem.objects.get(id=item_data['id'], playlist=playlist)
                item.ordem = item_data['ordem']
                item.save()
            except PlaylistItem.DoesNotExist:
                continue
        
        return Response({'status': 'playlist reordenada'})


class PlaylistItemViewSet(viewsets.ModelViewSet):
    """
    ViewSet para gerenciar itens de playlist
    """
    queryset = PlaylistItem.objects.select_related('playlist', 'video').all()
    serializer_class = PlaylistItemSerializer
    permission_classes = [CanManagePlaylists]
    
    def get_queryset(self):
        user = self.request.user
        qs = PlaylistItem.objects.select_related('playlist', 'video', 'video__cliente')
        if user.is_owner():
            return qs
        elif user.is_franchisee():
            return qs.filter(playlist__franqueado=user)
        return PlaylistItem.objects.none()


class DispositivoTVViewSet(viewsets.ModelViewSet):
    """
    ViewSet para gerenciar dispositivos de TV
    """
    queryset = DispositivoTV.objects.select_related('municipio', 'playlist_atual').all()
    serializer_class = DispositivoTVSerializer
    permission_classes = [IsFranchiseeOrOwner]
    lookup_field = 'identificador_unico'

    def get_queryset(self):
        user = self.request.user
        qs = DispositivoTV.objects.select_related('municipio', 'municipio__franqueado', 'playlist_atual')
        if user.is_owner():
            return qs
        elif user.is_franchisee():
            return qs.filter(municipio__franqueado=user)
        return DispositivoTV.objects.none()


class LogExibicaoViewSet(viewsets.ModelViewSet):
    """
    ViewSet para logs de exibição (somente leitura para usuários)
    """
    queryset = LogExibicao.objects.select_related('dispositivo', 'video', 'playlist').all()
    serializer_class = LogExibicaoSerializer
    permission_classes = [IsFranchiseeOrOwner]
    http_method_names = ['get', 'post']  # Apenas leitura e criação
    
    def get_queryset(self):
        user = self.request.user
        qs = LogExibicao.objects.select_related('dispositivo', 'video', 'playlist')
        if user.is_owner():
            return qs
        elif user.is_franchisee():
            return qs.filter(dispositivo__municipio__franqueado=user)
        return LogExibicao.objects.none()
    
    @action(detail=False, methods=['get'])
    def stats(self, request):
        """Estatísticas de exibição"""
        queryset = self.get_queryset()
        stats = {
            'total_exibicoes': queryset.count(),
            'exibicoes_completas': queryset.filter(completamente_exibido=True).count(),
            'videos_mais_exibidos': queryset.values('video__titulo').annotate(
                total=Count('id')
            ).order_by('-total')[:10]
        }
        return Response(stats)


# ─── Distribuição proporcional por tempo (Weighted Fair Queuing) ─────────────
def _distribuir_por_percentual(pairs, max_videos=200):
    """
    Weighted Fair Queuing baseado em duração real dos vídeos.

    Cada playlist acumula créditos de tempo proporcionais ao seu peso (percentual).
    A playlist com MAIS crédito acumulado ("mais devida") é a próxima a ser exibida.
    Ao exibir um vídeo, o crédito diminui pela duração real normalizada pela média
    global — vídeos mais longos "custam" mais crédito, garantindo equidade por tempo
    de tela e não apenas por contagem de vídeos.

    Ex: playlist A (80 %, vídeos 60 s) e B (20 %, vídeos 15 s):
        → A aparece ~4× para cada aparição de B (tempo de tela ≈ 80 %/20 %).
    """
    import itertools
    from math import lcm
    from functools import reduce

    valid = [(list(vids), pct) for vids, pct in pairs if vids and pct and pct > 0]
    if not valid:
        return [v for vids, _ in pairs for v in vids]
    if len(valid) == 1:
        return list(valid[0][0])

    total_pct = sum(pct for _, pct in valid)
    weights = [pct / total_pct for _, pct in valid]  # e.g. [0.8, 0.2]

    # Duração média global (todos os vídeos de todas as playlists)
    all_vids_flat = [v for vids, _ in valid for v in vids]
    all_durs = [
        v.get('duracao_segundos', 30)
        for v in all_vids_flat
        if isinstance(v, dict) and v.get('duracao_segundos', 0) > 0
    ]
    avg_dur = (sum(all_durs) / len(all_durs)) if all_durs else 30.0

    # Tamanho do ciclo: LCM dos tamanhos das playlists, escalado pelo menor peso,
    # limitado a max_videos para não gerar listas enormes
    lengths = [len(vids) for vids, _ in valid]
    try:
        cycle_len = reduce(lcm, lengths)
    except Exception:
        cycle_len = max(lengths)
    min_weight = min(weights)
    needed = int(max(lengths) / min_weight)  # garante ciclo completo de cada playlist
    total = min(max(needed, cycle_len), max_videos)

    iters = [itertools.cycle(vids) for vids, _ in valid]
    credits = [0.0] * len(valid)
    result = []

    for _ in range(total):
        # Cada playlist acumula crédito proporcional ao seu peso
        for i in range(len(valid)):
            credits[i] += weights[i]

        # Executa a playlist com maior crédito acumulado (a "mais devida")
        chosen = max(range(len(valid)), key=lambda i: credits[i])
        video = next(iters[chosen])

        # Desconta pela duração real normalizada:
        # vídeo de 60 s com média 30 s custa 2 créditos; de 15 s custa 0,5.
        dur = video.get('duracao_segundos', avg_dur) if isinstance(video, dict) else avg_dur
        if dur <= 0:
            dur = avg_dur
        credits[chosen] -= dur / avg_dur

        result.append(video)

    return result


# API específica para o App de TV
class TVAPIView(APIView):
    """
    API para o app de TV se autenticar e buscar playlist
    """
    permission_classes = [permissions.AllowAny]
    
    def post(self, request):
        """Autenticação de dispositivo e retorno da playlist"""
        serializer = DispositivoTVAuthSerializer(data=request.data)
        
        if not serializer.is_valid():
            return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)
        
        identificador = serializer.validated_data['identificador_unico']
        versao_app = serializer.validated_data.get('versao_app', '')
        
        try:
            dispositivo = DispositivoTV.objects.get(
                identificador_unico=identificador,
                ativo=True
            )
            
            # Atualiza última sincronização
            estava_offline = dispositivo.alerta_desconexao_enviado
            dispositivo.ultima_sincronizacao = timezone.now()
            if versao_app:
                dispositivo.versao_app = versao_app
            if estava_offline:
                dispositivo.alerta_desconexao_enviado = False
            dispositivo.save()

            # Notifica reconexão em background (sem bloquear resposta)
            if estava_offline:
                try:
                    from core.alerts import send_online_alert
                    import threading
                    threading.Thread(target=send_online_alert, args=(dispositivo,), daemon=True).start()
                except Exception:
                    pass
            
            # Retorna TODAS as playlists ativas no horário atual mescladas
            agendamentos_ativos = dispositivo.get_agendamentos_ativos_por_horario()

            if agendamentos_ativos:
                pairs = []  # (video_list, percentual)
                playlist_ids = []
                playlist_names = []

                for ag in agendamentos_ativos:
                    playlist = ag.playlist
                    if not playlist.ativa:
                        continue
                    playlist_ids.append(playlist.id)
                    playlist_names.append(playlist.nome)
                    pl_serializer = PlaylistTVSerializer(
                        playlist,
                        context={'request': request, 'dispositivo_id': dispositivo.id}
                    )
                    videos = pl_serializer.data.get('videos', [])
                    pairs.append((videos, ag.percentual))

                # Usar distribuição proporcional se houver percentuais variados
                has_varied = len(pairs) > 1 and any(pct != 100 for _, pct in pairs)
                if has_varied:
                    all_videos = _distribuir_por_percentual(pairs)
                else:
                    all_videos = [v for vids, _ in pairs for v in vids]

                return Response({
                    'dispositivo_id': dispositivo.id,
                    'dispositivo_nome': dispositivo.nome,
                    'municipio': str(dispositivo.municipio),
                    'playlist': {
                        'id': playlist_ids[0] if len(playlist_ids) == 1 else 0,
                        'nome': ' + '.join(playlist_names),
                        'duracao_total_segundos': sum(v.get('duracao_segundos', 0) for v in all_videos),
                        'videos': all_videos,
                        'playlists_mescladas': playlist_ids,
                    }
                })
            else:
                return Response({
                    'dispositivo_id': dispositivo.id,
                    'dispositivo_nome': dispositivo.nome,
                    'municipio': str(dispositivo.municipio),
                    'playlist': None,
                    'message': 'Nenhuma playlist ativa configurada'
                })
        
        except DispositivoTV.DoesNotExist:
            return Response(
                {'error': 'Dispositivo não encontrado ou inativo'},
                status=status.HTTP_404_NOT_FOUND
            )
        except Exception as e:
            import traceback as _tb
            import logging as _logging
            _logging.getLogger(__name__).error(
                "TVAPIView 500: %s\n%s", e, _tb.format_exc()
            )
            return Response(
                {'error': f'Erro interno: {type(e).__name__}: {e}'},
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )


class TVLogExibicaoView(APIView):
    """
    API para o app de TV registrar logs de exibição
    Formato esperado: {dispositivo_id, video_id, tempo_exibicao_segundos}
    """
    permission_classes = [permissions.AllowAny]

    def post(self, request):
        """Registra log de exibição e cria entrada parcial para o próximo item da fila"""
        dispositivo_id = request.data.get('dispositivo_id')
        video_id = request.data.get('video_id')
        tempo_exibicao_segundos = request.data.get('tempo_exibicao_segundos', 0)
        playlist_id = request.data.get('playlist_id')  # opcional — enviado pelo app

        if not dispositivo_id or not video_id:
            return Response(
                {'error': 'dispositivo_id e video_id são obrigatórios'},
                status=status.HTTP_400_BAD_REQUEST
            )

        try:
            dispositivo = DispositivoTV.objects.get(id=dispositivo_id)
            video = Video.objects.get(id=video_id)

            # Resolve a playlist do log:
            # 1. playlist_id enviado pelo app (mais preciso)
            # 2. playlist ativa no horário atual via agendamento
            # 3. playlist_atual como último fallback
            playlist = None
            if playlist_id:
                try:
                    playlist = Playlist.objects.get(id=playlist_id)
                except Exception:
                    pass
            if playlist is None:
                playlist = dispositivo.get_playlist_atual_por_horario()
            if playlist is None:
                playlist = dispositivo.playlist_atual

            data_hora_fim = timezone.now()
            data_hora_inicio = data_hora_fim - timezone.timedelta(seconds=int(tempo_exibicao_segundos))

            # Verifica se foi completamente exibido (≥90% do tempo ou duração desconhecida)
            if video.duracao_segundos and video.duracao_segundos > 0:
                porcentagem = (int(tempo_exibicao_segundos) / video.duracao_segundos) * 100
                completamente_exibido = porcentagem >= 90
            else:
                # Duração não cadastrada → assume completo
                completamente_exibido = True

            # Cria registro do log com início/fim corretos
            LogExibicao.objects.create(
                dispositivo=dispositivo,
                video=video,
                playlist=playlist,
                data_hora_inicio=data_hora_inicio,
                data_hora_fim=data_hora_fim,
                completamente_exibido=completamente_exibido
            )

            return Response(
                {'success': True, 'message': 'Log registrado com sucesso'},
                status=status.HTTP_201_CREATED
            )

        except DispositivoTV.DoesNotExist:
            return Response({'error': 'Dispositivo não encontrado'}, status=status.HTTP_404_NOT_FOUND)
        except Video.DoesNotExist:
            return Response({'error': 'Vídeo não encontrado'}, status=status.HTTP_404_NOT_FOUND)


def _criar_log_parcial_proximo(*args, **kwargs):
    """Removido — causava registros phantom de 'Em reprodução' para vídeos fora da playlist."""
    pass


class TVLogWebViewView(APIView):
    """
    API para o app de TV registrar execuções de conteúdo corporativo (WebView).
    Formato: {dispositivo_id, conteudo_corporativo_id, tipo_conteudo, titulo,
               duracao_segundos, completamente_exibido}
    """
    permission_classes = [permissions.AllowAny]

    def post(self, request):
        dispositivo_id = request.data.get('dispositivo_id')
        tipo_conteudo = request.data.get('tipo_conteudo', 'DESIGN')
        titulo = request.data.get('titulo', 'Conteúdo Corporativo')
        duracao_segundos = int(request.data.get('duracao_segundos', 0))
        corporativo_id = request.data.get('conteudo_corporativo_id')

        if not dispositivo_id:
            return Response({'error': 'dispositivo_id é obrigatório'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            dispositivo = DispositivoTV.objects.get(id=dispositivo_id)
        except DispositivoTV.DoesNotExist:
            return Response({'error': 'Dispositivo não encontrado'}, status=status.HTTP_404_NOT_FOUND)

        cc = None
        if corporativo_id:
            try:
                cc = ConteudoCorporativo.objects.get(id=corporativo_id)
                titulo = cc.titulo
                tipo_conteudo = cc.tipo
            except ConteudoCorporativo.DoesNotExist:
                pass

        data_hora_fim = timezone.now()
        data_hora_inicio = data_hora_fim - timezone.timedelta(seconds=duracao_segundos)

        # Duração prevista do conteúdo corporativo
        duracao_prevista = cc.duracao_segundos if (cc and cc.duracao_segundos) else duracao_segundos
        if duracao_prevista > 0:
            completamente_exibido = (duracao_segundos / duracao_prevista) >= 0.9
        else:
            completamente_exibido = True

        # Cria o log de exibição
        LogExibicaoWebView.objects.create(
            dispositivo=dispositivo,
            playlist=dispositivo.playlist_atual,
            conteudo_corporativo=cc,
            tipo_conteudo=tipo_conteudo,
            titulo=titulo,
            duracao_segundos=duracao_segundos,
            data_hora_inicio=data_hora_inicio,
            data_hora_fim=data_hora_fim,
            completamente_exibido=completamente_exibido,
        )

        return Response({'success': True}, status=status.HTTP_201_CREATED)


class TVCheckScheduleView(APIView):
    """
    API para o app de TV verificar se deve exibir conteúdo no momento atual
    """
    permission_classes = [permissions.AllowAny]
    
    def get(self, request, identificador_unico):
        """Verifica se o dispositivo deve estar exibindo conteúdo agora"""
        try:
            dispositivo = DispositivoTV.objects.get(
                identificador_unico=identificador_unico,
                ativo=True
            )
            
            # Verifica se está no horário de exibição
            should_display = dispositivo.esta_no_horario_exibicao()
            
            # Atualiza última sincronização
            dispositivo.ultima_sincronizacao = timezone.now()
            dispositivo.save(update_fields=['ultima_sincronizacao'])
            
            response_data = {
                'should_display': should_display,
                'current_time': timezone.localtime(timezone.now()).isoformat(),
                'dispositivo_nome': dispositivo.nome,
                'has_playlist': dispositivo.playlist_atual is not None,
                'horarios_funcionamento': [
                    {
                        'nome': h.nome,
                        'hora_inicio': h.hora_inicio.strftime('%H:%M'),
                        'hora_fim': h.hora_fim.strftime('%H:%M'),
                        'dias_semana': h.dias_semana or [],
                        'ativo': h.ativo,
                    }
                    for h in dispositivo.horarios_funcionamento.filter(ativo=True)
                ],
            }
            
            # Playlists ativas pelo horário (pode ser diferente da padrão)
            playlists_ativas = dispositivo.get_playlists_ativas_por_horario()
            if should_display and playlists_ativas:
                if len(playlists_ativas) == 1:
                    response_data['playlist_id'] = playlists_ativas[0].id
                    response_data['playlist_nome'] = playlists_ativas[0].nome
                else:
                    # Múltiplas playlists mescladas
                    response_data['playlist_id'] = 0  # 0 = múltiplas mescladas
                    response_data['playlist_nome'] = ' + '.join([p.nome for p in playlists_ativas])
                    response_data['playlists_mescladas'] = [p.id for p in playlists_ativas]
            
            # Retorna info dos agendamentos ativos com playlist vinculada
            agendamentos = dispositivo.agendamentos.filter(ativo=True).select_related('playlist')
            if agendamentos.exists():
                response_data['agendamentos'] = [
                    {
                        'nome': ag.nome,
                        'dias_semana': ag.dias_semana,
                        'hora_inicio': ag.hora_inicio.strftime('%H:%M') if ag.hora_inicio else None,
                        'hora_fim': ag.hora_fim.strftime('%H:%M') if ag.hora_fim else None,
                        'playlist_id': ag.playlist_id,
                        'playlist_nome': ag.playlist.nome if ag.playlist else None,
                        'prioridade': ag.prioridade,
                    }
                    for ag in agendamentos
                ]
            else:
                response_data['agendamentos'] = []
                response_data['message'] = 'Sem agendamentos: exibição 24/7'
            
            return Response(response_data)
        
        except DispositivoTV.DoesNotExist:
            return Response(
                {'error': 'Dispositivo não encontrado ou inativo'},
                status=status.HTTP_404_NOT_FOUND
            )


class TVVersionCheckView(APIView):
    """
    Retorna a versão mais recente ativa do APK para o app Android verificar atualizações.
    URL: GET /api/tv/version/
    Resposta:
      {
        "latest_version": "1.2.3",
        "download_url": "https://..../app/download/",
        "force_update": false,
        "size_bytes": 12345678,
        "notes": "..."
      }
    Se não há versão ativa, retorna 404.
    """
    permission_classes = [permissions.AllowAny]

    def get(self, request):
        app_version = AppVersion.get_versao_ativa()
        if not app_version:
            return Response({'error': 'Nenhuma versão disponível'}, status=status.HTTP_404_NOT_FOUND)

        download_url = request.build_absolute_uri('/app/download/')

        return Response({
            'latest_version': app_version.versao,
            'download_url': download_url,
            'force_update': app_version.force_update,
            'size_bytes': app_version.tamanho,
            'notes': app_version.notas_versao,
        })


class TVHeartbeatView(APIView):
    """
    Heartbeat leve: o app Android chama este endpoint periodicamente para indicar
    que o dispositivo está online.

    POST /api/tv/heartbeat/
    Body: { "identificador_unico": "<uuid>" }

    Atualiza ultima_sincronizacao.
    Se o dispositivo estava marcado como offline (alerta já enviado), envia e-mail
    de reconexão e reseta o flag.
    """
    permission_classes = [permissions.AllowAny]

    def post(self, request):
        identificador = request.data.get('identificador_unico', '').strip()
        if not identificador:
            return Response({'error': 'identificador_unico é obrigatório'}, status=status.HTTP_400_BAD_REQUEST)

        try:
            dispositivo = DispositivoTV.objects.select_related(
                'municipio', 'municipio__franqueado', 'franqueado'
            ).get(identificador_unico=identificador, ativo=True)
        except DispositivoTV.DoesNotExist:
            return Response({'error': 'Dispositivo não encontrado ou inativo'}, status=status.HTTP_404_NOT_FOUND)

        estava_offline = dispositivo.alerta_desconexao_enviado

        # Atualiza dados de presença
        now = timezone.now()
        update_fields = ['ultima_sincronizacao']
        dispositivo.ultima_sincronizacao = now

        if estava_offline:
            dispositivo.alerta_desconexao_enviado = False
            update_fields.append('alerta_desconexao_enviado')

        dispositivo.save(update_fields=update_fields)

        # Envia e-mail "voltou online" sem bloquear a resposta
        if estava_offline:
            try:
                from core.alerts import send_online_alert
                import threading
                threading.Thread(target=send_online_alert, args=(dispositivo,), daemon=True).start()
            except Exception:
                pass

        return Response({'status': 'ok', 'ts': now.isoformat()})


@method_decorator(xframe_options_exempt, name='dispatch')
class TVCorporativoHTMLView(APIView):
    """
    Retorna a página HTML completa de conteúdo corporativo para o app de TV.
    O Android WebView carrega esta URL e exibe por duracao_segundos.
    
    URL: /api/tv/corporativo/<tipo>/<playlist_id>/
    tipo: previsao_tempo | cotacoes | noticias
    """
    permission_classes = [permissions.AllowAny]

    def get(self, request, tipo, playlist_id):
        from django.shortcuts import render as django_render
        from .services import buscar_dados_corporativos

        tipo_upper = tipo.upper()
        valid = ['PREVISAO_TEMPO', 'COTACOES', 'NOTICIAS']
        if tipo_upper not in valid:
            return Response({'error': f'Tipo inválido. Válidos: {valid}'}, status=400)

        # Buscar município da playlist para previsão do tempo
        municipio = None
        try:
            playlist = Playlist.objects.select_related('municipio').get(pk=playlist_id)
            municipio = playlist.municipio
        except Playlist.DoesNotExist:
            pass

        # Buscar instância ConteudoCorporativo para obter orientação e configurações
        from .models import ConteudoCorporativo as CC, PlaylistItem
        conteudo = None
        try:
            # Preferir lookup direto via conteudo_id (passado pelo serializer na URL)
            conteudo_id_param = request.GET.get('conteudo_id')
            if conteudo_id_param:
                conteudo = CC.objects.filter(id=conteudo_id_param, ativo=True).first()
            if conteudo is None:
                conteudo = (
                    CC.objects.filter(tipo=tipo_upper, ativo=True)
                    .filter(playlistitem__playlist_id=playlist_id)
                    .first()
                    or CC.objects.filter(tipo=tipo_upper, ativo=True).first()
                )
        except Exception:
            pass

        dados = buscar_dados_corporativos(tipo_upper, municipio=municipio, conteudo=conteudo)
        orientacao = getattr(conteudo, 'orientacao', 'HORIZONTAL') if conteudo else 'HORIZONTAL'

        # dispositivo_id opcional via query param — para auto-log via JS na página
        dispositivo_id = request.GET.get('dispositivo_id', '')

        context = {
            'conteudo_tipo': tipo_upper,
            'dados': dados,
            'orientacao': orientacao,
            'dispositivo_id': dispositivo_id,
            'conteudo_id': conteudo.id if conteudo else '',
            'duracao_segundos': conteudo.duracao_segundos if (conteudo and conteudo.duracao_segundos) else 30,
        }
        return django_render(request, 'corporativo/conteudo_tv.html', context)


class DashboardStatsView(APIView):
    """
    View para estatísticas do dashboard
    """
    permission_classes = [permissions.IsAuthenticated]
    
    def get(self, request):
        user = request.user
        
        if user.is_owner():
            stats = {
                'total_franqueados': User.objects.filter(role='FRANCHISEE').count(),
                'total_clientes': Cliente.objects.count(),
                'total_municipios': Municipio.objects.count(),
                'total_videos': Video.objects.count(),
                'videos_pendentes': Video.objects.filter(status='PENDING').count(),
                'total_playlists': Playlist.objects.count(),
                'total_dispositivos': DispositivoTV.objects.count(),
                'dispositivos_ativos': DispositivoTV.objects.filter(ativo=True).count(),
            }
        elif user.is_franchisee():
            stats = {
                'total_clientes': Cliente.objects.filter(franqueado=user).count(),
                'total_municipios': Municipio.objects.filter(franqueado=user).count(),
                'total_videos': Video.objects.filter(cliente__franqueado=user).count(),
                'videos_pendentes': Video.objects.filter(
                    cliente__franqueado=user,
                    status='PENDING'
                ).count(),
                'total_playlists': Playlist.objects.filter(franqueado=user).count(),
                'total_dispositivos': DispositivoTV.objects.filter(
                    municipio__franqueado=user
                ).count(),
            }
        elif user.is_client():
            try:
                cliente = user.cliente_profile
                stats = {
                    'total_videos': Video.objects.filter(cliente=cliente).count(),
                    'videos_pendentes': Video.objects.filter(
                        cliente=cliente,
                        status='PENDING'
                    ).count(),
                    'videos_aprovados': Video.objects.filter(
                        cliente=cliente,
                        status='APPROVED'
                    ).count(),
                    'videos_rejeitados': Video.objects.filter(
                        cliente=cliente,
                        status='REJECTED'
                    ).count(),
                }
            except Cliente.DoesNotExist:
                stats = {}
        else:
            stats = {}
        
        return Response(stats)


# Web Views (Django Templates)
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth import login, logout, authenticate
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.contrib.auth.forms import AuthenticationForm
from django.urls import reverse
from django.db.models import Count, Sum, Q, Avg, Prefetch
from django.db.models.functions import TruncDate
from django.utils import timezone
from django.http import JsonResponse, Http404
from django.core.paginator import Paginator
from datetime import timedelta, datetime, time
import calendar
import os
from .forms import VideoForm, PlaylistForm, DispositivoTVForm, SegmentoForm, AppVersionForm, ConteudoCorporativoForm, ConfiguracaoAPIForm, HorarioFuncionamentoForm


def home_view(request):
    """Landing page pública com formulário de captura de leads."""
    if request.user.is_authenticated:
        return redirect('dashboard')

    if request.method == 'POST':
        nome     = request.POST.get('nome', '').strip()
        whatsapp = request.POST.get('whatsapp', '').strip()
        email    = request.POST.get('email', '').strip()
        cidade   = request.POST.get('cidade', '').strip()
        segmento = request.POST.get('segmento', '').strip()
        mensagem = request.POST.get('mensagem', '').strip()

        if not nome or not whatsapp:
            return JsonResponse({'success': False, 'error': 'Nome e WhatsApp são obrigatórios.'}, status=400)

        ip = (request.META.get('HTTP_X_FORWARDED_FOR') or '').split(',')[0].strip() or request.META.get('REMOTE_ADDR')
        LandingLead.objects.create(
            nome=nome, whatsapp=whatsapp, email=email,
            cidade=cidade, segmento=segmento, mensagem=mensagem, ip=ip or None,
        )
        return JsonResponse({'success': True})

    return render(request, 'landing/index.html')


def login_view(request):
    """View de login"""
    if request.user.is_authenticated:
        return redirect('dashboard')

    if request.method == 'POST':
        form = AuthenticationForm(request, data=request.POST)
        if form.is_valid():
            username = form.cleaned_data.get('username')
            password = form.cleaned_data.get('password')
            user = authenticate(username=username, password=password)
            if user is not None:
                login(request, user)
                messages.success(request, f'Bem-vindo, {user.get_full_name() or user.username}!')
                return redirect('dashboard')
            else:
                messages.error(request, 'Usuário ou senha inválidos.')
        else:
            messages.error(request, 'Dados inválidos.')
    else:
        form = AuthenticationForm()

    return render(request, 'auth/login.html', {'form': form})


def logout_view(request):
    """View de logout"""
    logout(request)
    messages.info(request, 'Você foi desconectado com sucesso.')
    return redirect('login')


@login_required
def dashboard_view(request):
    """Dashboard principal"""
    user = request.user
    context = {
        'now': timezone.now(),
        'video_stats': {'approved': 0, 'pending': 0, 'rejected': 0},  # Inicializar sempre
    }

    if user.is_owner():
        # Estatísticas para proprietário - consultas otimizadas
        context.update({
            'total_franchisees': User.objects.filter(role='FRANCHISEE').count(),
            'total_municipios': Municipio.objects.count(),
            'total_clients': Cliente.objects.count(),
            'total_devices': DispositivoTV.objects.count(),
        })
        
        # Pré-carregar dados dos municípios com anotações (evita N+1)
        municipios_annotated = Municipio.objects.select_related('franqueado').annotate(
            publico_total=Sum('dispositivos__publico_estimado_mes'),
            dispositivos_count=Count('dispositivos')
        ).order_by('nome')
        
        # Indexar por franqueado_id
        municipios_by_franqueado = {}
        for mun in municipios_annotated:
            municipios_by_franqueado.setdefault(mun.franqueado_id, []).append({
                'municipio': mun,
                'publico_total': mun.publico_total or 0,
                'dispositivos_count': mun.dispositivos_count or 0,
            })
        
        # Pré-carregar clientes por franqueado
        all_clientes = Cliente.objects.select_related('user', 'segmento').order_by('empresa')
        clientes_by_franqueado = {}
        for cli in all_clientes:
            clientes_by_franqueado.setdefault(cli.franqueado_id, []).append(cli)
        
        # Pré-carregar playlists por franqueado
        all_playlists = Playlist.objects.select_related('municipio').order_by('nome')
        playlists_by_franqueado = {}
        for pl in all_playlists:
            playlists_by_franqueado.setdefault(pl.franqueado_id, []).append(pl)
        
        # Pré-carregar contagens de dispositivos por franqueado
        dispositivos_por_franqueado = dict(
            DispositivoTV.objects.values_list('municipio__franqueado').annotate(
                total=Count('id')
            ).values_list('municipio__franqueado', 'total')
        )
        
        # Pré-carregar público total por franqueado
        publico_por_franqueado = dict(
            DispositivoTV.objects.values('municipio__franqueado').annotate(
                total=Sum('publico_estimado_mes')
            ).values_list('municipio__franqueado', 'total')
        )
        
        # Montar visão hierárquica sem queries adicionais
        franqueados = User.objects.filter(role='FRANCHISEE').order_by('username')
        franqueados_data = []
        
        for franqueado in franqueados:
            fid = franqueado.id
            muns = municipios_by_franqueado.get(fid, [])
            clis = clientes_by_franqueado.get(fid, [])
            pls = playlists_by_franqueado.get(fid, [])
            
            franqueados_data.append({
                'franqueado': franqueado,
                'municipios': muns,
                'clientes': clis,
                'playlists': pls,
                'stats': {
                    'total_municipios': len(muns),
                    'total_clientes': len(clis),
                    'total_playlists': len(pls),
                    'total_dispositivos': dispositivos_por_franqueado.get(fid, 0),
                    'publico_total': publico_por_franqueado.get(fid, 0) or 0,
                }
            })
        
        context['franqueados_data'] = franqueados_data
        
    elif user.is_franchisee():
        # Estatísticas para franqueado
        municipios = Municipio.objects.filter(franqueado=user)
        clientes = Cliente.objects.filter(franqueado=user)
        dispositivos = DispositivoTV.objects.filter(municipio__franqueado=user)

        context.update({
            'total_municipios': municipios.count(),
            'total_clients': clientes.count(),
            'total_devices': dispositivos.count(),
        })
    else:
        # Estatísticas para cliente
        try:
            cliente = user.cliente_profile
            videos = Video.objects.filter(cliente=cliente)
            
            # Informações do franqueado
            franqueado = cliente.franqueado
            
            # Municípios do franqueado com público - query otimizada com anotações
            municipios = Municipio.objects.filter(franqueado=franqueado).annotate(
                publico_total=Sum('dispositivos__publico_estimado_mes'),
                dispositivos_count=Count('dispositivos')
            ).order_by('nome')
            
            municipios_com_publico = [{
                'municipio': m,
                'publico_total': m.publico_total or 0,
                'dispositivos_count': m.dispositivos_count or 0,
            } for m in municipios]
            
            # Público total do franqueado
            publico_total_franqueado = sum(m['publico_total'] for m in municipios_com_publico)
            
            context.update({
                'total_videos': videos.count(),
                'approved_videos': videos.filter(status='APPROVED').count(),
                'cliente': cliente,
                'franqueado': franqueado,
                'municipios_franqueado': municipios_com_publico,
                'publico_total_franqueado': publico_total_franqueado,
            })
        except Cliente.DoesNotExist:
            context.update({
                'total_videos': 0,
                'approved_videos': 0,
            })

    # Atividades recentes (últimos 10 dias)
    recent_activities = []
    cutoff_date = timezone.now() - timedelta(days=10)

    if user.is_owner():
        # Atividades de todos
        recent_videos = Video.objects.filter(
            created_at__gte=cutoff_date
        ).select_related('cliente').order_by('-created_at')[:5]
        for video in recent_videos:
            recent_activities.append({
                'description': f'Novo vídeo "{video.titulo}" enviado por {video.cliente.empresa}',
                'timestamp': video.created_at
            })
    elif user.is_franchisee():
        # Atividades dos clientes do franqueado
        recent_videos = Video.objects.filter(
            cliente__franqueado=user,
            created_at__gte=cutoff_date
        ).select_related('cliente').order_by('-created_at')[:5]
        for video in recent_videos:
            recent_activities.append({
                'description': f'Novo vídeo "{video.titulo}" enviado por {video.cliente.empresa}',
                'timestamp': video.created_at
            })
    else:
        # Atividades do próprio cliente
        try:
            cliente = user.cliente_profile
            recent_videos = Video.objects.filter(
                cliente=cliente,
                created_at__gte=cutoff_date
            ).order_by('-created_at')[:5]
            for video in recent_videos:
                status_text = {
                    'PENDING': 'pendente',
                    'APPROVED': 'aprovado',
                    'REJECTED': 'rejeitado'
                }.get(video.status, video.status)
                recent_activities.append({
                    'description': f'Vídeo "{video.titulo}" está {status_text}',
                    'timestamp': video.created_at
                })
        except Cliente.DoesNotExist:
            pass

    context['recent_activities'] = sorted(recent_activities, key=lambda x: x['timestamp'], reverse=True)[:5]

    # Estatísticas do gráfico de pizza (para clientes)
    if not user.is_owner() and not user.is_franchisee():
        try:
            cliente = user.cliente_profile
            videos = Video.objects.filter(cliente=cliente)
            context['video_stats'] = {
                'approved': videos.filter(status='APPROVED').count(),
                'pending': videos.filter(status='PENDING').count(),
                'rejected': videos.filter(status='REJECTED').count(),
            }
        except Cliente.DoesNotExist:
            context['video_stats'] = {'approved': 0, 'pending': 0, 'rejected': 0}

    return render(request, 'dashboard/dashboard.html', context)


@login_required
def cliente_metricas_view(request):
    """Dashboard de métricas de visibilidade para clientes"""
    user = request.user
    
    # Verificar se é cliente
    if not user.is_client():
        messages.error(request, 'Esta página é exclusiva para clientes.')
        return redirect('dashboard')
    
    try:
        cliente = user.cliente_profile
    except Cliente.DoesNotExist:
        messages.error(request, 'Perfil de cliente não encontrado.')
        return redirect('dashboard')
    
    # Obter vídeos aprovados do cliente
    videos_cliente = Video.objects.filter(cliente=cliente, status='APPROVED', ativo=True)
    
    # Encontrar playlists que contêm vídeos do cliente
    playlist_ids = PlaylistItem.objects.filter(
        video__in=videos_cliente, 
        ativo=True
    ).values_list('playlist_id', flat=True).distinct()
    
    playlists = Playlist.objects.filter(id__in=playlist_ids, ativa=True)
    
    # Dispositivos que usam essas playlists - prefetch agendamentos para evitar N+1
    dispositivos = DispositivoTV.objects.filter(
        playlist_atual__in=playlists,
        ativo=True
    ).select_related('municipio').prefetch_related(
        Prefetch('agendamentos', queryset=AgendamentoExibicao.objects.filter(ativo=True))
    ).distinct()
    
    # Métricas básicas
    telas_ativas = dispositivos.count()
    
    # Público impactado (soma do público estimado de todos os dispositivos)
    publico_impactado = dispositivos.aggregate(
        total=Sum('publico_estimado_mes')
    )['total'] or 0
    
    # Calcular tempo total de exibição baseado nos agendamentos
    tempo_total_segundos = 0
    dispositivos_detalhes = []
    
    # Dias úteis no mês atual
    now = timezone.now()
    _, dias_no_mes = calendar.monthrange(now.year, now.month)
    
    for dispositivo in dispositivos:
        agendamentos = dispositivo.agendamentos.all()  # já prefetched
        
        horas_dia = 0
        if agendamentos.exists():
            # Calcular horas por semana baseado nos agendamentos
            for agendamento in agendamentos:
                hora_inicio = agendamento.hora_inicio
                hora_fim = agendamento.hora_fim
                
                # Calcular duração em horas
                inicio_segundos = hora_inicio.hour * 3600 + hora_inicio.minute * 60 + hora_inicio.second
                fim_segundos = hora_fim.hour * 3600 + hora_fim.minute * 60 + hora_fim.second
                duracao_segundos = fim_segundos - inicio_segundos
                
                if duracao_segundos > 0:
                    horas_diarias = duracao_segundos / 3600
                    dias_ativos = len(agendamento.dias_semana) if agendamento.dias_semana else 0
                    # Média de dias por mês (considerando semanas)
                    dias_mes = (dias_ativos / 7) * dias_no_mes
                    horas_dia += horas_diarias * dias_mes
        else:
            # Sem agendamento = 12h por dia, todos os dias
            horas_dia = 12 * dias_no_mes
        
        tempo_total_segundos += horas_dia * 3600
        
        dispositivos_detalhes.append({
            'dispositivo': dispositivo,
            'horas_mes': round(horas_dia, 1),
            'localizacao': dispositivo.localizacao or dispositivo.municipio.nome
        })
    
    tempo_total_horas = tempo_total_segundos / 3600
    
    # Duração média dos vídeos do cliente (em segundos)
    duracao_media_video = videos_cliente.aggregate(
        media=Avg('duracao_segundos')
    )['media'] or 15  # Default 15 segundos
    
    # Total de vídeos em todas as playlists relevantes
    total_videos_playlists = PlaylistItem.objects.filter(
        playlist__in=playlists,
        ativo=True
    ).count()
    
    # Proporção de vídeos do cliente nas playlists
    videos_cliente_em_playlists = PlaylistItem.objects.filter(
        playlist__in=playlists,
        video__in=videos_cliente,
        ativo=True
    ).count()
    
    # Evitar divisão por zero
    proporcao_cliente = videos_cliente_em_playlists / max(total_videos_playlists, 1)
    
    # Inserções totais: tempo_total (em segundos) / duração média do vídeo * proporção de vídeos do cliente
    insercoes_totais = int((tempo_total_segundos / max(duracao_media_video, 1)) * proporcao_cliente)
    
    # Anunciantes ativos (clientes únicos com vídeos nas mesmas playlists)
    anunciantes_ids = PlaylistItem.objects.filter(
        playlist__in=playlists,
        ativo=True
    ).values_list('video__cliente_id', flat=True).distinct()
    anunciantes_ativos = len(set(anunciantes_ids))
    
    # Tempo por anúncio (duração média dos vídeos do cliente)
    tempo_por_anuncio = duracao_media_video
    
    # Inserções por anunciante
    insercoes_por_anunciante = int(insercoes_totais / max(anunciantes_ativos, 1)) if anunciantes_ativos > 0 else insercoes_totais
    
    # Tempo corrido da marca (inserções * duração média)
    tempo_corrido_segundos = insercoes_totais * duracao_media_video
    tempo_corrido_horas = tempo_corrido_segundos / 3600
    
    # Dados do mês anterior para comparação (usando logs reais se disponíveis)
    mes_anterior = now.replace(day=1) - timedelta(days=1)
    logs_mes_anterior = LogExibicao.objects.filter(
        video__cliente=cliente,
        data_hora_inicio__year=mes_anterior.year,
        data_hora_inicio__month=mes_anterior.month
    ).count()
    
    logs_mes_atual = LogExibicao.objects.filter(
        video__cliente=cliente,
        data_hora_inicio__year=now.year,
        data_hora_inicio__month=now.month
    ).count()
    
    # Calcular variação percentual
    if logs_mes_anterior > 0:
        variacao_percentual = ((logs_mes_atual - logs_mes_anterior) / logs_mes_anterior) * 100
    else:
        variacao_percentual = 100 if logs_mes_atual > 0 else 0
    
    # Função auxiliar para formatar números grandes
    def formatar_numero(n):
        if n >= 1000000:
            return f"{n / 1000000:.1f}M"
        elif n >= 1000:
            return f"{n / 1000:.1f}K"
        return str(int(n))
    
    context = {
        'now': now,
        'cliente': cliente,
        'telas_ativas': telas_ativas,
        'publico_impactado': publico_impactado,
        'publico_impactado_formatado': formatar_numero(publico_impactado),
        'tempo_total_horas': round(tempo_total_horas, 1),
        'insercoes_totais': insercoes_totais,
        'insercoes_totais_formatado': formatar_numero(insercoes_totais),
        'anunciantes_ativos': anunciantes_ativos,
        'tempo_por_anuncio': round(tempo_por_anuncio, 0),
        'insercoes_por_anunciante': insercoes_por_anunciante,
        'insercoes_por_anunciante_formatado': formatar_numero(insercoes_por_anunciante),
        'tempo_corrido_horas': round(tempo_corrido_horas, 1),
        'variacao_percentual': round(variacao_percentual, 0),
        'dispositivos_detalhes': dispositivos_detalhes,
        'videos_cliente': videos_cliente,
        'playlists': playlists,
        'duracao_media_video': round(duracao_media_video, 0),
        'proporcao_cliente': round(proporcao_cliente * 100, 1),
        'logs_mes_atual': logs_mes_atual,
    }
    
    return render(request, 'dashboard/cliente_metricas.html', context)


# Video Views
@login_required
def video_list_view(request):
    """Lista de vídeos"""
    user = request.user
    videos = Video.objects.select_related(
        'cliente', 'cliente__user', 'cliente__segmento'
    ).annotate(
        qrcode_clicks_count=Count('qrcode_clicks')
    )

    # Filtros
    search = request.GET.get('search', '')
    cliente_filter = request.GET.get('cliente', '')
    status_filter = request.GET.get('status', '')
    orphaned_filter = request.GET.get('orphaned', '')

    if search:
        videos = videos.filter(Q(titulo__icontains=search) | Q(descricao__icontains=search))

    if cliente_filter:
        videos = videos.filter(cliente_id=cliente_filter)

    if status_filter:
        videos = videos.filter(status=status_filter.upper())
    
    # Filtro para arquivos órfãos (apenas para OWNER)
    if orphaned_filter == 'true' and user.is_owner():
        orphaned_ids = []
        for video in videos.only('id', 'arquivo'):
            if not video.arquivo_existe():
                orphaned_ids.append(video.id)
        videos = videos.filter(id__in=orphaned_ids)

    # Controle de permissões
    if user.is_franchisee():
        clientes_ids = Cliente.objects.filter(franqueado=user).values_list('id', flat=True)
        videos = videos.filter(cliente_id__in=clientes_ids)
    elif not user.is_owner():
        try:
            cliente = user.cliente_profile
            videos = videos.filter(cliente=cliente)
        except Cliente.DoesNotExist:
            videos = Video.objects.none()

    # Paginação
    paginator = Paginator(videos.order_by('-created_at'), 12)
    page_number = request.GET.get('page')
    page_obj = paginator.get_page(page_number)

    # Clientes para filtro (se aplicável)
    clientes = Cliente.objects.only('id', 'empresa')
    if user.is_franchisee():
        clientes = clientes.filter(franqueado=user)

    context = {
        'videos': page_obj,
        'clientes': clientes,
    }

    return render(request, 'videos/video_list.html', context)


@login_required
def video_create_view(request):
    """Criar novo vídeo"""
    user = request.user
    
    # Verificar se é cliente e buscar perfil
    if user.is_client():
        try:
            cliente = user.cliente_profile
        except Cliente.DoesNotExist:
            messages.error(request, 'Perfil de cliente não encontrado.')
            return redirect('video_list')
        
        if request.method == 'POST':
            form = VideoForm(request.POST, request.FILES, user=user)
            if form.is_valid():
                video = form.save(commit=False)
                video.cliente = cliente
                video.status = 'PENDING'  # Garante que clientes sempre enviam como pendente
                video.save()
                messages.success(request, 'Vídeo enviado com sucesso! Aguarde aprovação.')
                return redirect('video_list')
            else:
                import logging
                logger = logging.getLogger(__name__)
                logger.warning(f'[VIDEO_UPLOAD] Form inválido para cliente {cliente}: {form.errors}')
        else:
            form = VideoForm(user=user)
        
        return render(request, 'videos/video_form.html', {'form': form})
    
    # Owner e Franqueado podem fazer upload para qualquer cliente
    elif user.is_owner() or user.is_franchisee():
        # Filtrar clientes disponíveis
        if user.is_owner():
            clientes = Cliente.objects.filter(ativo=True).select_related('user').order_by('empresa')
        else:
            clientes = Cliente.objects.filter(franqueado=user, ativo=True).select_related('user').order_by('empresa')
        
        if request.method == 'POST':
            cliente_id = request.POST.get('cliente')
            titulo = request.POST.get('titulo')
            arquivo = request.FILES.get('arquivo')
            url_externa = request.POST.get('url_externa', '').strip() or None

            erros = []
            if not cliente_id:
                erros.append('cliente')
            if not titulo:
                erros.append('título')
            if not arquivo and not url_externa:
                erros.append('arquivo de vídeo ou URL externa')

            if erros:
                messages.error(request, f'Campo(s) obrigatório(s) faltando: {", ".join(erros)}.')
            else:
                try:
                    cliente = Cliente.objects.get(id=cliente_id)

                    # Verificar se franqueado tem permissão
                    if user.is_franchisee() and cliente.franqueado != user:
                        messages.error(request, 'Você não tem permissão para enviar vídeos para este cliente.')
                        return redirect('video_list')

                    descricao = request.POST.get('descricao')
                    qrcode_url_destino = request.POST.get('qrcode_url_destino', '').strip() or None
                    qrcode_descricao = request.POST.get('qrcode_descricao', '').strip() or None
                    texto_tarja = request.POST.get('texto_tarja', '').strip() or None
                    orientacao = request.POST.get('orientacao', 'HORIZONTAL')
                    duracao_segundos = int(request.POST.get('duracao_segundos') or 0)

                    create_kwargs = dict(
                        cliente=cliente,
                        titulo=titulo,
                        descricao=descricao,
                        url_externa=url_externa,
                        qrcode_url_destino=qrcode_url_destino,
                        qrcode_descricao=qrcode_descricao,
                        texto_tarja=texto_tarja,
                        orientacao=orientacao,
                        duracao_segundos=duracao_segundos,
                        status='PENDING',
                    )
                    if arquivo:
                        create_kwargs['arquivo'] = arquivo

                    video = Video.objects.create(**create_kwargs)
                    messages.success(request, f'Vídeo enviado com sucesso para {cliente.empresa}!')
                    return redirect('video_list')
                except Cliente.DoesNotExist:
                    messages.error(request, 'Cliente não encontrado.')
        
        return render(request, 'videos/video_form.html', {
            'clientes': clientes,
            'show_cliente_select': True,
        })
    
    else:
        messages.error(request, 'Você não tem permissão para enviar vídeos.')
        return redirect('video_list')


@login_required
def video_bulk_upload_view(request):
    """Upload em lote de vídeos (owner e franqueado apenas).
    GET  → renderiza a página de bulk upload.
    POST → endpoint AJAX: recebe UM arquivo por vez, retorna JSON.
    """
    from django.http import JsonResponse
    user = request.user

    if not user.is_owner() and not user.is_franchisee():
        if request.method == 'POST':
            return JsonResponse({'success': False, 'error': 'Sem permissão.'}, status=403)
        messages.error(request, 'Acesso negado.')
        return redirect('video_list')

    if user.is_owner():
        clientes = Cliente.objects.filter(ativo=True).select_related('user').order_by('empresa')
    else:
        clientes = Cliente.objects.filter(franqueado=user, ativo=True).select_related('user').order_by('empresa')

    if request.method == 'POST':
        cliente_id  = request.POST.get('cliente')
        titulo      = request.POST.get('titulo', '').strip()
        arquivo     = request.FILES.get('arquivo')
        orientacao  = request.POST.get('orientacao', 'HORIZONTAL')
        descricao   = request.POST.get('descricao', '').strip() or None
        texto_tarja = request.POST.get('texto_tarja', '').strip() or None
        qrcode_url  = request.POST.get('qrcode_url_destino', '').strip() or None

        if not cliente_id:
            return JsonResponse({'success': False, 'error': 'Cliente não selecionado.'}, status=400)
        if not titulo:
            return JsonResponse({'success': False, 'error': 'Título obrigatório.'}, status=400)
        if not arquivo:
            return JsonResponse({'success': False, 'error': 'Arquivo não enviado.'}, status=400)

        try:
            cliente = Cliente.objects.get(id=cliente_id)
        except Cliente.DoesNotExist:
            return JsonResponse({'success': False, 'error': 'Cliente não encontrado.'}, status=400)

        if user.is_franchisee() and cliente.franqueado != user:
            return JsonResponse({'success': False, 'error': 'Sem permissão para este cliente.'}, status=403)

        video = Video.objects.create(
            cliente=cliente,
            titulo=titulo,
            descricao=descricao,
            arquivo=arquivo,
            orientacao=orientacao,
            texto_tarja=texto_tarja,
            qrcode_url_destino=qrcode_url,
            status='PENDING',
        )
        return JsonResponse({'success': True, 'id': video.id, 'titulo': video.titulo})

    return render(request, 'videos/video_bulk_upload.html', {
        'clientes': clientes,
    })


@login_required
def video_update_view(request, pk):
    """Atualizar vídeo"""
    user = request.user
    video = get_object_or_404(Video, pk=pk)

    # Verificar permissões
    if not user.is_owner():
        if user.is_franchisee():
            if video.cliente.franqueado != user:
                messages.error(request, 'Você não tem permissão para editar este vídeo.')
                return redirect('video_list')
        elif not user.is_client() or video.cliente.user != user:
            messages.error(request, 'Você não tem permissão para editar este vídeo.')
            return redirect('video_list')

    if request.method == 'POST':
        form = VideoForm(request.POST, request.FILES, instance=video, user=user)
        if form.is_valid():
            updated = form.save(commit=False)
            # Status é gerenciado via aprovar/rejeitar — nunca altera pelo form de edição
            updated.status = video.status
            updated.save()
            messages.success(request, 'Vídeo atualizado com sucesso!')
            return redirect('video_list')
        else:
            import logging as _log
            _log.getLogger(__name__).warning(f'[VIDEO_UPDATE] Form inválido pk={pk}: {form.errors}')
            messages.error(request, 'Não foi possível salvar. Verifique os campos destacados abaixo.')
    else:
        form = VideoForm(instance=video, user=user)

    return render(request, 'videos/video_form.html', {'form': form, 'video': video})


@login_required
def video_approve_view(request, pk):
    """Aprovar vídeo"""
    user = request.user
    
    print(f"DEBUG: Tentando aprovar vídeo {pk} por usuário {user.username}")
    
    # Apenas owners e franqueados podem aprovar
    if not user.is_owner() and not user.is_franchisee():
        print(f"DEBUG: Usuário {user.username} não tem permissão (role: {user.role})")
        return JsonResponse({'success': False, 'error': 'Sem permissão'}, status=403)
    
    video = get_object_or_404(Video, pk=pk)
    print(f"DEBUG: Video encontrado: {video.titulo}, Status atual: {video.status}")
    
    # Franqueados só podem aprovar vídeos de seus clientes
    if user.is_franchisee() and video.cliente.franqueado != user:
        print(f"DEBUG: Franqueado não gerencia este cliente")
        return JsonResponse({'success': False, 'error': 'Sem permissão'}, status=403)
    
    video.status = 'APPROVED'
    video.save()
    print(f"DEBUG: Video {video.titulo} aprovado! Novo status: {video.status}")
    
    messages.success(request, f'Vídeo "{video.titulo}" aprovado com sucesso!')
    return JsonResponse({'success': True})


@login_required
def video_reject_view(request, pk):
    """Rejeitar vídeo"""
    user = request.user
    
    # Apenas owners e franqueados podem rejeitar
    if not user.is_owner() and not user.is_franchisee():
        return JsonResponse({'success': False, 'error': 'Sem permissão'}, status=403)
    
    video = get_object_or_404(Video, pk=pk)
    
    # Franqueados só podem rejeitar vídeos de seus clientes
    if user.is_franchisee() and video.cliente.franqueado != user:
        return JsonResponse({'success': False, 'error': 'Sem permissão'}, status=403)
    
    video.status = 'REJECTED'
    video.save()
    
    messages.success(request, f'Vídeo "{video.titulo}" rejeitado.')
    return JsonResponse({'success': True})


@login_required
def video_delete_view(request, pk):
    """Deletar vídeo"""
    import os
    user = request.user
    video = get_object_or_404(Video, pk=pk)
    
    # Verificar permissões
    if not user.is_owner():
        if user.is_franchisee():
            if video.cliente.franqueado != user:
                return JsonResponse({'success': False, 'error': 'Sem permissão'}, status=403)
        elif not user.is_client() or video.cliente.user != user:
            return JsonResponse({'success': False, 'error': 'Sem permissão'}, status=403)
    
    titulo = video.titulo
    arquivo_existe = video.arquivo and os.path.exists(video.arquivo.path)
    
    # Deleta o registro (e o arquivo se existir)
    video.delete()
    
    if arquivo_existe:
        messages.success(request, f'Vídeo "{titulo}" excluído com sucesso!')
    else:
        messages.success(request, f'Registro do vídeo "{titulo}" excluído (arquivo já não existia no servidor)')
    
    return JsonResponse({'success': True})


@login_required
def video_convert_mp4_view(request, pk):
    """Converte um vídeo para MP4 1080p FireTV-safe (V6):
    H.264 Main 4.0, 1080×1920 (vertical) ou 1920×1080 (horizontal),
    BT.709, VBV 5M, GOP 60, brand mp42.
    """
    import subprocess
    import shutil
    import tempfile
    from django.conf import settings

    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Método não permitido'}, status=405)

    user = request.user
    video = get_object_or_404(Video, pk=pk)

    # Verificar permissões (owner, franqueado do cliente, ou o próprio cliente)
    if not user.is_owner():
        if user.is_franchisee():
            if video.cliente.franqueado != user:
                return JsonResponse({'success': False, 'error': 'Sem permissão'}, status=403)
        elif not user.is_client() or video.cliente.user != user:
            return JsonResponse({'success': False, 'error': 'Sem permissão'}, status=403)

    # Verificar se o arquivo existe
    if not video.arquivo or not video.arquivo_existe():
        return JsonResponse({'success': False, 'error': 'Arquivo de vídeo não encontrado no servidor'}, status=404)

    # Verificar se ffmpeg está disponível
    if not shutil.which('ffmpeg'):
        return JsonResponse({'success': False, 'error': 'ffmpeg não está instalado no servidor'}, status=500)

    tmp_input = None
    tmp_output = None
    is_remote = False
    try:
        from django.core.files import File as DjangoFile

        # ── Obter input (local ou download do R2) ──────────────────────────
        try:
            input_path = video.arquivo.path  # storage local
        except NotImplementedError:
            is_remote = True
            ext = os.path.splitext(video.arquivo.name)[1] or '.mp4'
            fd, tmp_input = tempfile.mkstemp(suffix=ext, dir='/tmp')
            os.close(fd)
            with video.arquivo.open('rb') as src:
                with open(tmp_input, 'wb') as dst:
                    shutil.copyfileobj(src, dst)
            input_path = tmp_input

        fd, tmp_output = tempfile.mkstemp(suffix='.mp4', dir='/tmp')
        os.close(fd)

        input_basename = os.path.splitext(os.path.basename(video.arquivo.name))[0]
        output_filename = f'{input_basename}.mp4'

        # Detectar orientação e resolução real do vídeo via ffprobe
        orientacao, orig_w, orig_h = Video._detectar_orientacao_video(input_path)
        scale_filter = Video._calcular_scale_filter(orig_w, orig_h, orientacao)

        # Bitrate 1080p — VBV calibrado para Fire TV (V6 validado)
        bitrate = '5M'
        maxrate = '5M'
        bufsize = '10M'

        # Pipeline 1080p FireTV-safe — Main 4.0 + BT.709 + VBV + GOP 60 + mp42
        cmd = [
            'ffmpeg', '-y',
            '-i', input_path,
            '-vf', scale_filter,
            '-map_metadata', '-1',
            '-c:v', 'libx264',
            '-profile:v', 'main',
            '-level', '4.0',
            '-pix_fmt', 'yuv420p',
            '-r', '30',
            '-b:v', bitrate,
            '-maxrate', maxrate,
            '-bufsize', bufsize,
            '-g', '60',
            '-keyint_min', '60',
            '-preset', 'medium',
            '-color_range', 'tv',
            '-colorspace', 'bt709',
            '-color_primaries', 'bt709',
            '-color_trc', 'bt709',
            '-c:a', 'aac',
            '-b:a', '160k',
            '-ar', '44100',
            '-movflags', '+faststart',
            '-brand', 'mp42',
            '-tag:v', 'avc1',
            '-vsync', 'cfr',
            tmp_output
        ]

        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=900  # 15 minutos
        )

        if result.returncode != 0:
            return JsonResponse({
                'success': False,
                'error': f'Erro na conversão: {result.stderr[:500]}'
            }, status=500)

        new_size = os.path.getsize(tmp_output)

        if is_remote:
            # ── Upload para R2 ───────────────────────────────────
            storage_name = os.path.splitext(video.arquivo.name)[0] + '.mp4'
            try:
                video.arquivo.storage.delete(video.arquivo.name)
            except Exception:
                pass
            with open(tmp_output, 'rb') as f:
                saved_name = video.arquivo.storage.save(storage_name, DjangoFile(f))
            Video.objects.filter(pk=video.pk).update(
                arquivo=saved_name,
                orientacao=orientacao,
            )
        else:
            # ── Storage local: substituir no disco ────────────────────
            input_dir = os.path.dirname(input_path)
            output_path = os.path.join(input_dir, output_filename)
            if os.path.exists(input_path):
                try:
                    os.remove(input_path)
                except OSError:
                    pass
            shutil.move(tmp_output, output_path)
            tmp_output = None  # já movido, não precisa remover no finally
            relative_path = os.path.relpath(output_path, settings.MEDIA_ROOT).replace('\\', '/')
            Video.objects.filter(pk=video.pk).update(
                arquivo=relative_path,
                orientacao=orientacao,
            )

        return JsonResponse({
            'success': True,
            'message': f'Vídeo convertido para MP4 com sucesso! ({orientacao})',
            'new_filename': output_filename,
            'new_size': new_size,
            'orientacao': orientacao,
        })

    except subprocess.TimeoutExpired:
        return JsonResponse({
            'success': False,
            'error': 'Conversão excedeu o tempo limite de 15 minutos. Tente com um vídeo menor.'
        }, status=504)
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': f'Erro inesperado: {str(e)}'
        }, status=500)
    finally:
        for tmp in (tmp_input, tmp_output):
            if tmp and os.path.exists(tmp):
                try:
                    os.remove(tmp)
                except OSError:
                    pass


@login_required
def video_qrcode_metricas_view(request, pk):
    """Métricas de cliques do QR Code de um vídeo"""
    user = request.user
    video = get_object_or_404(Video, pk=pk)
    
    # Verificar permissões
    if not user.is_owner():
        if user.is_franchisee():
            if video.cliente.franqueado != user:
                messages.error(request, 'Você não tem permissão para ver estas métricas.')
                return redirect('video_list')
        elif video.cliente.user != user:
            messages.error(request, 'Você não tem permissão para ver estas métricas.')
            return redirect('video_list')
    
    # Buscar cliques
    clicks = video.qrcode_clicks.all().order_by('-created_at')
    
    # Estatísticas
    total_clicks = clicks.count()
    
    # Cliques por dia (últimos 30 dias)
    hoje = timezone.now().date()
    inicio_periodo = hoje - timedelta(days=30)
    
    clicks_por_dia = (
        clicks.filter(created_at__date__gte=inicio_periodo)
        .annotate(data=TruncDate('created_at'))
        .values('data')
        .annotate(count=Count('id'))
        .order_by('data')
    )
    
    # Preparar dados para o gráfico
    dias_labels = []
    dias_valores = []
    clicks_dict = {item['data']: item['count'] for item in clicks_por_dia}
    
    for i in range(31):  # 31 dias para incluir hoje
        dia = inicio_periodo + timedelta(days=i)
        dias_labels.append(dia.strftime('%d/%m'))
        dias_valores.append(clicks_dict.get(dia, 0))
    
    # Top IPs (para detectar possível fraude)
    top_ips = (
        clicks.exclude(ip_address__isnull=True)
        .values('ip_address')
        .annotate(total=Count('id'))
        .order_by('-total')[:10]
    )
    
    # Paginação dos cliques
    paginator = Paginator(clicks, 20)
    page_number = request.GET.get('page')
    page_obj = paginator.get_page(page_number)
    
    # Converter listas para JSON para o gráfico (JavaScript precisa de JSON válido)
    import json
    
    context = {
        'video': video,
        'clicks': page_obj,
        'total_clicks': total_clicks,
        'dias_labels': json.dumps(dias_labels),
        'dias_valores': json.dumps(dias_valores),
        'top_ips': top_ips,
    }
    
    return render(request, 'videos/video_qrcode_metricas.html', context)


# Playlist Views
@login_required
def playlist_list_view(request):
    """Lista de playlists"""
    user = request.user
    playlists = Playlist.objects.select_related(
        'municipio', 'franqueado'
    ).prefetch_related(
        Prefetch('items', queryset=PlaylistItem.objects.filter(ativo=True).select_related('video').only(
            'id', 'playlist_id', 'video_id', 'ativo', 'video__titulo'
        )),
        'dispositivos'
    )

    # Controle de permissões
    if user.is_franchisee():
        playlists = playlists.filter(franqueado=user)
    elif not user.is_owner():
        messages.error(request, 'Você não tem permissão para ver playlists.')
        return redirect('dashboard')

    # Paginação
    paginator = Paginator(playlists.order_by('-created_at'), 9)
    page_number = request.GET.get('page')
    page_obj = paginator.get_page(page_number)

    context = {
        'playlists': page_obj,
    }

    return render(request, 'playlists/playlist_list.html', context)


@login_required
def playlist_create_view(request):
    """Criar nova playlist"""
    user = request.user

    # Verificar permissões
    if not user.is_franchisee() and not user.is_owner():
        messages.error(request, 'Apenas franqueados podem criar playlists.')
        return redirect('playlist_list')

    # Obter vídeos disponíveis
    available_videos = Video.objects.filter(status='APPROVED')
    if user.is_franchisee():
        clientes_ids = Cliente.objects.filter(franqueado=user).values_list('id', flat=True)
        available_videos = available_videos.filter(cliente_id__in=clientes_ids)

    # Conteúdos corporativos disponíveis
    conteudos_corporativos = ConteudoCorporativo.objects.filter(ativo=True)

    if request.method == 'POST':
        form = PlaylistForm(request.POST, user=user)
        selected_videos = request.POST.getlist('selected_videos')
        # Filtrar valores vazios
        selected_videos = [v for v in selected_videos if v and v.strip()]

        if form.is_valid():
            playlist = form.save(commit=False)
            # Definir franqueado baseado no município selecionado ou no usuário
            if user.is_franchisee():
                playlist.franqueado = user
            elif playlist.municipio:
                playlist.franqueado = playlist.municipio.franqueado
            playlist.save()

            # Adicionar itens selecionados (vídeos + corporativos)
            if selected_videos:
                for order, item_id in enumerate(selected_videos, 1):
                    if item_id.startswith('corp_'):
                        corp_id = int(item_id.replace('corp_', ''))
                        PlaylistItem.objects.create(
                            playlist=playlist,
                            conteudo_corporativo_id=corp_id,
                            ordem=order
                        )
                    else:
                        PlaylistItem.objects.create(
                            playlist=playlist,
                            video_id=int(item_id),
                            ordem=order
                        )
                messages.success(request, f'Playlist criada com sucesso com {len(selected_videos)} item(ns)!')
            else:
                messages.success(request, 'Playlist criada com sucesso! Você pode adicionar itens depois.')

            return redirect('playlist_list')
    else:
        form = PlaylistForm(user=user)

    context = {
        'form': form,
        'available_videos': available_videos,
        'conteudos_corporativos': conteudos_corporativos,
    }

    return render(request, 'playlists/playlist_form.html', context)


@login_required
def playlist_detail_view(request, pk):
    """Detalhes da playlist"""
    playlist = get_object_or_404(Playlist, pk=pk)
    user = request.user
    
    # Verificar permissões
    if user.is_franchisee() and playlist.franqueado != user:
        messages.error(request, 'Você não tem permissão para ver esta playlist.')
        return redirect('playlist_list')
    
    # Obter itens da playlist ordenados
    items = playlist.items.all().select_related('video', 'video__cliente', 'conteudo_corporativo').order_by('ordem')
    
    context = {
        'playlist': playlist,
        'items': items,
    }
    
    return render(request, 'playlists/playlist_detail.html', context)


@login_required
def playlist_update_view(request, pk):
    """Editar playlist"""
    playlist = get_object_or_404(Playlist, pk=pk)
    user = request.user
    
    # Verificar permissões
    if user.is_franchisee() and playlist.franqueado != user:
        messages.error(request, 'Você não tem permissão para editar esta playlist.')
        return redirect('playlist_list')
    
    # Obter vídeos disponíveis
    available_videos = Video.objects.filter(status='APPROVED')
    if user.is_franchisee():
        clientes_ids = Cliente.objects.filter(franqueado=user).values_list('id', flat=True)
        available_videos = available_videos.filter(cliente_id__in=clientes_ids)

    # Conteúdos corporativos disponíveis
    conteudos_corporativos = ConteudoCorporativo.objects.filter(ativo=True)
    
    if request.method == 'POST':
        form = PlaylistForm(request.POST, instance=playlist, user=user)
        selected_videos = request.POST.getlist('selected_videos')
        # Filtrar valores vazios
        selected_videos = [v for v in selected_videos if v and v.strip()]
        
        if form.is_valid():
            playlist = form.save(commit=False)
            # Definir franqueado baseado no município selecionado ou no usuário
            if user.is_franchisee():
                playlist.franqueado = user
            elif playlist.municipio:
                playlist.franqueado = playlist.municipio.franqueado
            playlist.save()
            
            # Remover itens antigos
            playlist.items.all().delete()
            
            # Adicionar itens selecionados (vídeos + corporativos)
            if selected_videos:
                for order, item_id in enumerate(selected_videos, 1):
                    if item_id.startswith('corp_'):
                        corp_id = int(item_id.replace('corp_', ''))
                        PlaylistItem.objects.create(
                            playlist=playlist,
                            conteudo_corporativo_id=corp_id,
                            ordem=order
                        )
                    else:
                        PlaylistItem.objects.create(
                            playlist=playlist,
                            video_id=int(item_id),
                            ordem=order
                        )
                messages.success(request, f'Playlist atualizada com sucesso com {len(selected_videos)} item(ns)!')
            else:
                messages.success(request, 'Playlist atualizada com sucesso!')
            
            return redirect('playlist_detail', pk=playlist.pk)
    else:
        form = PlaylistForm(instance=playlist, user=user)
    
    # Obter itens atuais da playlist
    current_items = playlist.items.all().select_related('video', 'conteudo_corporativo').order_by('ordem')
    
    context = {
        'form': form,
        'playlist': playlist,
        'available_videos': available_videos,
        'conteudos_corporativos': conteudos_corporativos,
        'current_items': current_items,
    }
    
    return render(request, 'playlists/playlist_form.html', context)


@login_required
def playlist_delete_view(request, pk):
    """Deletar playlist"""
    playlist = get_object_or_404(Playlist, pk=pk)
    user = request.user
    
    # Verificar permissões
    if user.is_franchisee() and playlist.franqueado != user:
        messages.error(request, 'Você não tem permissão para deletar esta playlist.')
        return redirect('playlist_list')
    
    if request.method == 'POST':
        playlist_nome = playlist.nome
        playlist.delete()
        messages.success(request, f'Playlist "{playlist_nome}" deletada com sucesso!')
        return redirect('playlist_list')
    
    context = {
        'playlist': playlist,
    }
    
    return render(request, 'playlists/playlist_confirm_delete.html', context)


# Dispositivo Views
@login_required
def dispositivo_list_view(request):
    """Lista de dispositivos TV"""
    user = request.user
    dispositivos = DispositivoTV.objects.select_related(
        'municipio', 'municipio__franqueado', 'playlist_atual'
    ).annotate(
        total_exibicoes=Count('logs_exibicao')
    )

    # Filtros
    search = request.GET.get('search', '')
    municipio_filter = request.GET.get('municipio', '')
    status_filter = request.GET.get('status', '')
    playlist_filter = request.GET.get('playlist', '')

    if search:
        dispositivos = dispositivos.filter(
            Q(nome__icontains=search) | Q(localizacao__icontains=search)
        )

    if municipio_filter:
        dispositivos = dispositivos.filter(municipio_id=municipio_filter)

    if status_filter:
        dispositivos = dispositivos.filter(ativo=(status_filter == 'ativo'))

    if playlist_filter:
        dispositivos = dispositivos.filter(playlist_atual_id=playlist_filter)

    # Controle de permissões
    if user.is_franchisee():
        dispositivos = dispositivos.filter(
            Q(municipio__franqueado=user) | Q(franqueado=user)
        )
    elif not user.is_owner():
        messages.error(request, 'Você não tem permissão para ver dispositivos.')
        return redirect('dashboard')

    # Estatísticas - baseadas no queryset filtrado
    # Total de exibições: contar logs dos dispositivos filtrados
    total_exibicoes = LogExibicao.objects.filter(
        dispositivo__in=dispositivos
    ).count()

    # Status real de conexão (requer avaliação Python por dispositivo)
    all_dispositivos_list = list(dispositivos.prefetch_related('agendamentos'))
    count_transmitindo = 0
    count_fora_horario = 0
    count_desconectado = 0
    for _d in all_dispositivos_list:
        _st = _d.status_conexao()
        if _st == 'transmitindo':
            count_transmitindo += 1
        elif _st == 'fora_horario':
            count_fora_horario += 1
        else:
            count_desconectado += 1

    context = {
        'dispositivos_ativos': dispositivos.filter(ativo=True).count(),
        'dispositivos_inativos': dispositivos.filter(ativo=False).count(),
        'total_exibicoes': total_exibicoes,
        'tempo_total_exibicao': '0h',
        'count_transmitindo': count_transmitindo,
        'count_fora_horario': count_fora_horario,
        'count_desconectado': count_desconectado,
    }

    # Paginação
    paginator = Paginator(dispositivos.order_by('-created_at'), 9)
    page_number = request.GET.get('page')
    page_obj = paginator.get_page(page_number)

    # Municipios e playlists para filtro (scoped)
    if user.is_owner():
        filter_municipios = Municipio.objects.only('id', 'nome', 'estado')
    else:
        # Municípios das TVs que o franqueado gerencia (pelo município ou atribuição direta)
        municipio_ids = dispositivos.values_list('municipio_id', flat=True)
        filter_municipios = Municipio.objects.filter(id__in=municipio_ids).only('id', 'nome', 'estado')
    
    context.update({
        'dispositivos': page_obj,
        'municipios': filter_municipios,
        'playlists': Playlist.objects.filter(ativa=True).only('id', 'nome'),
    })

    return render(request, 'dispositivos/dispositivo_list.html', context)


@login_required
def dispositivo_create_view(request):
    """Criar novo dispositivo TV"""
    user = request.user

    # Verificar permissões
    if not user.is_franchisee() and not user.is_owner():
        messages.error(request, 'Apenas franqueados podem cadastrar dispositivos.')
        return redirect('dispositivo_list')

    if request.method == 'POST':
        form = DispositivoTVForm(request.POST)
        if user.is_franchisee() and 'franqueado' in form.fields:
            del form.fields['franqueado']
        if form.is_valid():
            dispositivo = form.save(commit=False)
            # Gerar identificador único se não existir
            if not dispositivo.identificador_unico:
                import uuid
                dispositivo.identificador_unico = str(uuid.uuid4())
            dispositivo.save()
            messages.success(request, 'Dispositivo cadastrado com sucesso!')
            return redirect('dispositivo_list')
    else:
        form = DispositivoTVForm()
        if user.is_franchisee() and 'franqueado' in form.fields:
            del form.fields['franqueado']

    return render(request, 'dispositivos/dispositivo_form.html', {'form': form, 'is_owner': user.is_owner()})


@login_required
def dispositivo_detail_view(request, pk):
    """Detalhes do dispositivo TV"""
    dispositivo = get_object_or_404(DispositivoTV, pk=pk)
    user = request.user
    
    # Verificar permissões
    if user.is_franchisee() and dispositivo.municipio.franqueado != user and dispositivo.franqueado != user:
        messages.error(request, 'Você não tem permissão para ver este dispositivo.')
        return redirect('dispositivo_list')
    
    # Contar agendamentos ativos
    agendamentos_ativos = dispositivo.agendamentos.filter(ativo=True).select_related('playlist')
    agendamentos_ativos_count = agendamentos_ativos.count()
    
    # Horário de funcionamento do dispositivo
    no_horario = dispositivo.esta_no_horario_exibicao()
    horarios_funcionamento = dispositivo.horarios_funcionamento.all()

    # Buscar logs recentes — vídeos e webview unidos e ordenados por data
    logs_video = list(
        dispositivo.logs_exibicao.select_related('video', 'video__cliente', 'playlist')
        .filter(data_hora_fim__isnull=False)  # exclui registros ghost sem fim
        .order_by('-data_hora_inicio')[:30]
    )
    logs_wv = list(
        dispositivo.logs_webview.select_related('conteudo_corporativo', 'playlist')
        .order_by('-data_hora_inicio')[:30]
    )
    logs_recentes = sorted(
        logs_video + logs_wv,
        key=lambda x: x.data_hora_inicio,
        reverse=True
    )[:20]

    context = {
        'dispositivo': dispositivo,
        'agendamentos_ativos_count': agendamentos_ativos_count,
        'no_horario': no_horario,
        'horarios_funcionamento': horarios_funcionamento,
        'logs_recentes': logs_recentes,
        'is_owner': user.is_owner(),
    }
    
    return render(request, 'dispositivos/dispositivo_detail.html', context)


@login_required
def dispositivo_tv_preview_view(request, pk):
    """Preview do que está passando na TV — para mostrar ao cliente no celular/tablet"""
    dispositivo = get_object_or_404(DispositivoTV, pk=pk)
    user = request.user

    if user.is_franchisee() and dispositivo.municipio.franqueado != user and dispositivo.franqueado != user:
        messages.error(request, 'Sem permissão para visualizar este dispositivo.')
        return redirect('dispositivo_list')

    # Reutiliza a mesma lógica da TVAPIView para montar a fila de vídeos
    agendamentos_ativos = dispositivo.get_agendamentos_ativos_por_horario()
    videos = []

    if agendamentos_ativos:
        pairs = []
        for ag in agendamentos_ativos:
            playlist = ag.playlist
            if not playlist.ativa:
                continue
            pl_serializer = PlaylistTVSerializer(
                playlist,
                context={'request': request, 'dispositivo_id': dispositivo.id}
            )
            pl_videos = pl_serializer.data.get('videos', [])
            pairs.append((pl_videos, ag.percentual))

        has_varied = len(pairs) > 1 and any(pct != 100 for _, pct in pairs)
        if has_varied:
            videos = _distribuir_por_percentual(pairs)
        else:
            videos = [v for vids, _ in pairs for v in vids]

    import json as _json
    videos_json = _json.dumps(videos)

    return render(request, 'dispositivos/dispositivo_tv_preview.html', {
        'dispositivo': dispositivo,
        'videos_json': videos_json,
        'total_videos': len(videos),
    })


@login_required
def dispositivo_update_view(request, pk):
    """Editar dispositivo TV"""
    dispositivo = get_object_or_404(DispositivoTV, pk=pk)
    user = request.user
    
    # Verificar permissões
    if user.is_franchisee() and dispositivo.municipio.franqueado != user and dispositivo.franqueado != user:
        messages.error(request, 'Você não tem permissão para editar este dispositivo.')
        return redirect('dispositivo_list')
    
    if request.method == 'POST':
        form = DispositivoTVForm(request.POST, instance=dispositivo)
        # Franqueado não pode alterar o campo franqueado — apenas o dono pode
        if user.is_franchisee() and 'franqueado' in form.fields:
            del form.fields['franqueado']
        if form.is_valid():
            dispositivo = form.save()
            messages.success(request, 'Dispositivo atualizado com sucesso!')
            return redirect('dispositivo_detail', pk=dispositivo.pk)
    else:
        form = DispositivoTVForm(instance=dispositivo)
        if user.is_franchisee() and 'franqueado' in form.fields:
            del form.fields['franqueado']
    
    context = {
        'form': form,
        'dispositivo': dispositivo,
        'is_owner': user.is_owner(),
    }
    
    return render(request, 'dispositivos/dispositivo_form.html', context)


@login_required
def dispositivo_delete_view(request, pk):
    """Deletar dispositivo TV"""
    dispositivo = get_object_or_404(DispositivoTV, pk=pk)
    user = request.user
    
    # Verificar permissões
    if user.is_franchisee() and dispositivo.municipio.franqueado != user and dispositivo.franqueado != user:
        messages.error(request, 'Você não tem permissão para deletar este dispositivo.')
        return redirect('dispositivo_list')
    
    if request.method == 'POST':
        dispositivo_nome = dispositivo.nome
        dispositivo.delete()
        messages.success(request, f'Dispositivo "{dispositivo_nome}" deletado com sucesso!')
        return redirect('dispositivo_list')
    
    context = {
        'dispositivo': dispositivo,
    }
    
    return render(request, 'dispositivos/dispositivo_confirm_delete.html', context)


# Agendamento Views
@login_required
def agendamento_create_view(request, dispositivo_pk):
    """Criar novo agendamento de exibição"""
    from .forms import AgendamentoExibicaoForm
    from .models import AgendamentoExibicao
    
    dispositivo = get_object_or_404(DispositivoTV, pk=dispositivo_pk)
    user = request.user
    
    # Verificar permissões
    if user.is_franchisee() and dispositivo.municipio.franqueado != user and dispositivo.franqueado != user:
        messages.error(request, 'Você não tem permissão para criar agendamentos para este dispositivo.')
        return redirect('dispositivo_detail', pk=dispositivo_pk)
    
    if request.method == 'POST':
        form = AgendamentoExibicaoForm(request.POST, user=user)
        if form.is_valid():
            agendamento = form.save(commit=False)
            agendamento.dispositivo = dispositivo
            agendamento.save()
            messages.success(request, 'Agendamento criado com sucesso!')
            return redirect('dispositivo_detail', pk=dispositivo_pk)
    else:
        form = AgendamentoExibicaoForm(user=user)
    
    context = {
        'form': form,
        'dispositivo': dispositivo,
        'title': 'Vincular Playlist',
        'button_text': 'Vincular',
    }
    
    return render(request, 'agendamentos/agendamento_form.html', context)


@login_required
def agendamento_update_view(request, dispositivo_pk, pk):
    """Atualizar vínculo playlist-dispositivo"""
    from .forms import AgendamentoExibicaoForm
    from .models import AgendamentoExibicao
    
    dispositivo = get_object_or_404(DispositivoTV, pk=dispositivo_pk)
    agendamento = get_object_or_404(AgendamentoExibicao, pk=pk, dispositivo=dispositivo)
    user = request.user
    
    # Verificar permissões
    if user.is_franchisee() and dispositivo.municipio.franqueado != user and dispositivo.franqueado != user:
        messages.error(request, 'Você não tem permissão para editar este agendamento.')
        return redirect('dispositivo_detail', pk=dispositivo_pk)
    if user.is_franchisee() and agendamento.playlist and agendamento.playlist.franqueado != user:
        messages.error(request, 'Você não pode editar um agendamento de playlist adicionada pelo administrador.')
        return redirect('dispositivo_detail', pk=dispositivo_pk)
    
    if request.method == 'POST':
        form = AgendamentoExibicaoForm(request.POST, instance=agendamento, user=user)
        if form.is_valid():
            form.save()
            messages.success(request, 'Playlist atualizada com sucesso!')
            return redirect('dispositivo_detail', pk=dispositivo_pk)
    else:
        form = AgendamentoExibicaoForm(instance=agendamento, user=user)
    
    context = {
        'form': form,
        'dispositivo': dispositivo,
        'agendamento': agendamento,
        'title': 'Editar Playlist Vinculada',
        'button_text': 'Salvar Alterações',
    }
    
    return render(request, 'agendamentos/agendamento_form.html', context)


@login_required
def agendamento_delete_view(request, dispositivo_pk, pk):
    """Deletar agendamento de exibição"""
    from .models import AgendamentoExibicao
    
    dispositivo = get_object_or_404(DispositivoTV, pk=dispositivo_pk)
    agendamento = get_object_or_404(AgendamentoExibicao, pk=pk, dispositivo=dispositivo)
    user = request.user
    
    # Verificar permissões
    if user.is_franchisee() and dispositivo.municipio.franqueado != user and dispositivo.franqueado != user:
        messages.error(request, 'Você não tem permissão para deletar este agendamento.')
        return redirect('dispositivo_detail', pk=dispositivo_pk)
    if user.is_franchisee() and agendamento.playlist and agendamento.playlist.franqueado != user:
        messages.error(request, 'Você não pode desvincular uma playlist adicionada pelo administrador.')
        return redirect('dispositivo_detail', pk=dispositivo_pk)
    
    if request.method == 'POST':
        playlist_nome = agendamento.playlist.nome if agendamento.playlist else agendamento.nome
        agendamento.delete()
        messages.success(request, f'Playlist "{playlist_nome}" desvinculada com sucesso!')
        return redirect('dispositivo_detail', pk=dispositivo_pk)
    
    context = {
        'agendamento': agendamento,
        'dispositivo': dispositivo,
    }
    
    return render(request, 'agendamentos/agendamento_confirm_delete.html', context)


# ===== Horário de Funcionamento CRUD =====

@login_required
def horario_create_view(request, dispositivo_pk):
    """Criar horário de funcionamento para um dispositivo"""
    dispositivo = get_object_or_404(DispositivoTV, pk=dispositivo_pk)
    user = request.user

    if user.is_franchisee() and dispositivo.municipio.franqueado != user and dispositivo.franqueado != user:
        messages.error(request, 'Você não tem permissão para este dispositivo.')
        return redirect('dispositivo_detail', pk=dispositivo_pk)

    if request.method == 'POST':
        form = HorarioFuncionamentoForm(request.POST)
        if form.is_valid():
            horario = form.save(commit=False)
            horario.dispositivo = dispositivo
            horario.save()
            messages.success(request, f'Horário "{horario.nome}" criado com sucesso!')
            return redirect('dispositivo_detail', pk=dispositivo_pk)
    else:
        form = HorarioFuncionamentoForm()

    context = {
        'form': form,
        'dispositivo': dispositivo,
    }
    return render(request, 'horarios/horario_form.html', context)


@login_required
def horario_update_view(request, dispositivo_pk, pk):
    """Editar horário de funcionamento"""
    dispositivo = get_object_or_404(DispositivoTV, pk=dispositivo_pk)
    horario = get_object_or_404(HorarioFuncionamento, pk=pk, dispositivo=dispositivo)
    user = request.user

    if user.is_franchisee() and dispositivo.municipio.franqueado != user and dispositivo.franqueado != user:
        messages.error(request, 'Você não tem permissão para este dispositivo.')
        return redirect('dispositivo_detail', pk=dispositivo_pk)

    if request.method == 'POST':
        form = HorarioFuncionamentoForm(request.POST, instance=horario)
        if form.is_valid():
            form.save()
            messages.success(request, f'Horário "{horario.nome}" atualizado com sucesso!')
            return redirect('dispositivo_detail', pk=dispositivo_pk)
    else:
        form = HorarioFuncionamentoForm(instance=horario)

    context = {
        'form': form,
        'dispositivo': dispositivo,
        'horario': horario,
    }
    return render(request, 'horarios/horario_form.html', context)


@login_required
def horario_delete_view(request, dispositivo_pk, pk):
    """Deletar horário de funcionamento"""
    dispositivo = get_object_or_404(DispositivoTV, pk=dispositivo_pk)
    horario = get_object_or_404(HorarioFuncionamento, pk=pk, dispositivo=dispositivo)
    user = request.user

    if user.is_franchisee() and dispositivo.municipio.franqueado != user and dispositivo.franqueado != user:
        messages.error(request, 'Você não tem permissão para este dispositivo.')
        return redirect('dispositivo_detail', pk=dispositivo_pk)

    if request.method == 'POST':
        nome = horario.nome
        horario.delete()
        messages.success(request, f'Horário "{nome}" excluído com sucesso!')
        return redirect('dispositivo_detail', pk=dispositivo_pk)

    context = {
        'horario': horario,
        'dispositivo': dispositivo,
    }
    return render(request, 'horarios/horario_confirm_delete.html', context)


# User/Franchisee Views
@login_required
def user_list_view(request):
    """Lista de usuários"""
    user = request.user
    
    # Apenas proprietários podem ver a lista de usuários
    if not user.is_owner():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    users = User.objects.all().order_by('-created_at')
    
    # Filtros
    search = request.GET.get('search', '')
    role_filter = request.GET.get('role', '')
    
    if search:
        users = users.filter(
            Q(username__icontains=search) |
            Q(email__icontains=search) |
            Q(first_name__icontains=search) |
            Q(last_name__icontains=search)
        )
    
    if role_filter:
        users = users.filter(role=role_filter)
    
    context = {
        'users': users,
        'search': search,
        'role_filter': role_filter,
        'roles': User.ROLE_CHOICES,
    }
    
    return render(request, 'users/user_list.html', context)


@login_required
def user_create_view(request):
    """Criar novo usuário"""
    user = request.user
    
    # Apenas proprietários podem criar usuários
    if not user.is_owner():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    if request.method == 'POST':
        username = request.POST.get('username')
        email = request.POST.get('email')
        password = request.POST.get('password')
        password_confirm = request.POST.get('password_confirm')
        first_name = request.POST.get('first_name', '')
        last_name = request.POST.get('last_name', '')
        role = request.POST.get('role')
        
        # Validações
        if not username or not password:
            messages.error(request, 'Username e senha são obrigatórios.')
        elif password != password_confirm:
            messages.error(request, 'As senhas não coincidem.')
        elif User.objects.filter(username=username).exists():
            messages.error(request, 'Este username já está em uso.')
        elif email and User.objects.filter(email=email).exists():
            messages.error(request, 'Este email já está em uso.')
        else:
            # Criar usuário
            new_user = User.objects.create_user(
                username=username,
                email=email,
                password=password,
                first_name=first_name,
                last_name=last_name,
                role=role
            )
            messages.success(request, f'Usuário {username} criado com sucesso!')
            return redirect('user_list')
    
    context = {
        'roles': User.ROLE_CHOICES,
    }
    
    return render(request, 'users/user_form.html', context)


@login_required
def user_update_view(request, pk):
    """Editar usuário existente"""
    user = request.user
    
    # Apenas proprietários podem editar usuários
    if not user.is_owner():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    edit_user = get_object_or_404(User, pk=pk)
    
    if request.method == 'POST':
        email = request.POST.get('email')
        first_name = request.POST.get('first_name', '')
        last_name = request.POST.get('last_name', '')
        role = request.POST.get('role')
        is_active = request.POST.get('is_active') == 'on'
        password = request.POST.get('password')
        password_confirm = request.POST.get('password_confirm')
        
        # Validações
        if email and User.objects.filter(email=email).exclude(pk=pk).exists():
            messages.error(request, 'Este email já está em uso.')
        elif password and password != password_confirm:
            messages.error(request, 'As senhas não coincidem.')
        else:
            # Atualizar usuário
            edit_user.email = email
            edit_user.first_name = first_name
            edit_user.last_name = last_name
            edit_user.role = role
            edit_user.is_active = is_active
            
            if password:
                edit_user.set_password(password)
            
            edit_user.save()
            messages.success(request, f'Usuário {edit_user.username} atualizado com sucesso!')
            return redirect('user_list')
    
    context = {
        'edit_user': edit_user,
        'roles': User.ROLE_CHOICES,
    }
    
    return render(request, 'users/user_form.html', context)


@login_required
def user_delete_view(request, pk):
    """Deletar usuário"""
    user = request.user
    
    # Apenas proprietários podem deletar usuários
    if not user.is_owner():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    delete_user = get_object_or_404(User, pk=pk)
    
    # Não pode deletar a si mesmo
    if delete_user == user:
        messages.error(request, 'Você não pode deletar sua própria conta.')
        return redirect('user_list')
    
    if request.method == 'POST':
        username = delete_user.username
        delete_user.delete()
        messages.success(request, f'Usuário {username} deletado com sucesso!')
        return redirect('user_list')
    
    context = {
        'delete_user': delete_user,
    }
    
    return render(request, 'users/user_confirm_delete.html', context)


# Municipio Views
@login_required
def municipio_list_view(request):
    """Lista de municípios"""
    user = request.user
    
    # Apenas proprietários e franqueados podem ver municípios
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    municipios = Municipio.objects.all().select_related('franqueado').order_by('nome')
    
    # Franqueados veem apenas seus municípios
    if user.is_franchisee():
        municipios = municipios.filter(franqueado=user)
    
    # Filtros
    search = request.GET.get('search', '')
    franqueado_filter = request.GET.get('franqueado', '')
    
    if search:
        municipios = municipios.filter(
            Q(nome__icontains=search) |
            Q(estado__icontains=search)
        )
    
    if franqueado_filter and user.is_owner():
        municipios = municipios.filter(franqueado_id=franqueado_filter)
    
    # Lista de franqueados para filtro (apenas para owner)
    franqueados = User.objects.filter(role='FRANCHISEE').order_by('username') if user.is_owner() else []
    
    context = {
        'municipios': municipios,
        'search': search,
        'franqueado_filter': franqueado_filter,
        'franqueados': franqueados,
    }
    
    return render(request, 'municipios/municipio_list.html', context)


def _parse_coord(value):
    """Converte texto de coordenada (latitude/longitude) para Decimal com 6 casas.
    Suporta valores com muitas casas decimais vindos do Google Maps."""
    from decimal import Decimal, InvalidOperation, ROUND_HALF_UP
    if not value or not str(value).strip():
        return None
    try:
        # Converte via float para lidar com notação científica, depois quantiza
        d = Decimal(str(round(float(value), 6)))
        return d.quantize(Decimal('0.000001'), rounding=ROUND_HALF_UP)
    except (ValueError, InvalidOperation):
        return None


@login_required
def municipio_create_view(request):
    """Criar novo município"""
    user = request.user
    
    # Apenas proprietários podem criar municípios
    if not user.is_owner():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    if request.method == 'POST':
        nome = request.POST.get('nome')
        estado = request.POST.get('estado')
        franqueado_id = request.POST.get('franqueado')
        latitude = _parse_coord(request.POST.get('latitude'))
        longitude = _parse_coord(request.POST.get('longitude'))
        
        # Validações
        if not nome or not estado:
            messages.error(request, 'Nome e Estado são obrigatórios.')
        elif not franqueado_id:
            messages.error(request, 'Franqueado é obrigatório.')
        else:
            # Criar município
            municipio = Municipio.objects.create(
                nome=nome,
                estado=estado,
                franqueado_id=franqueado_id,
                latitude=latitude,
                longitude=longitude,
            )
            messages.success(request, f'Município {nome}/{estado} criado com sucesso!')
            return redirect('municipio_list')
    
    # Lista de franqueados para atribuição
    franqueados = User.objects.filter(role='FRANCHISEE').order_by('username')
    
    # Lista de UFs brasileiras
    ufs = [
        'AC', 'AL', 'AP', 'AM', 'BA', 'CE', 'DF', 'ES', 'GO', 'MA',
        'MT', 'MS', 'MG', 'PA', 'PB', 'PR', 'PE', 'PI', 'RJ', 'RN',
        'RS', 'RO', 'RR', 'SC', 'SP', 'SE', 'TO'
    ]
    
    context = {
        'franqueados': franqueados,
        'ufs': ufs,
    }
    
    return render(request, 'municipios/municipio_form.html', context)


@login_required
def municipio_update_view(request, pk):
    """Editar município existente"""
    user = request.user
    
    # Apenas proprietários podem editar municípios
    if not user.is_owner():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    municipio = get_object_or_404(Municipio, pk=pk)
    
    if request.method == 'POST':
        nome = request.POST.get('nome')
        estado = request.POST.get('estado')
        franqueado_id = request.POST.get('franqueado')
        latitude = _parse_coord(request.POST.get('latitude'))
        longitude = _parse_coord(request.POST.get('longitude'))
        
        # Validações
        if not nome or not estado:
            messages.error(request, 'Nome e Estado são obrigatórios.')
        elif not franqueado_id:
            messages.error(request, 'Franqueado é obrigatório.')
        else:
            # Atualizar município
            municipio.nome = nome
            municipio.estado = estado
            municipio.franqueado_id = franqueado_id
            municipio.latitude = latitude
            municipio.longitude = longitude
            municipio.save()
            messages.success(request, f'Município {nome}/{estado} atualizado com sucesso!')
            return redirect('municipio_list')
    
    # Lista de franqueados para atribuição
    franqueados = User.objects.filter(role='FRANCHISEE').order_by('username')
    
    # Lista de UFs brasileiras
    ufs = [
        'AC', 'AL', 'AP', 'AM', 'BA', 'CE', 'DF', 'ES', 'GO', 'MA',
        'MT', 'MS', 'MG', 'PA', 'PB', 'PR', 'PE', 'PI', 'RJ', 'RN',
        'RS', 'RO', 'RR', 'SC', 'SP', 'SE', 'TO'
    ]
    
    context = {
        'municipio': municipio,
        'franqueados': franqueados,
        'ufs': ufs,
    }
    
    return render(request, 'municipios/municipio_form.html', context)


@login_required
def municipio_delete_view(request, pk):
    """Deletar município"""
    user = request.user
    
    # Apenas proprietários podem deletar municípios
    if not user.is_owner():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    municipio = get_object_or_404(Municipio, pk=pk)
    
    if request.method == 'POST':
        nome = municipio.nome
        estado = municipio.estado
        municipio.delete()
        messages.success(request, f'Município {nome}/{estado} deletado com sucesso!')
        return redirect('municipio_list')
    
    # Contar clientes e dispositivos relacionados
    total_clientes = municipio.clientes.count()
    total_dispositivos = municipio.dispositivos.count()
    
    context = {
        'municipio': municipio,
        'total_clientes': total_clientes,
        'total_dispositivos': total_dispositivos,
    }
    
    return render(request, 'municipios/municipio_confirm_delete.html', context)


# Segmento Views
@login_required
def segmento_list_view(request):
    """Lista de segmentos"""
    user = request.user
    
    # Apenas proprietários e franqueados podem ver segmentos
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    segmentos = Segmento.objects.all().order_by('nome')
    
    # Filtro de busca
    search = request.GET.get('search', '')
    if search:
        segmentos = segmentos.filter(nome__icontains=search)
    
    # Filtro ativo/inativo
    status_filter = request.GET.get('status', '')
    if status_filter == 'ativo':
        segmentos = segmentos.filter(ativo=True)
    elif status_filter == 'inativo':
        segmentos = segmentos.filter(ativo=False)
    
    # Paginação
    paginator = Paginator(segmentos, 15)
    page_number = request.GET.get('page')
    segmentos_page = paginator.get_page(page_number)
    
    context = {
        'segmentos': segmentos_page,
        'total_segmentos': segmentos.count(),
        'search': search,
        'status_filter': status_filter,
    }
    
    return render(request, 'segmentos/segmento_list.html', context)


@login_required
def segmento_create_view(request):
    """Criar novo segmento"""
    user = request.user
    
    print(f"DEBUG: segmento_create_view chamada. Método: {request.method}, User: {user.username}")
    
    # Proprietários e franqueados podem criar segmentos
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    if request.method == 'POST':
        print(f"DEBUG: POST data: {request.POST}")
        print(f"DEBUG: Is AJAX: {request.headers.get('X-Requested-With') == 'XMLHttpRequest'}")
        
        form = SegmentoForm(request.POST)
        if form.is_valid():
            segmento = form.save()
            print(f"DEBUG: Segmento criado com sucesso: ID={segmento.id}, Nome={segmento.nome}")
            messages.success(request, f'Segmento {segmento.nome} criado com sucesso!')
            
            # Se for AJAX, retornar JSON
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                from django.http import JsonResponse
                return JsonResponse({
                    'success': True,
                    'segmento': {
                        'id': segmento.id,
                        'nome': segmento.nome
                    }
                })
            
            return redirect('segmento_list')
        else:
            print(f"DEBUG: Formulário inválido. Erros: {form.errors}")
            
            # Se for AJAX, retornar erros
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                from django.http import JsonResponse
                return JsonResponse({
                    'success': False,
                    'errors': form.errors
                }, status=400)
            
            messages.error(request, 'Erro ao criar segmento. Verifique os dados.')
    else:
        form = SegmentoForm()
    
    context = {
        'form': form,
        'title': 'Novo Segmento',
        'button_text': 'Criar Segmento'
    }
    
    return render(request, 'segmentos/segmento_form.html', context)


@login_required
def segmento_update_view(request, pk):
    """Editar segmento"""
    user = request.user
    
    # Apenas proprietários podem editar segmentos
    if not user.is_owner():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    segmento = get_object_or_404(Segmento, pk=pk)
    
    if request.method == 'POST':
        form = SegmentoForm(request.POST, instance=segmento)
        if form.is_valid():
            segmento = form.save()
            messages.success(request, f'Segmento {segmento.nome} atualizado com sucesso!')
            return redirect('segmento_list')
        else:
            messages.error(request, 'Erro ao atualizar segmento. Verifique os dados.')
    else:
        form = SegmentoForm(instance=segmento)
    
    context = {
        'form': form,
        'segmento': segmento,
        'title': f'Editar Segmento: {segmento.nome}',
        'button_text': 'Salvar Alterações'
    }
    
    return render(request, 'segmentos/segmento_form.html', context)


@login_required
def segmento_delete_view(request, pk):
    """Deletar segmento"""
    user = request.user
    
    # Apenas proprietários podem deletar segmentos
    if not user.is_owner():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    segmento = get_object_or_404(Segmento, pk=pk)
    
    if request.method == 'POST':
        nome = segmento.nome
        
        # Verificar se há clientes usando este segmento
        total_clientes = segmento.clientes.count()
        if total_clientes > 0:
            messages.error(request, f'Não é possível deletar o segmento {nome} pois existem {total_clientes} cliente(s) vinculado(s).')
            return redirect('segmento_list')
        
        segmento.delete()
        messages.success(request, f'Segmento {nome} deletado com sucesso!')
        return redirect('segmento_list')
    
    # Contar clientes relacionados
    total_clientes = segmento.clientes.count()
    
    context = {
        'segmento': segmento,
        'total_clientes': total_clientes,
    }
    
    return render(request, 'segmentos/segmento_confirm_delete.html', context)


# Cliente Views
@login_required
def cliente_list_view(request):
    """Lista de clientes"""
    user = request.user
    
    # Apenas proprietários e franqueados podem ver clientes
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    clientes = Cliente.objects.select_related(
        'user', 'franqueado', 'segmento'
    ).prefetch_related('municipios').order_by('-created_at')
    
    # Franqueados veem apenas seus clientes
    if user.is_franchisee():
        clientes = clientes.filter(franqueado=user)
    
    # Filtros
    search = request.GET.get('search', '')
    franqueado_filter = request.GET.get('franqueado', '')
    municipio_filter = request.GET.get('municipio', '')
    
    if search:
        clientes = clientes.filter(
            Q(empresa__icontains=search) |
            Q(user__username__icontains=search) |
            Q(user__email__icontains=search) |
            Q(user__first_name__icontains=search) |
            Q(user__last_name__icontains=search)
        )
    
    if franqueado_filter and user.is_owner():
        clientes = clientes.filter(franqueado_id=franqueado_filter)
    
    if municipio_filter:
        clientes = clientes.filter(municipios__id=municipio_filter).distinct()
    
    # Lista de franqueados e municípios para filtros
    franqueados = User.objects.filter(role='FRANCHISEE').order_by('username') if user.is_owner() else []
    
    if user.is_owner():
        municipios_filter = Municipio.objects.all().order_by('nome')
    else:
        municipios_filter = Municipio.objects.filter(franqueado=user).order_by('nome')
    
    context = {
        'clientes': clientes,
        'search': search,
        'franqueado_filter': franqueado_filter,
        'municipio_filter': municipio_filter,
        'franqueados': franqueados,
        'municipios_filter': municipios_filter,
    }
    
    return render(request, 'clientes/cliente_list.html', context)


@login_required
def cliente_create_view(request):
    """Criar novo cliente"""
    user = request.user
    
    # Apenas proprietários e franqueados podem criar clientes
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    if request.method == 'POST':
        # Dados do usuário
        username = request.POST.get('username')
        email = request.POST.get('email')
        password = request.POST.get('password')
        password_confirm = request.POST.get('password_confirm')
        first_name = request.POST.get('first_name', '')
        last_name = request.POST.get('last_name', '')
        
        # Dados do cliente
        empresa = request.POST.get('empresa')
        segmento_id = request.POST.get('segmento')
        franqueado_id = request.POST.get('franqueado') if user.is_owner() else user.id
        municipios_ids = request.POST.getlist('municipios')
        observacoes = request.POST.get('observacoes', '')
        contrato = request.FILES.get('contrato')
        
        # Validações
        if not username or not password or not empresa or not segmento_id:
            messages.error(request, 'Username, senha, nome da empresa e segmento são obrigatórios.')
        elif password != password_confirm:
            messages.error(request, 'As senhas não coincidem.')
        elif User.objects.filter(username=username).exists():
            messages.error(request, 'Este username já está em uso.')
        elif email and User.objects.filter(email=email).exists():
            messages.error(request, 'Este email já está em uso.')
        elif not municipios_ids:
            messages.error(request, 'Selecione pelo menos um município.')
        else:
            # Validar regra: uma marca por segmento por cidade
            for municipio_id in municipios_ids:
                existe = Cliente.objects.filter(
                    segmento_id=segmento_id,
                    municipios__id=municipio_id
                )
                if existe.exists():
                    cliente_existente = existe.first()
                    municipio = Municipio.objects.get(id=municipio_id)
                    messages.error(request, 
                        f'Já existe o cliente "{cliente_existente.empresa}" '
                        f'no segmento "{cliente_existente.segmento.nome}" '
                        f'no município {municipio.nome}/{municipio.estado}. '
                        'Apenas uma marca por segmento por cidade é permitida.'
                    )
                    # Buscar segmentos para renderizar novamente
                    segmentos = Segmento.objects.filter(ativo=True).order_by('nome')
                    
                    if user.is_owner():
                        franqueados = User.objects.filter(role='FRANCHISEE').order_by('username')
                        municipios = Municipio.objects.all().select_related('franqueado').order_by('estado', 'nome')
                    else:
                        franqueados = []
                        municipios = Municipio.objects.filter(franqueado=user).order_by('estado', 'nome')
                    
                    return render(request, 'clientes/cliente_form.html', {
                        'franqueados': franqueados,
                        'municipios': municipios,
                        'segmentos': segmentos,
                    })
            
            # Criar usuário
            new_user = User.objects.create_user(
                username=username,
                email=email,
                password=password,
                first_name=first_name,
                last_name=last_name,
                role='CLIENT'
            )
            
            # Criar cliente
            cliente = Cliente.objects.create(
                user=new_user,
                empresa=empresa,
                segmento_id=segmento_id,
                franqueado_id=franqueado_id,
                observacoes=observacoes,
                contrato=contrato
            )
            
            # Vincular municípios
            cliente.municipios.set(municipios_ids)
            
            messages.success(request, f'Cliente {empresa} criado com sucesso!')
            return redirect('cliente_list')
    
    # Lista de franqueados e municípios para seleção
    if user.is_owner():
        franqueados = User.objects.filter(role='FRANCHISEE').order_by('username')
        municipios = Municipio.objects.all().select_related('franqueado').order_by('estado', 'nome')
    else:
        franqueados = []
        municipios = Municipio.objects.filter(franqueado=user).order_by('estado', 'nome')
    
    # Buscar segmentos ativos
    segmentos = Segmento.objects.filter(ativo=True).order_by('nome')
    
    context = {
        'franqueados': franqueados,
        'municipios': municipios,
        'segmentos': segmentos,
    }
    
    return render(request, 'clientes/cliente_form.html', context)


@login_required
def cliente_update_view(request, pk):
    """Editar cliente existente"""
    user = request.user
    
    # Apenas proprietários e franqueados podem editar clientes
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    cliente = get_object_or_404(Cliente, pk=pk)
    
    # Franqueados só podem editar seus próprios clientes
    if user.is_franchisee() and cliente.franqueado != user:
        messages.error(request, 'Acesso negado.')
        return redirect('cliente_list')
    
    if request.method == 'POST':
        # Dados do usuário
        email = request.POST.get('email')
        first_name = request.POST.get('first_name', '')
        last_name = request.POST.get('last_name', '')
        password = request.POST.get('password')
        password_confirm = request.POST.get('password_confirm')
        
        # Dados do cliente
        empresa = request.POST.get('empresa')
        segmento_id = request.POST.get('segmento')
        franqueado_id = request.POST.get('franqueado') if user.is_owner() else cliente.franqueado_id
        municipios_ids = request.POST.getlist('municipios')
        observacoes = request.POST.get('observacoes', '')
        is_active = request.POST.get('is_active') == 'on'
        contrato = request.FILES.get('contrato')
        remover_contrato = request.POST.get('remover_contrato') == 'on'
        
        # Validações
        if not empresa or not segmento_id:
            messages.error(request, 'Nome da empresa e segmento são obrigatórios.')
        elif email and User.objects.filter(email=email).exclude(pk=cliente.user.pk).exists():
            messages.error(request, 'Este email já está em uso.')
        elif password and password != password_confirm:
            messages.error(request, 'As senhas não coincidem.')
        elif not municipios_ids:
            messages.error(request, 'Selecione pelo menos um município.')
        else:
            # Validar regra: uma marca por segmento por cidade
            for municipio_id in municipios_ids:
                existe = Cliente.objects.filter(
                    segmento_id=segmento_id,
                    municipios__id=municipio_id
                ).exclude(pk=cliente.pk)
                
                if existe.exists():
                    cliente_existente = existe.first()
                    municipio = Municipio.objects.get(id=municipio_id)
                    messages.error(request, 
                        f'Já existe o cliente "{cliente_existente.empresa}" '
                        f'no segmento "{cliente_existente.segmento.nome}" '
                        f'no município {municipio.nome}/{municipio.estado}. '
                        'Apenas uma marca por segmento por cidade é permitida.'
                    )
                    # Buscar segmentos para renderizar novamente
                    segmentos = Segmento.objects.filter(ativo=True).order_by('nome')
                    
                    if user.is_owner():
                        franqueados = User.objects.filter(role='FRANCHISEE').order_by('username')
                        municipios = Municipio.objects.all().select_related('franqueado').order_by('estado', 'nome')
                    else:
                        franqueados = []
                        municipios = Municipio.objects.filter(franqueado=user).order_by('estado', 'nome')
                    
                    return render(request, 'clientes/cliente_form.html', {
                        'cliente': cliente,
                        'franqueados': franqueados,
                        'municipios': municipios,
                        'segmentos': segmentos,
                    })
            
            # Atualizar usuário
            cliente.user.email = email
            cliente.user.first_name = first_name
            cliente.user.last_name = last_name
            cliente.user.is_active = is_active
            
            if password:
                cliente.user.set_password(password)
            
            cliente.user.save()
            
            # Atualizar cliente
            cliente.empresa = empresa
            cliente.segmento_id = segmento_id
            cliente.franqueado_id = franqueado_id
            cliente.observacoes = observacoes
            cliente.ativo = is_active
            
            # Atualizar contrato
            if remover_contrato:
                if cliente.contrato:
                    cliente.contrato.delete()
                cliente.contrato = None
            elif contrato:
                # Deletar contrato antigo se existir
                if cliente.contrato:
                    cliente.contrato.delete()
                cliente.contrato = contrato
            
            cliente.save()
            
            # Atualizar municípios
            cliente.municipios.set(municipios_ids)
            
            messages.success(request, f'Cliente {empresa} atualizado com sucesso!')
            return redirect('cliente_list')
    
    # Lista de franqueados e municípios para seleção
    if user.is_owner():
        franqueados = User.objects.filter(role='FRANCHISEE').order_by('username')
        municipios = Municipio.objects.all().select_related('franqueado').order_by('estado', 'nome')
    else:
        franqueados = []
        municipios = Municipio.objects.filter(franqueado=user).order_by('estado', 'nome')
    
    # Buscar segmentos ativos
    segmentos = Segmento.objects.filter(ativo=True).order_by('nome')
    
    context = {
        'cliente': cliente,
        'franqueados': franqueados,
        'municipios': municipios,
        'segmentos': segmentos,
    }
    
    return render(request, 'clientes/cliente_form.html', context)


@login_required
def cliente_delete_view(request, pk):
    """Deletar cliente"""
    user = request.user
    
    # Apenas proprietários e franqueados podem deletar clientes
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    cliente = get_object_or_404(Cliente, pk=pk)
    
    # Franqueados só podem deletar seus próprios clientes
    if user.is_franchisee() and cliente.franqueado != user:
        messages.error(request, 'Acesso negado.')
        return redirect('cliente_list')
    
    if request.method == 'POST':
        empresa = cliente.empresa
        username = cliente.user.username
        
        # Deletar usuário e cliente (CASCADE)
        cliente.user.delete()
        
        messages.success(request, f'Cliente {empresa} ({username}) deletado com sucesso!')
        return redirect('cliente_list')
    
    # Contar vídeos relacionados
    total_videos = cliente.videos.count()
    
    context = {
        'cliente': cliente,
        'total_videos': total_videos,
    }
    
    return render(request, 'clientes/cliente_confirm_delete.html', context)


# =====================================
# APP MANAGEMENT VIEWS (OWNER ONLY)
# =====================================

@login_required
def app_upload_view(request):
    """Upload e gerenciamento de versões do aplicativo (OWNER ONLY)"""
    user = request.user
    
    # Apenas proprietários podem fazer upload
    if not user.is_owner():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    if request.method == 'POST':
        form = AppVersionForm(request.POST, request.FILES)
        if form.is_valid():
            app_version = form.save(commit=False)
            app_version.uploaded_by = user
            app_version.save()
            messages.success(request, f'Versão {app_version.versao} enviada com sucesso!')
            return redirect('app_upload')
    else:
        form = AppVersionForm()
    
    # Listar todas as versões
    versions = AppVersion.objects.all()
    
    context = {
        'form': form,
        'versions': versions,
    }
    
    return render(request, 'app/app_upload.html', context)


@login_required
def app_version_toggle_view(request, pk):
    """Ativar/desativar uma versão do app (OWNER ONLY)"""
    user = request.user
    
    if not user.is_owner():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    app_version = get_object_or_404(AppVersion, pk=pk)
    app_version.ativo = not app_version.ativo
    app_version.save()
    
    status = 'ativada' if app_version.ativo else 'desativada'
    messages.success(request, f'Versão {app_version.versao} {status} com sucesso!')
    
    return redirect('app_upload')


@login_required
def app_version_delete_view(request, pk):
    """Deletar uma versão do app (OWNER ONLY)"""
    user = request.user
    
    if not user.is_owner():
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')
    
    app_version = get_object_or_404(AppVersion, pk=pk)
    
    if request.method == 'POST':
        versao = app_version.versao
        # Deletar arquivo do storage
        if app_version.arquivo_apk:
            app_version.arquivo_apk.delete()
        app_version.delete()
        messages.success(request, f'Versão {versao} deletada com sucesso!')
        return redirect('app_upload')
    
    return redirect('app_upload')


def app_download_view(request):
    """Download público da versão ativa do aplicativo"""
    from django.http import FileResponse, Http404
    import os
    
    # Buscar versão ativa mais recente
    app_version = AppVersion.get_versao_ativa()
    
    if not app_version:
        messages.error(request, 'Nenhuma versão disponível para download no momento.')
        return redirect('login')
    
    # Incrementar contador de downloads
    app_version.downloads += 1
    app_version.save(update_fields=['downloads'])
    
    # Verificar se arquivo existe
    if not app_version.arquivo_apk:
        raise Http404("Arquivo não encontrado")

    # Com storage remoto (R2): redirecionar direto para a URL do bucket
    # Com storage local: servir o arquivo
    try:
        app_version.arquivo_apk.path  # levanta NotImplementedError se R2
        local_path = app_version.arquivo_apk.path
        if not os.path.exists(local_path):
            raise Http404("Arquivo não encontrado")
        response = FileResponse(
            open(local_path, 'rb'),
            content_type='application/vnd.android.package-archive'
        )
        response['Content-Disposition'] = f'attachment; filename="MediaExpandTV-v{app_version.versao}.apk"'
        return response
    except NotImplementedError:
        # Storage remoto — redirecionar para URL do R2
        from django.shortcuts import redirect as _redirect
        return _redirect(app_version.arquivo_apk.url)


# ============================================================
# QR CODE TRACKING
# ============================================================

def qrcode_redirect_view(request, tracking_code):
    """
    View pública que rastreia cliques no QR Code e redireciona ao destino.
    URL: /r/<tracking_code>/
    """
    from django.shortcuts import redirect as django_redirect
    
    try:
        video = Video.objects.get(qrcode_tracking_code=tracking_code)
    except Video.DoesNotExist:
        raise Http404("Link não encontrado")
    
    if not video.qrcode_url_destino:
        raise Http404("Link de destino não configurado")
    
    # Registrar o clique
    QRCodeClick.objects.create(
        video=video,
        tracking_code=tracking_code,
        ip_address=get_client_ip(request),
        user_agent=request.META.get('HTTP_USER_AGENT', '')[:500],
        referer=request.META.get('HTTP_REFERER', '')[:500] or None,
    )
    
    # Redirecionar ao destino do cliente
    return django_redirect(video.qrcode_url_destino)


def get_client_ip(request):
    """Obtém o IP real do cliente, considerando proxies"""
    x_forwarded_for = request.META.get('HTTP_X_FORWARDED_FOR')
    if x_forwarded_for:
        return x_forwarded_for.split(',')[0].strip()
    return request.META.get('REMOTE_ADDR')


# ============================================================
# STREAMING DE MÍDIA COM SUPORTE A RANGE REQUESTS
# ============================================================

def serve_media_streaming(request, path):
    """
    Serve arquivos de mídia com suporte a HTTP Range Requests (RFC 7233).
    Permite streaming progressivo de vídeos grandes sem carregar tudo na memória.
    
    - Suporta Range requests (HTTP 206 Partial Content)
    - Streaming em chunks de 8MB para não sobrecarregar memória
    - Headers corretos para players de vídeo (Accept-Ranges, Content-Range)
    - Cache de 7 dias com ETag + Last-Modified → 304 Not Modified para re-requisições
    """
    import os
    import mimetypes
    from django.http import StreamingHttpResponse, HttpResponse, Http404
    from django.conf import settings
    
    # Construir caminho completo e validar contra path traversal
    full_path = os.path.join(str(settings.MEDIA_ROOT), path)
    full_path = os.path.normpath(full_path)
    media_root = os.path.normpath(str(settings.MEDIA_ROOT))

    # Com storage remoto (R2): o arquivo está no bucket, redirecionar direto
    from django.core.files.storage import default_storage
    try:
        default_storage.path(path)  # levanta NotImplementedError se S3/R2
    except NotImplementedError:
        from django.shortcuts import redirect as _sr
        return _sr(default_storage.url(path))

    if not full_path.startswith(media_root):
        raise Http404("Acesso negado")

    if not os.path.isfile(full_path):
        raise Http404("Arquivo não encontrado")
    
    # Detectar content type
    content_type, _ = mimetypes.guess_type(full_path)
    if not content_type:
        content_type = 'application/octet-stream'
    
    file_size = os.path.getsize(full_path)
    mtime = os.path.getmtime(full_path)

    # ETag baseado em tamanho + mtime (sem ler o arquivo)
    etag = f'"{file_size:x}-{int(mtime):x}"'

    # Last-Modified header
    from email.utils import formatdate
    last_modified = formatdate(mtime, usegmt=True)

    # ── Conditional request: If-None-Match ──
    if request.META.get('HTTP_IF_NONE_MATCH') == etag:
        resp = HttpResponse(status=304)
        resp['ETag'] = etag
        resp['Cache-Control'] = 'public, max-age=604800'
        return resp

    # ── Conditional request: If-Modified-Since ──
    ims = request.META.get('HTTP_IF_MODIFIED_SINCE')
    if ims:
        from email.utils import parsedate
        import time
        parsed = parsedate(ims)
        if parsed and time.mktime(parsed) >= mtime:
            resp = HttpResponse(status=304)
            resp['ETag'] = etag
            resp['Cache-Control'] = 'public, max-age=604800'
            return resp

    # Iterator que lê o arquivo em chunks de 8MB (não carrega tudo na memória)
    def file_iterator(file_path, start, end, chunk_size=8 * 1024 * 1024):
        with open(file_path, 'rb') as f:
            f.seek(start)
            remaining = end - start + 1
            while remaining > 0:
                data = f.read(min(chunk_size, remaining))
                if not data:
                    break
                remaining -= len(data)
                yield data
    
    # Verificar se é um Range request
    range_header = request.META.get('HTTP_RANGE', '').strip()
    
    if range_header.startswith('bytes='):
        range_val = range_header[6:].split('-')
        start = int(range_val[0]) if range_val[0] else 0
        end = int(range_val[1]) if range_val[1] else file_size - 1
        end = min(end, file_size - 1)
        
        # Validar range
        if start >= file_size:
            response = HttpResponse(status=416)  # Range Not Satisfiable
            response['Content-Range'] = f'bytes */{file_size}'
            return response
        
        content_length = end - start + 1
        
        response = StreamingHttpResponse(
            file_iterator(full_path, start, end),
            status=206,
            content_type=content_type,
        )
        response['Content-Range'] = f'bytes {start}-{end}/{file_size}'
        response['Content-Length'] = content_length
    else:
        # Sem Range header - resposta completa com streaming
        response = StreamingHttpResponse(
            file_iterator(full_path, 0, file_size - 1),
            content_type=content_type,
        )
        response['Content-Length'] = file_size
    
    response['Accept-Ranges'] = 'bytes'
    response['ETag'] = etag
    response['Last-Modified'] = last_modified
    response['Cache-Control'] = 'public, max-age=604800'  # 7 dias
    return response


# ══════════════════════════════════════════════════════════
#  VIEWS: Conteúdo Corporativo & Configuração de API
# ══════════════════════════════════════════════════════════

def _conteudos_visiveis(user):
    """Retorna queryset de conteúdos corporativos visíveis para o usuário.
    - Owner: tudo.
    - Franqueado: os seus + templates (is_template=True).
    """
    if user.is_owner():
        return ConteudoCorporativo.objects.all()
    # Franqueado
    return ConteudoCorporativo.objects.filter(
        Q(franqueado=user) | Q(is_template=True)
    )


@login_required
def conteudo_corporativo_list_view(request):
    """Lista de conteúdos corporativos"""
    user = request.user
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Você não tem permissão para acessar esta página.')
        return redirect('dashboard')

    conteudos = _conteudos_visiveis(user)
    search = request.GET.get('search', '')
    tipo_filter = request.GET.get('tipo', '')
    if search:
        conteudos = conteudos.filter(titulo__icontains=search)
    if tipo_filter:
        conteudos = conteudos.filter(tipo=tipo_filter)

    context = {
        'conteudos': conteudos,
        'tipos': ConteudoCorporativo.TIPO_CHOICES,
    }
    return render(request, 'corporativo/conteudo_list.html', context)


@login_required
def conteudo_corporativo_create_view(request):
    """Criar novo conteúdo corporativo"""
    user = request.user
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Você não tem permissão para criar conteúdo corporativo.')
        return redirect('dashboard')

    if request.method == 'POST':
        form = ConteudoCorporativoForm(request.POST)
        if form.is_valid():
            conteudo = form.save(commit=False)
            if user.is_franchisee():
                conteudo.franqueado = user
                conteudo.is_template = False  # franqueado não pode criar templates
            conteudo.save()
            messages.success(request, 'Conteúdo corporativo criado com sucesso!')
            return redirect('conteudo_corporativo_list')
    else:
        form = ConteudoCorporativoForm()

    return render(request, 'corporativo/conteudo_form.html', {'form': form})


@login_required
def conteudo_corporativo_update_view(request, pk):
    """Editar conteúdo corporativo"""
    user = request.user
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Sem permissão.')
        return redirect('dashboard')

    conteudo = get_object_or_404(ConteudoCorporativo, pk=pk)
    # Franqueado só pode editar o que é dele (não pode editar templates do owner)
    if user.is_franchisee() and conteudo.franqueado != user:
        messages.error(request, 'Você não tem permissão para editar este conteúdo.')
        return redirect('conteudo_corporativo_list')

    if request.method == 'POST':
        form = ConteudoCorporativoForm(request.POST, instance=conteudo)
        if form.is_valid():
            updated = form.save(commit=False)
            if user.is_franchisee():
                updated.franqueado = user  # garante que não pode mudar o dono
                updated.is_template = False
            updated.save()
            messages.success(request, 'Conteúdo corporativo atualizado!')
            return redirect('conteudo_corporativo_list')
    else:
        form = ConteudoCorporativoForm(instance=conteudo)

    return render(request, 'corporativo/conteudo_form.html', {'form': form, 'conteudo': conteudo})


@login_required
def conteudo_corporativo_delete_view(request, pk):
    """Deletar conteúdo corporativo"""
    user = request.user
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Sem permissão.')
        return redirect('dashboard')

    conteudo = get_object_or_404(ConteudoCorporativo, pk=pk)
    if user.is_franchisee() and conteudo.franqueado != user:
        messages.error(request, 'Você não tem permissão para deletar este conteúdo.')
        return redirect('conteudo_corporativo_list')

    if request.method == 'POST':
        conteudo.delete()
        messages.success(request, 'Conteúdo corporativo deletado!')
        return redirect('conteudo_corporativo_list')

    return render(request, 'corporativo/conteudo_confirm_delete.html', {'conteudo': conteudo})


@login_required
def conteudo_corporativo_preview_view(request, pk):
    """Preview de dados ao vivo de um conteúdo corporativo.
    Para previsão do tempo, o município é o da playlist vinculada — igual ao que
    o app Android recebe. O usuário pode escolher qual playlist simular.
    """
    user = request.user
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Sem permissão.')
        return redirect('dashboard')

    conteudo = get_object_or_404(ConteudoCorporativo, pk=pk)
    # Franqueado só pode visualizar o que é dele ou templates
    if user.is_franchisee() and conteudo.franqueado != user and not conteudo.is_template:
        messages.error(request, 'Você não tem permissão para visualizar este conteúdo.')
        return redirect('conteudo_corporativo_list')

    # Se for um DESIGN, redireciona para a visualização TV
    if conteudo.tipo == 'DESIGN':
        return redirect('design_render_tv', pk=pk)
    
    from .services import buscar_dados_corporativos

    # Playlists que contêm este conteúdo
    playlists_qs = Playlist.objects.filter(
        items__conteudo_corporativo=conteudo
    ).select_related('municipio').distinct()
    if user.is_franchisee():
        playlists_qs = playlists_qs.filter(franqueado=user)

    municipio = None
    playlist_selecionada = None
    if conteudo.tipo == 'PREVISAO_TEMPO':
        playlist_id = request.GET.get('playlist_id')
        if playlist_id:
            playlist_selecionada = playlists_qs.filter(pk=playlist_id).first()
        if not playlist_selecionada:
            # Preferir playlist com município com coordenadas
            playlist_selecionada = (
                playlists_qs.exclude(municipio__latitude=None).first()
                or playlists_qs.first()
            )
        if playlist_selecionada:
            municipio = playlist_selecionada.municipio

    dados = buscar_dados_corporativos(conteudo.tipo, municipio=municipio, conteudo=conteudo)
    config = ConfiguracaoAPI.get_config()

    context = {
        'conteudo': conteudo,
        'dados': dados,
        'config': config,
        'municipio': municipio,
        'playlists': playlists_qs if conteudo.tipo == 'PREVISAO_TEMPO' else [],
        'playlist_selecionada': playlist_selecionada,
        'conteudo_tipo': conteudo.tipo,
        'orientacao': conteudo.orientacao,
    }
    return render(request, 'corporativo/conteudo_preview.html', context)


@login_required
@xframe_options_exempt
def conteudo_corporativo_render_view(request, pk):
    """
    Renderiza o conteudo_tv.html real (1920×1080 ou 1080×1920) para ser embutido
    como iframe no preview. Idêntico ao que o app Android recebe.
    """
    from django.shortcuts import render as django_render
    from .services import buscar_dados_corporativos

    conteudo = get_object_or_404(ConteudoCorporativo, pk=pk)

    municipio = None
    playlist_id = request.GET.get('playlist_id')
    if playlist_id:
        try:
            playlist = Playlist.objects.select_related('municipio').get(pk=playlist_id)
            municipio = playlist.municipio
        except Playlist.DoesNotExist:
            pass
    if not municipio and conteudo.tipo == 'PREVISAO_TEMPO':
        playlist = (
            Playlist.objects.filter(items__conteudo_corporativo=conteudo)
            .select_related('municipio')
            .exclude(municipio__latitude=None)
            .first()
        )
        if playlist:
            municipio = playlist.municipio

    dados = buscar_dados_corporativos(conteudo.tipo, municipio=municipio, conteudo=conteudo)

    context = {
        'conteudo_tipo': conteudo.tipo,
        'dados': dados,
        'orientacao': conteudo.orientacao,
    }
    return django_render(request, 'corporativo/conteudo_tv.html', context)


@login_required
def configuracao_api_view(request):
    """Configuração das APIs externas (apenas OWNER)"""
    user = request.user
    if not user.is_owner():
        messages.error(request, 'Apenas o administrador pode configurar as APIs.')
        return redirect('dashboard')

    config = ConfiguracaoAPI.get_config()
    config.resetar_contadores_se_necessario()

    if request.method == 'POST':
        form = ConfiguracaoAPIForm(request.POST, instance=config)
        if form.is_valid():
            form.save()
            messages.success(request, 'Configurações de API atualizadas!')
            return redirect('configuracao_api')
    else:
        form = ConfiguracaoAPIForm(instance=config)

    context = {
        'form': form,
        'config': config,
    }
    return render(request, 'corporativo/configuracao_api.html', context)


# ═══════════════════════════════════════════════════════════
#  DESIGN EDITOR — Editor Canva/PPT dentro do Corporativo
# ═══════════════════════════════════════════════════════════

@login_required
def design_editor_view(request, pk=None):
    """
    Editor visual (Fabric.js) para criar ou editar um design corporativo.
    Se pk=None → novo design. Se pk → editar existente.
    """
    import json as json_mod
    user = request.user
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Sem permissão.')
        return redirect('dashboard')

    conteudo = None
    design_json_str = 'null'
    if pk:
        conteudo = get_object_or_404(ConteudoCorporativo, pk=pk, tipo='DESIGN')
        # Franqueado só pode editar o que é dele
        if user.is_franchisee() and conteudo.franqueado != user:
            messages.error(request, 'Você não tem permissão para editar este design.')
            return redirect('design_list')
        if conteudo.design_json:
            design_json_str = json_mod.dumps(conteudo.design_json)

    # Categorias de template para seletor
    categorias = (
        ConteudoCorporativo.objects
        .filter(tipo='DESIGN', is_template=True, template_categoria__gt='')
        .values_list('template_categoria', flat=True)
        .distinct()
        .order_by('template_categoria')
    )

    context = {
        'conteudo': conteudo,
        'design_json_str': design_json_str,
        'categorias': list(categorias),
    }
    return render(request, 'corporativo/design_editor.html', context)


@login_required
def design_save_api(request, pk=None):
    """
    AJAX POST → salva ou atualiza design (JSON canvas + thumbnail base64).
    Retorna JSON {success, id, message}.
    """
    import json, base64, uuid
    from django.core.files.base import ContentFile
    from django.http import JsonResponse

    if request.method != 'POST':
        return JsonResponse({'success': False, 'message': 'Método não permitido'}, status=405)

    user = request.user
    if not user.is_owner() and not user.is_franchisee():
        return JsonResponse({'success': False, 'message': 'Sem permissão'}, status=403)

    try:
        data = json.loads(request.body)
    except json.JSONDecodeError:
        return JsonResponse({'success': False, 'message': 'JSON inválido'}, status=400)

    titulo = data.get('titulo', '').strip()
    design_json = data.get('design_json')
    thumbnail_b64 = data.get('thumbnail')  # data:image/png;base64,...
    largura = int(data.get('largura', 1920))
    altura = int(data.get('altura', 1080))
    duracao = int(data.get('duracao_segundos', 15))
    is_template = bool(data.get('is_template', False))
    template_cat = data.get('template_categoria', '')

    if not titulo:
        return JsonResponse({'success': False, 'message': 'Título é obrigatório'}, status=400)
    if not design_json:
        return JsonResponse({'success': False, 'message': 'Design vazio'}, status=400)

    if pk:
        conteudo = get_object_or_404(ConteudoCorporativo, pk=pk, tipo='DESIGN')
        # Franqueado só pode editar o que é dele
        if user.is_franchisee() and conteudo.franqueado != user:
            return JsonResponse({'success': False, 'message': 'Sem permissão para editar este design'}, status=403)
    else:
        conteudo = ConteudoCorporativo(tipo='DESIGN')
        if user.is_franchisee():
            conteudo.franqueado = user

    conteudo.titulo = titulo
    conteudo.design_json = design_json
    conteudo.design_largura = largura
    conteudo.design_altura = altura
    conteudo.duracao_segundos = duracao
    # Apenas owner pode marcar como template
    conteudo.is_template = is_template if user.is_owner() else False
    conteudo.template_categoria = template_cat if user.is_owner() else ''

    # Salvar thumbnail PNG
    if thumbnail_b64 and ',' in thumbnail_b64:
        fmt, imgstr = thumbnail_b64.split(',', 1)
        ext = 'png'
        filename = f'design_{uuid.uuid4().hex[:8]}.{ext}'
        conteudo.design_thumbnail.save(filename, ContentFile(base64.b64decode(imgstr)), save=False)

    conteudo.save()

    return JsonResponse({
        'success': True,
        'id': conteudo.pk,
        'message': 'Design salvo com sucesso!',
    })


@login_required
def design_list_view(request):
    """Lista de designs criados (não-templates) com grid de thumbnails."""
    user = request.user
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Sem permissão.')
        return redirect('dashboard')

    if user.is_owner():
        designs = ConteudoCorporativo.objects.filter(tipo='DESIGN', is_template=False)
    else:
        # Franqueado vê apenas seus próprios designs (não-templates)
        designs = ConteudoCorporativo.objects.filter(tipo='DESIGN', is_template=False, franqueado=user)

    designs = designs.order_by('-updated_at')
    search = request.GET.get('search', '')
    if search:
        designs = designs.filter(titulo__icontains=search)

    context = {
        'designs': designs,
        'search': search,
    }
    return render(request, 'corporativo/design_list.html', context)


@login_required
def design_template_gallery_view(request):
    """Galeria de modelos (templates reutilizáveis)."""
    user = request.user
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Sem permissão.')
        return redirect('dashboard')

    templates = ConteudoCorporativo.objects.filter(tipo='DESIGN', is_template=True, ativo=True).order_by('template_categoria', 'titulo')

    categoria = request.GET.get('categoria', '')
    search = request.GET.get('search', '')
    if categoria:
        templates = templates.filter(template_categoria=categoria)
    if search:
        templates = templates.filter(titulo__icontains=search)

    categorias = (
        ConteudoCorporativo.objects
        .filter(tipo='DESIGN', is_template=True, template_categoria__gt='')
        .values_list('template_categoria', flat=True)
        .distinct()
        .order_by('template_categoria')
    )

    context = {
        'templates': templates,
        'categorias': list(categorias),
        'categoria_selecionada': categoria,
        'search': search,
    }
    return render(request, 'corporativo/design_template_gallery.html', context)


@login_required
def design_duplicate_view(request, pk):
    """Duplica um design existente (ou template) como novo design pessoal."""
    user = request.user
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Sem permissão.')
        return redirect('dashboard')

    original = get_object_or_404(ConteudoCorporativo, pk=pk, tipo='DESIGN')
    # Franqueado só pode duplicar seus próprios designs ou templates
    if user.is_franchisee() and original.franqueado != user and not original.is_template:
        messages.error(request, 'Você não tem permissão para duplicar este design.')
        return redirect('design_list')

    novo = ConteudoCorporativo(
        titulo=f'{original.titulo} (Cópia)',
        tipo='DESIGN',
        duracao_segundos=original.duracao_segundos,
        design_json=original.design_json,
        design_largura=original.design_largura,
        design_altura=original.design_altura,
        is_template=False,
        template_categoria='',
        ativo=True,
        franqueado=user if user.is_franchisee() else None,
    )
    # Copiar thumbnail (pode não existir no storage)
    if original.design_thumbnail:
        from django.core.files.base import ContentFile
        import uuid
        try:
            novo.design_thumbnail.save(
                f'design_{uuid.uuid4().hex[:8]}.png',
                ContentFile(original.design_thumbnail.read()),
                save=False,
            )
        except (FileNotFoundError, Exception):
            pass  # thumbnail ausente no storage — continua sem ela
    novo.save()
    messages.success(request, f'Design duplicado! Editando "{novo.titulo}".')
    return redirect('design_editor_edit', pk=novo.pk)


@login_required
def design_delete_view(request, pk):
    """Deletar design."""
    user = request.user
    if not user.is_owner() and not user.is_franchisee():
        messages.error(request, 'Sem permissão.')
        return redirect('dashboard')

    conteudo = get_object_or_404(ConteudoCorporativo, pk=pk, tipo='DESIGN')
    if user.is_franchisee() and conteudo.franqueado != user:
        messages.error(request, 'Você não tem permissão para deletar este design.')
        return redirect('design_list')

    if request.method == 'POST':
        titulo = conteudo.titulo
        conteudo.delete()
        messages.success(request, f'Design "{titulo}" removido.')
        return redirect('design_list')

    return render(request, 'corporativo/design_confirm_delete.html', {'conteudo': conteudo})


@xframe_options_exempt
def design_render_tv_view(request, pk):
    """
    Renderiza o design como HTML/SVG estático para o app de TV (WebView).
    Usa Fabric.js para gerar canvas a partir do JSON salvo.
    Suporta multi-page com transições e áudio.
    """
    conteudo = get_object_or_404(ConteudoCorporativo, pk=pk, tipo='DESIGN')
    import json as json_mod
    design_data = conteudo.design_json or {}
    dispositivo_id = request.GET.get('dispositivo_id', '')
    context = {
        'conteudo': conteudo,
        'design_json': json_mod.dumps(design_data),
        'dispositivo_id': dispositivo_id,
        'conteudo_id': conteudo.id,
        'duracao_segundos': conteudo.duracao_segundos or 30,
    }
    return render(request, 'corporativo/design_tv_render.html', context)


@login_required
def design_import_pptx_view(request):
    """
    AJAX POST: Recebe upload de arquivo .pptx, converte slides para
    estrutura multi-page JSON compatível com o editor de design.
    Retorna JSON com as pages prontas para inserção.
    """
    import json as json_mod
    import base64
    import io
    from django.http import JsonResponse

    if request.method != 'POST':
        return JsonResponse({'success': False, 'message': 'Método não permitido'}, status=405)

    user = request.user
    if not user.is_owner() and not user.is_franchisee():
        return JsonResponse({'success': False, 'message': 'Sem permissão'}, status=403)

    pptx_file = request.FILES.get('pptx') or request.FILES.get('pptx_file')
    if not pptx_file:
        return JsonResponse({'success': False, 'message': 'Nenhum arquivo enviado'}, status=400)

    if not pptx_file.name.lower().endswith('.pptx'):
        return JsonResponse({'success': False, 'message': 'Arquivo deve ser .pptx'}, status=400)

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation(pptx_file)

        # Slide dimensions (EMUs → px at 96 DPI)
        slide_w_px = int(prs.slide_width / 914400 * 96)
        slide_h_px = int(prs.slide_height / 914400 * 96)

        pages = []
        for slide_idx, slide in enumerate(prs.slides):
            fabric_objects = []

            # Extract background color
            bg_color = '#ffffff'
            try:
                bg_fill = slide.background.fill
                if bg_fill and bg_fill.type is not None:
                    if bg_fill.fore_color and bg_fill.fore_color.rgb:
                        bg_color = '#' + str(bg_fill.fore_color.rgb)
            except Exception:
                pass

            def extract_shape(shape, fabric_objects):
                """Recursively extract objects from a shape (handles groups)."""
                from pptx.enum.shapes import MSO_SHAPE_TYPE

                left_px = int((shape.left or 0) / 914400 * 96)
                top_px = int((shape.top or 0) / 914400 * 96)
                width_px = int((shape.width or 0) / 914400 * 96)
                height_px = int((shape.height or 0) / 914400 * 96)

                # Handle GROUP shapes — recurse into children
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    try:
                        for child in shape.shapes:
                            extract_shape(child, fabric_objects)
                    except Exception:
                        pass
                    return

                # Shape fill color
                fill_color = '#cccccc'
                try:
                    if hasattr(shape, 'fill') and shape.fill and shape.fill.type is not None:
                        if shape.fill.fore_color and shape.fill.fore_color.rgb:
                            fill_color = '#' + str(shape.fill.fore_color.rgb)
                except Exception:
                    pass

                # --- IMAGE: detect via hasattr(shape, 'image') — most reliable ---
                is_picture = False
                try:
                    is_picture = (shape.shape_type == MSO_SHAPE_TYPE.PICTURE or
                                  shape.shape_type == MSO_SHAPE_TYPE.LINKED_PICTURE)
                except Exception:
                    pass
                if not is_picture:
                    # fallback: try accessing shape.image directly
                    try:
                        _ = shape.image
                        is_picture = True
                    except Exception:
                        pass

                if is_picture:
                    try:
                        image = shape.image
                        img_bytes = image.blob
                        content_type = (image.content_type or 'image/png').split(';')[0].strip()
                        if content_type not in ('image/png', 'image/jpeg', 'image/gif', 'image/webp'):
                            content_type = 'image/png'

                        b64 = base64.b64encode(img_bytes).decode('utf-8')
                        data_url = f'data:{content_type};base64,{b64}'

                        # Get natural pixel dimensions via PIL
                        try:
                            from PIL import Image as PILImage
                            pil_img = PILImage.open(io.BytesIO(img_bytes))
                            natural_w, natural_h = pil_img.size
                            pil_img.close()
                        except Exception:
                            natural_w = width_px or 100
                            natural_h = height_px or 100

                        natural_w = max(natural_w, 1)
                        natural_h = max(natural_h, 1)
                        target_w = width_px if width_px > 0 else natural_w
                        target_h = height_px if height_px > 0 else natural_h

                        fabric_objects.append({
                            'type': 'image',
                            'version': '5.3.1',
                            'originX': 'left',
                            'originY': 'top',
                            'left': left_px,
                            'top': top_px,
                            'width': natural_w,
                            'height': natural_h,
                            'scaleX': round(target_w / natural_w, 6),
                            'scaleY': round(target_h / natural_h, 6),
                            'angle': 0,
                            'opacity': 1,
                            'flipX': False,
                            'flipY': False,
                            'src': data_url,
                            'crossOrigin': 'anonymous',
                            'filters': [],
                        })
                        return  # done with this shape
                    except Exception:
                        pass  # fall through to text/shape handling

                if shape.has_text_frame:
                    # TEXT SHAPE
                    for para in shape.text_frame.paragraphs:
                        full_text = para.text.strip()
                        if not full_text:
                            continue

                        # Gather text properties from first run
                        font_size = 24
                        font_family = 'Inter'
                        font_color = '#333333'
                        font_weight = 'normal'
                        font_style = 'normal'
                        text_align = 'left'
                        underline = False

                        if para.runs:
                            run = para.runs[0]
                            if run.font.size:
                                font_size = int(run.font.size / 12700)  # EMU → pt
                            if run.font.name:
                                font_family = run.font.name
                            if run.font.bold:
                                font_weight = 'bold'
                            if run.font.italic:
                                font_style = 'italic'
                            if run.font.underline:
                                underline = True
                            try:
                                if run.font.color and run.font.color.rgb:
                                    font_color = '#' + str(run.font.color.rgb)
                            except Exception:
                                pass

                        if para.alignment:
                            align_map = {
                                PP_ALIGN.LEFT: 'left',
                                PP_ALIGN.CENTER: 'center',
                                PP_ALIGN.RIGHT: 'right',
                            }
                            text_align = align_map.get(para.alignment, 'left')

                        fabric_objects.append({
                            'type': 'i-text',
                            'text': full_text,
                            'left': left_px,
                            'top': top_px,
                            'width': width_px,
                            'fontFamily': font_family,
                            'fontSize': font_size,
                            'fill': font_color,
                            'fontWeight': font_weight,
                            'fontStyle': font_style,
                            'underline': underline,
                            'textAlign': text_align,
                            'editable': True,
                        })
                        # Offset subsequent paragraphs vertically
                        top_px += int(font_size * 1.4)

                else:
                    # GEOMETRIC SHAPE → rect (fallback for non-text, non-image shapes)
                    stroke_color = 'transparent'
                    stroke_width = 0
                    try:
                        if shape.line and shape.line.fill and shape.line.fill.fore_color:
                            stroke_color = '#' + str(shape.line.fill.fore_color.rgb)
                            stroke_width = int((shape.line.width or 0) / 12700)
                    except Exception:
                        pass

                    rx = 0
                    fabric_objects.append({
                        'type': 'rect',
                        'left': left_px,
                        'top': top_px,
                        'width': width_px,
                        'height': height_px,
                        'fill': fill_color,
                        'stroke': stroke_color,
                        'strokeWidth': stroke_width,
                        'rx': rx,
                        'ry': rx,
                    })

            for shape in slide.shapes:
                extract_shape(shape, fabric_objects)

            # Build Fabric-compatible canvas JSON for this slide
            fabric_json = {
                'version': '5.3.1',
                'objects': fabric_objects,
                'background': bg_color,
            }

            pages.append({
                'id': f'page_{slide_idx + 1}',
                'name': f'Slide {slide_idx + 1}',
                'duration': 5,
                'transition': 'fade',
                'transitionDuration': 0.5,
                'audioUrl': '',
                'fabricJson': fabric_json,
                'animations': [],
            })

        return JsonResponse({
            'success': True,
            'canvasWidth': slide_w_px,
            'canvasHeight': slide_h_px,
            'pages': pages,
            'totalSlides': len(pages),
            'message': f'{len(pages)} slides importados com sucesso!',
        })

    except Exception as e:
        import traceback
        traceback.print_exc()
        return JsonResponse({
            'success': False,
            'message': f'Erro ao processar PPTX: {str(e)}'
        }, status=400)


@login_required
def design_audio_upload_view(request):
    """
    AJAX POST: Recebe upload de arquivo de áudio (mp3, wav, ogg).
    Salva no media/designs/audio/ e retorna a URL.
    """
    import uuid
    from django.http import JsonResponse

    if request.method != 'POST':
        return JsonResponse({'success': False, 'message': 'Método não permitido'}, status=405)

    user = request.user
    if not user.is_owner() and not user.is_franchisee():
        return JsonResponse({'success': False, 'message': 'Sem permissão'}, status=403)

    audio_file = request.FILES.get('audio') or request.FILES.get('audio_file')
    if not audio_file:
        return JsonResponse({'success': False, 'message': 'Nenhum arquivo enviado'}, status=400)

    allowed_ext = ['.mp3', '.wav', '.ogg', '.m4a', '.aac']
    ext = os.path.splitext(audio_file.name)[1].lower()
    if ext not in allowed_ext:
        return JsonResponse({
            'success': False,
            'message': f'Formato não suportado. Use: {", ".join(allowed_ext)}'
        }, status=400)

    # Save file
    filename = f'audio_{uuid.uuid4().hex[:8]}{ext}'
    save_dir = os.path.join('media', 'designs', 'audio')
    os.makedirs(save_dir, exist_ok=True)
    save_path = os.path.join(save_dir, filename)

    with open(save_path, 'wb') as f:
        for chunk in audio_file.chunks():
            f.write(chunk)

    audio_url = f'/media/designs/audio/{filename}'
    return JsonResponse({
        'success': True,
        'audioUrl': audio_url,
        'filename': audio_file.name,
        'message': 'Áudio enviado com sucesso!',
    })


def design_video_upload_view(request):
    """
    AJAX POST: Recebe upload de arquivo de vídeo (mp4, webm, mov).
    Salva no media/designs/videos/ e retorna a URL.
    """
    import uuid
    from django.http import JsonResponse

    if request.method != 'POST':
        return JsonResponse({'success': False, 'message': 'Método não permitido'}, status=405)

    user = request.user
    if not user.is_owner() and not user.is_franchisee():
        return JsonResponse({'success': False, 'message': 'Sem permissão'}, status=403)

    video_file = request.FILES.get('video') or request.FILES.get('video_file')
    if not video_file:
        return JsonResponse({'success': False, 'message': 'Nenhum arquivo enviado'}, status=400)

    allowed_ext = ['.mp4', '.webm', '.mov', '.avi', '.m4v']
    ext = os.path.splitext(video_file.name)[1].lower()
    if ext not in allowed_ext:
        return JsonResponse({
            'success': False,
            'message': f'Formato não suportado. Use: {", ".join(allowed_ext)}'
        }, status=400)

    # Limite de 200 MB
    max_size = 200 * 1024 * 1024
    if video_file.size > max_size:
        return JsonResponse({'success': False, 'message': 'Arquivo muito grande (máx 200 MB)'}, status=400)

    filename = f'video_{uuid.uuid4().hex[:8]}{ext}'
    from django.core.files.storage import default_storage
    from django.core.files.base import ContentFile

    storage_path = f'designs/videos/{filename}'
    default_storage.save(storage_path, ContentFile(video_file.read()))

    raw_url = default_storage.url(storage_path)
    # Se URL relativa (storage local), tornar absoluta
    if raw_url.startswith('/'):
        video_url = request.build_absolute_uri(raw_url)
        if 'railway.app' in video_url:
            video_url = video_url.replace('http://', 'https://')
    else:
        video_url = raw_url  # R2 já retorna URL absoluta
    return JsonResponse({
        'success': True,
        'videoUrl': video_url,
        'filename': video_file.name,
        'message': 'Vídeo enviado com sucesso!',
    })


@login_required
def design_video_library_view(request):
    """
    AJAX GET: Retorna lista de vídeos aprovados disponíveis para o usuário
    usar como vídeo de fundo no design editor.
    """
    from django.conf import settings

    user = request.user
    if not user.is_owner() and not user.is_franchisee():
        return JsonResponse({'success': False, 'message': 'Sem permissão'}, status=403)

    videos_qs = Video.objects.select_related('cliente').filter(
        arquivo__isnull=False,
        ativo=True,
        status__in=['APPROVED', 'PENDING'],
    ).exclude(arquivo='').order_by('-created_at')

    if user.is_franchisee():
        clientes_ids = Cliente.objects.filter(franqueado=user).values_list('id', flat=True)
        videos_qs = videos_qs.filter(cliente_id__in=clientes_ids)

    result = []
    for v in videos_qs[:100]:  # limite de 100
        if not v.arquivo_existe():
            continue
        relative = v.arquivo.url  # já inclui MEDIA_URL
        url = request.build_absolute_uri(relative)
        if 'railway.app' in url:
            url = url.replace('http://', 'https://')
        result.append({
            'id': v.pk,
            'titulo': v.titulo,
            'cliente': v.cliente.empresa,
            'url': url,
            'extensao': v.extensao,
            'tamanho': v.file_size_bytes,
        })

    return JsonResponse({'success': True, 'videos': result})


# ═══════════════════════════════════════════════════════════════
#  FREE IMAGE BANK — Proxy to Pixabay API
# ═══════════════════════════════════════════════════════════════

@login_required
def design_search_images_view(request):
    """
    AJAX GET: Busca imagens no Pixabay (gratuito).
    Params: q (query), page, per_page, image_type (photo, illustration, vector), category
    Returns JSON with results array of {id, thumbnail, preview, fullUrl, tags, width, height, user}.
    """
    from django.conf import settings
    import urllib.parse
    import urllib.request
    import urllib.error
    import json as json_mod

    if request.method != 'GET':
        return JsonResponse({'success': False, 'message': 'Método não permitido'}, status=405)

    try:
        # Tenta pegar a chave do settings (que lê do .env via python-decouple)
        PIXABAY_API_KEY = getattr(settings, 'PIXABAY_API_KEY', '')
        if not PIXABAY_API_KEY:
            # Fallback: tenta pegar diretamente da variável de ambiente
            PIXABAY_API_KEY = os.environ.get('PIXABAY_API_KEY', '')
        
        print(f"[DEBUG] PIXABAY_API_KEY configured: {bool(PIXABAY_API_KEY)}")
        if PIXABAY_API_KEY:
            print(f"[DEBUG] API Key: {PIXABAY_API_KEY[:10]}...{PIXABAY_API_KEY[-5:]}")

        query = request.GET.get('q', '').strip()
        page = int(request.GET.get('page', 1))
        per_page = min(int(request.GET.get('per_page', 40)), 200)
        image_type = request.GET.get('image_type', 'all')  # photo, illustration, vector, all
        category = request.GET.get('category', '')
        colors = request.GET.get('colors', '')  # for filtering by color
        editors_choice = request.GET.get('editors_choice', 'false')

        if not PIXABAY_API_KEY:
            print("[DEBUG] No Pixabay API key found, using fallback (Lorem Picsum)")
            # Return curated fallback using Lorem Picsum (no API key needed)
            return _fallback_image_search(query, page, per_page)
    except Exception as e:
        print(f"[ERROR] Error in initial setup: {type(e).__name__}: {str(e)}")
        import traceback
        traceback.print_exc()
        return JsonResponse({
            'success': False,
            'message': f'Erro na configuração: {str(e)}',
            'results': [],
        }, status=500)

    # Build Pixabay API URL
    params = {
        'key': PIXABAY_API_KEY,
        'q': query or 'nature',
        'page': page,
        'per_page': per_page,
        'image_type': image_type if image_type in ('photo', 'illustration', 'vector') else 'all',
        'lang': 'pt',
        'safesearch': 'true',
        'editors_choice': editors_choice,
    }
    if category:
        params['category'] = category
    if colors:
        params['colors'] = colors

    url = 'https://pixabay.com/api/?' + urllib.parse.urlencode(params)
    
    print(f"[DEBUG] Calling Pixabay API: {url.replace(PIXABAY_API_KEY, 'KEY_HIDDEN')}")

    try:
        req = urllib.request.Request(url, headers={'User-Agent': 'MediaExpand/1.0'})
        with urllib.request.urlopen(req, timeout=10) as resp:
            data = json_mod.loads(resp.read().decode())

        print(f"[DEBUG] Pixabay response: totalHits={data.get('totalHits', 0)}, hits={len(data.get('hits', []))}")

        results = []
        for hit in data.get('hits', []):
            results.append({
                'id': hit.get('id'),
                'thumbnail': hit.get('previewURL', ''),
                'preview': hit.get('webformatURL', ''),
                'fullUrl': hit.get('largeImageURL', hit.get('webformatURL', '')),
                'tags': hit.get('tags', ''),
                'width': hit.get('imageWidth', 0),
                'height': hit.get('imageHeight', 0),
                'user': hit.get('user', ''),
                'source': 'pixabay',
            })

        return JsonResponse({
            'success': True,
            'results': results,
            'total': data.get('totalHits', 0),
            'page': page,
            'per_page': per_page,
        })
    except urllib.error.HTTPError as e:
        error_body = e.read().decode() if hasattr(e, 'read') else ''
        print(f"[ERROR] Pixabay HTTP {e.code}: {e.reason}")
        print(f"[ERROR] Response body: {error_body}")
        
        if e.code == 400:
            message = 'API Key inválida ou parâmetros incorretos'
        elif e.code == 429:
            message = 'Limite de requisições excedido. Usando fallback...'
            return _fallback_image_search(query, page, per_page)
        else:
            message = f'Erro da API Pixabay ({e.code}): {e.reason}'
        
        return JsonResponse({
            'success': False,
            'message': message,
            'results': [],
            'fallback_available': True,
        }, status=200)  # Retorna 200 para o cliente processar o fallback
    except urllib.error.URLError as e:
        print(f"[ERROR] Pixabay connection error: {str(e)}")
        return JsonResponse({
            'success': False,
            'message': f'Erro de conexão com Pixabay. Usando fallback...',
            'results': [],
        }, status=200)
    except Exception as e:
        print(f"[ERROR] Pixabay unexpected error: {type(e).__name__}: {str(e)}")
        import traceback
        traceback.print_exc()
        
        return JsonResponse({
            'success': False,
            'message': f'Erro ao buscar imagens: {str(e)}',
            'results': [],
        }, status=500)


def _fallback_image_search(query, page, per_page):
    """
    Fallback: Return curated images from Lorem Picsum (no API key needed).
    Not searchable by keyword, but provides free high-quality photos.
    """
    import urllib.request
    import urllib.error
    import json as json_mod

    print(f"[DEBUG] Using fallback (Lorem Picsum) for query: {query}")

    try:
        url = f'https://picsum.photos/v2/list?page={page}&limit={per_page}'
        req = urllib.request.Request(url, headers={'User-Agent': 'MediaExpand/1.0'})
        with urllib.request.urlopen(req, timeout=10) as resp:
            items = json_mod.loads(resp.read().decode())

        print(f"[DEBUG] Fallback returned {len(items)} images")

        results = []
        for item in items:
            pid = item.get('id', '')
            results.append({
                'id': pid,
                'thumbnail': f'https://picsum.photos/id/{pid}/300/200',
                'preview': f'https://picsum.photos/id/{pid}/600/400',
                'fullUrl': f'https://picsum.photos/id/{pid}/1920/1080',
                'tags': item.get('author', ''),
                'width': item.get('width', 1920),
                'height': item.get('height', 1080),
                'user': item.get('author', ''),
                'source': 'picsum',
            })

        return JsonResponse({
            'success': True,
            'results': results,
            'total': 1000,
            'page': page,
            'per_page': per_page,
            'message': 'Usando imagens gratuitas (Lorem Picsum)',
        })
    except Exception as e:
        print(f"[ERROR] Fallback error: {type(e).__name__}: {str(e)}")
        import traceback
        traceback.print_exc()
        
        return JsonResponse({
            'success': False,
            'message': f'Erro no fallback: {str(e)}',
            'results': [],
        }, status=500)


# ═══════════════════════════════════════════════════════════════
#  FREE ICON LIBRARY — Proxy to Iconify API (100k+ icons, no key)
# ═══════════════════════════════════════════════════════════════

@login_required
def design_search_icons_view(request):
    """
    AJAX GET: Busca ícones no Iconify (gratuito, sem API key).
    Params: q (query), prefix (icon set), limit
    Returns JSON with icons array of {name, prefix, svg, tags}.
    """
    import urllib.parse
    import urllib.request
    import json as json_mod

    if request.method != 'GET':
        return JsonResponse({'success': False, 'message': 'Método não permitido'}, status=405)

    query = request.GET.get('q', '').strip()
    prefix = request.GET.get('prefix', '')  # e.g. 'mdi', 'fa', 'lucide'
    limit = min(int(request.GET.get('limit', 60)), 999)

    if not query:
        # Return popular icon sets info
        return JsonResponse({
            'success': True,
            'icons': [],
            'total': 0,
            'collections': [
                {'prefix': 'mdi', 'name': 'Material Design Icons', 'total': 7400},
                {'prefix': 'fa-solid', 'name': 'Font Awesome Solid', 'total': 1400},
                {'prefix': 'fa-regular', 'name': 'Font Awesome Regular', 'total': 162},
                {'prefix': 'lucide', 'name': 'Lucide Icons', 'total': 1500},
                {'prefix': 'tabler', 'name': 'Tabler Icons', 'total': 4700},
                {'prefix': 'ph', 'name': 'Phosphor Icons', 'total': 7500},
                {'prefix': 'ri', 'name': 'Remix Icons', 'total': 2800},
                {'prefix': 'bi', 'name': 'Bootstrap Icons', 'total': 2000},
                {'prefix': 'carbon', 'name': 'Carbon Icons', 'total': 2100},
                {'prefix': 'ion', 'name': 'Ionicons', 'total': 1300},
                {'prefix': 'fluent', 'name': 'Fluent UI Icons', 'total': 4200},
                {'prefix': 'heroicons', 'name': 'Heroicons', 'total': 600},
            ]
        })

    # Search icons via Iconify API
    params = {
        'query': query,
        'limit': limit,
    }
    if prefix:
        params['prefixes'] = prefix

    url = 'https://api.iconify.design/search?' + urllib.parse.urlencode(params)

    try:
        req = urllib.request.Request(url, headers={'User-Agent': 'MediaExpand/1.0'})
        with urllib.request.urlopen(req, timeout=10) as resp:
            data = json_mod.loads(resp.read().decode())

        icons_list = data.get('icons', [])
        total = data.get('total', len(icons_list))

        # For each icon, build the SVG URL (Iconify serves SVGs directly)
        results = []
        for icon_name in icons_list:
            # icon_name format: "prefix:name" e.g. "mdi:home"
            parts = icon_name.split(':', 1)
            if len(parts) == 2:
                pfx, name = parts
            else:
                pfx, name = '', icon_name

            svg_url = f'https://api.iconify.design/{pfx}/{name}.svg'
            results.append({
                'name': name,
                'prefix': pfx,
                'fullName': icon_name,
                'svgUrl': svg_url,
                'previewUrl': f'https://api.iconify.design/{pfx}/{name}.svg?height=64',
            })

        return JsonResponse({
            'success': True,
            'icons': results,
            'total': total,
        })
    except Exception as e:
        return JsonResponse({
            'success': False,
            'message': f'Erro ao buscar ícones: {str(e)}',
            'icons': [],
        }, status=500)


@login_required
def design_get_icon_svg_view(request):
    """
    AJAX GET: Proxy to fetch SVG content from Iconify.
    Params: icon (full name like "mdi:home"), color, size
    Returns the SVG string directly for embedding into Fabric.js canvas.
    """
    import urllib.parse
    import urllib.request

    icon = request.GET.get('icon', '').strip()
    color = request.GET.get('color', '#333333')
    size = int(request.GET.get('size', 256))

    if not icon or ':' not in icon:
        return JsonResponse({'success': False, 'message': 'Ícone inválido'}, status=400)

    parts = icon.split(':', 1)
    pfx, name = parts[0], parts[1]

    svg_url = f'https://api.iconify.design/{pfx}/{name}.svg?color={urllib.parse.quote(color)}&height={size}&width={size}'

    try:
        req = urllib.request.Request(svg_url, headers={'User-Agent': 'MediaExpand/1.0'})
        with urllib.request.urlopen(req, timeout=10) as resp:
            svg_content = resp.read().decode()

        return JsonResponse({
            'success': True,
            'svg': svg_content,
            'icon': icon,
        })
    except Exception as e:
        return JsonResponse({
            'success': False,
            'message': f'Erro ao carregar ícone: {str(e)}',
        }, status=500)


# ═══════════════════════════════════════════════════════════════
#  FREE PNG/STICKER SEARCH — Proxy to Sticker/PNG APIs
# ═══════════════════════════════════════════════════════════════

@login_required
def design_search_stickers_view(request):
    """
    AJAX GET: Busca stickers/PNGs transparentes.
    Uses Pixabay with transparent filter or vector type.
    Params: q (query), page, per_page
    """
    from django.conf import settings
    import urllib.parse
    import urllib.request
    import urllib.error
    import json as json_mod

    if request.method != 'GET':
        return JsonResponse({'success': False, 'message': 'Método não permitido'}, status=405)

    try:
        PIXABAY_API_KEY = getattr(settings, 'PIXABAY_API_KEY', '')
        if not PIXABAY_API_KEY:
            PIXABAY_API_KEY = os.environ.get('PIXABAY_API_KEY', '')
        
        print(f"[DEBUG] Stickers - PIXABAY_API_KEY configured: {bool(PIXABAY_API_KEY)}")

        query = request.GET.get('q', '').strip()
        page = int(request.GET.get('page', 1))
        per_page = min(int(request.GET.get('per_page', 40)), 200)

        if not PIXABAY_API_KEY:
            print("[DEBUG] Stickers - No API key, returning empty results")
            return JsonResponse({
                'success': True,
                'results': [],
                'total': 0,
                'message': 'Configure PIXABAY_API_KEY para buscar PNGs transparentes',
            })
    except Exception as e:
        print(f"[ERROR] Stickers - Setup error: {type(e).__name__}: {str(e)}")
        import traceback
        traceback.print_exc()
        return JsonResponse({
            'success': False,
            'message': f'Erro na configuração: {str(e)}',
            'results': [],
        }, status=500)

    # Search for vectors/illustrations (usually have transparency)
    params = {
        'key': PIXABAY_API_KEY,
        'q': query or 'icon',
        'page': page,
        'per_page': per_page,
        'image_type': 'vector',
        'lang': 'pt',
        'safesearch': 'true',
    }

    url = 'https://pixabay.com/api/?' + urllib.parse.urlencode(params)
    print(f"[DEBUG] Stickers - Calling Pixabay API for vectors")

    try:
        req = urllib.request.Request(url, headers={'User-Agent': 'MediaExpand/1.0'})
        with urllib.request.urlopen(req, timeout=10) as resp:
            data = json_mod.loads(resp.read().decode())
        
        print(f"[DEBUG] Stickers - Got {len(data.get('hits', []))} results")

        results = []
        for hit in data.get('hits', []):
            results.append({
                'id': hit.get('id'),
                'thumbnail': hit.get('previewURL', ''),
                'preview': hit.get('webformatURL', ''),
                'fullUrl': hit.get('largeImageURL', hit.get('webformatURL', '')),
                'tags': hit.get('tags', ''),
                'width': hit.get('imageWidth', 0),
                'height': hit.get('imageHeight', 0),
                'user': hit.get('user', ''),
                'source': 'pixabay_vector',
            })

        return JsonResponse({
            'success': True,
            'results': results,
            'total': data.get('totalHits', 0),
            'page': page,
            'per_page': per_page,
        })
    except urllib.error.HTTPError as e:
        error_body = e.read().decode() if hasattr(e, 'read') else ''
        print(f"[ERROR] Stickers HTTP {e.code}: {e.reason}")
        print(f"[ERROR] Response: {error_body}")
        return JsonResponse({
            'success': False,
            'message': f'Erro da API ({e.code}): {e.reason}',
            'results': [],
        }, status=200)
    except Exception as e:
        print(f"[ERROR] Stickers error: {type(e).__name__}: {str(e)}")
        import traceback
        traceback.print_exc()
        return JsonResponse({
            'success': False,
            'message': f'Erro ao buscar PNGs: {str(e)}',
            'results': [],
        }, status=500)


# =============================================================================
# LAB DE CODIFICAÇÃO DE VÍDEO — página provisória para testar variantes FireTV
# =============================================================================

import os
import subprocess
import threading
import uuid as _uuid

# Dicionário em memória para rastrear status dos jobs de codificação.
# Chave: job_id (str). Valor: dict com lista de variantes e seus status.
_LAB_JOBS: dict = {}
_LAB_JOBS_LOCK = threading.Lock()

# Variantes de codificação a testar no Fire TV Stick
LAB_VARIANTS = [
    {
        'key': 'v1_480p_baseline31',
        'label': 'V1 — 480p Baseline 3.1 (referência atual)',
        'res_v': (480, 854),    # largura × altura para vídeo VERTICAL
        'res_h': (854, 480),    # largura × altura para vídeo HORIZONTAL
        'profile': 'baseline',
        'level': '3.1',
        'bitrate': '2M', 'maxrate': '2M', 'bufsize': '4M',
        'gop': None,
        'color': False,   # sem flags -colorspace/-color_primaries/-color_trc
        'brand': None,
        'extra_vf': '',
    },
    {
        'key': 'v2_540p_baseline31',
        'label': 'V2 — 540p Baseline 3.1',
        'res_v': (540, 960),
        'res_h': (960, 540),
        'profile': 'baseline',
        'level': '3.1',
        'bitrate': '2M', 'maxrate': '2M', 'bufsize': '4M',
        'gop': None,
        'color': False,
        'brand': None,
        'extra_vf': '',
    },
    {
        'key': 'v3_720p_baseline31',
        'label': 'V3 — 720p Baseline 3.1',
        'res_v': (720, 1280),
        'res_h': (1280, 720),
        'profile': 'baseline',
        'level': '3.1',
        'bitrate': '3M', 'maxrate': '3M', 'bufsize': '6M',
        'gop': None,
        'color': False,
        'brand': None,
        'extra_vf': '',
    },
    {
        'key': 'v4_720p_main31_color',
        'label': 'V4 — 720p Main 3.1 + VBV + BT.709',
        'res_v': (720, 1280),
        'res_h': (1280, 720),
        'profile': 'main',
        'level': '3.1',
        'bitrate': '3M', 'maxrate': '3M', 'bufsize': '6M',
        'gop': 60,
        'color': True,
        'brand': 'mp42',
        'extra_vf': ',setsar=1',
    },
    {
        'key': 'v5_720p_main40_color',
        'label': 'V5 — 720p Main 4.0 + VBV + BT.709',
        'res_v': (720, 1280),
        'res_h': (1280, 720),
        'profile': 'main',
        'level': '4.0',
        'bitrate': '3M', 'maxrate': '3M', 'bufsize': '6M',
        'gop': 60,
        'color': True,
        'brand': 'mp42',
        'extra_vf': ',setsar=1',
    },
    {
        'key': 'v6_1080p_main40_color',
        'label': 'V6 — 1080p Main 4.0 + VBV + BT.709 (sugestão GPT)',
        'res_v': (1080, 1920),
        'res_h': (1920, 1080),
        'profile': 'main',
        'level': '4.0',
        'bitrate': '5M', 'maxrate': '5M', 'bufsize': '10M',
        'gop': 60,
        'color': True,
        'brand': 'mp42',
        'extra_vf': ',setsar=1',
    },
    {
        'key': 'v7_1080p_main40_nocolor',
        'label': 'V7 — 1080p Main 4.0 + VBV sem BT.709',
        'res_v': (1080, 1920),
        'res_h': (1920, 1080),
        'profile': 'main',
        'level': '4.0',
        'bitrate': '5M', 'maxrate': '5M', 'bufsize': '10M',
        'gop': 60,
        'color': False,
        'brand': None,
        'extra_vf': '',
    },
    {
        'key': 'v8_1080p_baseline40',
        'label': 'V8 — 1080p Baseline 4.0 sem color',
        'res_v': (1080, 1920),
        'res_h': (1920, 1080),
        'profile': 'baseline',
        'level': '4.0',
        'bitrate': '4M', 'maxrate': '4M', 'bufsize': '8M',
        'gop': None,
        'color': False,
        'brand': None,
        'extra_vf': '',
    },
]


def _lab_build_ffmpeg_cmd(variant, input_path, output_path, orient):
    """Constrói o comando ffmpeg para uma variante do lab."""
    w, h = variant['res_v'] if orient == 'VERTICAL' else variant['res_h']
    vf = f'scale={w}:{h}:flags=lanczos,format=yuv420p{variant["extra_vf"]}'

    cmd = [
        'ffmpeg', '-y',
        '-i', input_path,
        '-vf', vf,
        '-map_metadata', '-1',
        '-c:v', 'libx264',
        '-profile:v', variant['profile'],
        '-level', variant['level'],
        '-pix_fmt', 'yuv420p',
        '-r', '30',
        '-b:v', variant['bitrate'],
        '-maxrate', variant['maxrate'],
        '-bufsize', variant['bufsize'],
        '-preset', 'medium',
        '-vsync', 'cfr',
    ]

    if variant['gop']:
        cmd += ['-g', str(variant['gop']), '-keyint_min', str(variant['gop'])]

    if variant['color']:
        cmd += [
            '-color_range', 'tv',
            '-colorspace', 'bt709',
            '-color_primaries', 'bt709',
            '-color_trc', 'bt709',
        ]

    cmd += ['-c:a', 'aac', '-b:a', '160k', '-ar', '44100']

    movflags = '+faststart'
    cmd += ['-movflags', movflags]

    if variant['brand']:
        cmd += ['-brand', variant['brand'], '-tag:v', 'avc1']

    cmd.append(output_path)
    return cmd


def _lab_run_variant(job_id, variant_key, client_id, title, input_path, orient):
    """Executa ffmpeg para uma variante em background e cria o Video resultante."""
    import shutil
    import tempfile
    from django.conf import settings

    variant = next(v for v in LAB_VARIANTS if v['key'] == variant_key)

    def _set_status(s, msg='', video_id=None):
        with _LAB_JOBS_LOCK:
            for item in _LAB_JOBS[job_id]['variants']:
                if item['key'] == variant_key:
                    item['status'] = s
                    item['msg'] = msg
                    if video_id:
                        item['video_id'] = video_id
                    break

    _set_status('running', 'Codificando com ffmpeg...')

    try:
        # Sempre usar /tmp para output (funciona com local e R2)
        out_filename = f'lab_{job_id[:8]}_{variant_key}.mp4'
        out_path = os.path.join('/tmp', out_filename)

        cmd = _lab_build_ffmpeg_cmd(variant, input_path, out_path, orient)

        result = subprocess.run(cmd, capture_output=True, text=True, timeout=900)

        if result.returncode != 0:
            _set_status('error', result.stderr[-400:] if result.stderr else 'ffmpeg falhou')
            return

        if not os.path.exists(out_path) or os.path.getsize(out_path) == 0:
            _set_status('error', 'Arquivo de saída vazio ou não gerado')
            return

        # Cria um Video no banco para ser acessado nas playlists.
        # Salva SEM arquivo para não disparar _normalizar_video() no save(),
        # depois atualiza o arquivo via update() (que não chama o hook).
        from .models import Cliente
        try:
            cliente = Cliente.objects.get(pk=client_id)
        except Cliente.DoesNotExist:
            _set_status('error', f'Cliente ID {client_id} não encontrado')
            return

        # Upload para storage (R2 ou local)
        from django.core.files.storage import default_storage
        from django.core.files import File as DjangoFile
        storage_path = f'lab_outputs/{out_filename}'
        try:
            with open(out_path, 'rb') as f:
                saved_name = default_storage.save(storage_path, DjangoFile(f))
        finally:
            if os.path.exists(out_path):
                try:
                    os.remove(out_path)
                except OSError:
                    pass

        video = Video(
            cliente=cliente,
            titulo=f'[LAB {variant["key"].upper()}] {title}',
            descricao=f'Variante: {variant["label"]}\nJob: {job_id}',
            status='APPROVED',
            ativo=True,
            orientacao=orient,
        )
        video.save()

        Video.objects.filter(pk=video.pk).update(arquivo=saved_name, orientacao=orient)

        try:
            file_size_mb = round(default_storage.size(saved_name) / 1024 / 1024, 1)
        except Exception:
            file_size_mb = 0
        _set_status('done', f'{file_size_mb} MB — Video ID {video.pk}', video_id=video.pk)

    except subprocess.TimeoutExpired:
        _set_status('error', 'Timeout após 15 minutos')
    except Exception as e:
        _set_status('error', str(e)[:300])


@login_required
def lab_video_encode_view(request):
    """Página provisória para testar variantes de codificação de vídeo no Fire TV."""
    import shutil

    ffmpeg_ok = bool(shutil.which('ffmpeg'))
    clientes = Cliente.objects.order_by('empresa')

    if request.method == 'POST':
        arquivo = request.FILES.get('arquivo')
        titulo = request.POST.get('titulo', '').strip()
        cliente_id = request.POST.get('cliente_id')

        if not arquivo or not titulo or not cliente_id:
            return render(request, 'lab/video_encode.html', {
                'clientes': clientes,
                'ffmpeg_ok': ffmpeg_ok,
                'error': 'Preencha todos os campos obrigatórios.',
                'variants': LAB_VARIANTS,
            })

        # Salva o original em /tmp (s\u00f3 precisa existir durante processamento ffmpeg)

        job_id = str(_uuid.uuid4())
        ext = os.path.splitext(arquivo.name)[1].lower() or '.mp4'
        input_filename = f'lab_{job_id[:8]}_original{ext}'
        # Usa /tmp para o input original (só precisa existir durante processamento)
        input_path = os.path.join('/tmp', input_filename)

        with open(input_path, 'wb') as f:
            for chunk in arquivo.chunks():
                f.write(chunk)

        # Detecta orientação do arquivo original
        orient, orig_w, orig_h = Video._detectar_orientacao_video(input_path)

        # Registra o job em memória
        with _LAB_JOBS_LOCK:
            _LAB_JOBS[job_id] = {
                'titulo': titulo,
                'cliente_id': int(cliente_id),
                'orient': orient,
                'input_path': input_path,
                'variants': [
                    {
                        'key': v['key'],
                        'label': v['label'],
                        'status': 'queued',
                        'msg': '',
                        'video_id': None,
                    }
                    for v in LAB_VARIANTS
                ],
            }

        # Dispara uma thread por variante (paralelas)
        for v in LAB_VARIANTS:
            t = threading.Thread(
                target=_lab_run_variant,
                args=(job_id, v['key'], int(cliente_id), titulo, input_path, orient),
                daemon=True,
            )
            t.start()

        return redirect(f'/lab/video-encode/{job_id}/')

    return render(request, 'lab/video_encode.html', {
        'clientes': clientes,
        'ffmpeg_ok': ffmpeg_ok,
        'variants': LAB_VARIANTS,
    })


@login_required
def lab_video_job_view(request, job_id):
    """Página de progresso de um job de codificação. Exibe status de cada variante."""
    clientes = Cliente.objects.order_by('empresa')
    ffmpeg_ok = True

    with _LAB_JOBS_LOCK:
        job = _LAB_JOBS.get(job_id)

    if not job:
        return render(request, 'lab/video_encode.html', {
            'clientes': clientes,
            'ffmpeg_ok': ffmpeg_ok,
            'variants': LAB_VARIANTS,
            'error': f'Job {job_id} não encontrado (servidor pode ter reiniciado).',
        })

    return render(request, 'lab/video_job.html', {
        'job_id': job_id,
        'job': job,
    })


@login_required
def lab_video_status_api(request, job_id):
    """Retorna JSON com status de todas as variantes de um job."""
    with _LAB_JOBS_LOCK:
        job = _LAB_JOBS.get(job_id)

    if not job:
        return JsonResponse({'error': 'Job não encontrado'}, status=404)

    total = len(job['variants'])
    done = sum(1 for v in job['variants'] if v['status'] == 'done')
    errors = sum(1 for v in job['variants'] if v['status'] == 'error')

    return JsonResponse({
        'job_id': job_id,
        'titulo': job['titulo'],
        'orient': job['orient'],
        'total': total,
        'done': done,
        'errors': errors,
        'finished': (done + errors) == total,
        'variants': job['variants'],
    })


# ─────────────────────────────────────────────────────────────────────────────
#  CAMPANHAS
# ─────────────────────────────────────────────────────────────────────────────

def _campanha_qs(user):
    """Queryset de campanhas visíveis para o usuário."""
    if user.is_owner():
        return Campanha.objects.select_related('franqueado').prefetch_related('config_cupom')
    return Campanha.objects.filter(franqueado=user).select_related('franqueado').prefetch_related('config_cupom')


@login_required
def campanha_list_view(request):
    if not (request.user.is_owner() or request.user.is_franchisee()):
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')

    campanhas = _campanha_qs(request.user).order_by('-criado_em')
    # Atualiza status de expiradas automaticamente
    now = timezone.now()
    campanhas.filter(status='ATIVA', data_fim__lt=now).update(status='ENCERRADA')
    campanhas = campanhas.all()  # re-evaluate
    return render(request, 'campanhas/campanha_list.html', {'campanhas': campanhas})


@login_required
def campanha_create_view(request):
    """Passo 1: cria a campanha (nome + tipo). Redireciona para configurar."""
    if not (request.user.is_owner() or request.user.is_franchisee()):
        messages.error(request, 'Acesso negado.')
        return redirect('dashboard')

    if request.method == 'POST':
        nome = request.POST.get('nome', '').strip()
        tipo = request.POST.get('tipo', 'CUPOM')
        if not nome:
            messages.error(request, 'Informe o nome da campanha.')
            return render(request, 'campanhas/campanha_create.html', {'nome': nome, 'tipo': tipo})

        franqueado = request.user if request.user.is_franchisee() else None
        if request.user.is_owner():
            fid = request.POST.get('franqueado')
            if fid:
                try:
                    franqueado = User.objects.get(pk=fid, role='FRANCHISEE')
                except User.DoesNotExist:
                    pass
            if franqueado is None:
                # owner pode criar sem franqueado? Usar o próprio owner como placeholder.
                franqueado = request.user

        campanha = Campanha.objects.create(nome=nome, tipo=tipo, franqueado=franqueado)

        # Cria configuração padrão conforme o tipo
        if tipo == 'CUPOM':
            CampanhaCupomConfig.objects.create(campanha=campanha)
        elif tipo == 'ROLETA':
            CampanhaRoletaConfig.objects.create(campanha=campanha)
        elif tipo == 'CARTA':
            CampanhaCartaConfig.objects.create(campanha=campanha)
        elif tipo == 'ALERTA':
            CampanhaAlertaConfig.objects.create(campanha=campanha)
        elif tipo == 'SORTEIO':
            CampanhaSorteioConfig.objects.create(campanha=campanha)

        messages.success(request, f'Campanha "{nome}" criada. Configure os detalhes abaixo.')
        return redirect('campanha_configure', pk=campanha.pk)

    franqueados = User.objects.filter(role='FRANCHISEE', is_active=True).order_by('first_name', 'username') if request.user.is_owner() else None
    return render(request, 'campanhas/campanha_create.html', {
        'franqueados': franqueados,
        'tipo_choices': Campanha.TIPO_CHOICES,
    })


@login_required
def campanha_configure_view(request, pk):
    """Passo 2: configura os detalhes específicos do tipo de campanha."""
    campanha = get_object_or_404(Campanha, pk=pk)
    user = request.user
    if not user.is_owner() and campanha.franqueado != user:
        messages.error(request, 'Sem permissão.')
        return redirect('campanha_list')

    if campanha.tipo == 'CUPOM':
        config, _ = CampanhaCupomConfig.objects.get_or_create(campanha=campanha)

        if request.method == 'POST':
            config.modo_codigo   = request.POST.get('modo_codigo', 'CODIGO_UNICO')
            config.codigo_unico  = request.POST.get('codigo_unico', '').strip()
            config.prefixo_codigo = request.POST.get('prefixo_codigo', '').strip()
            config.capturar_nome     = 'capturar_nome'     in request.POST
            config.capturar_cpf      = 'capturar_cpf'      in request.POST
            config.capturar_telefone = 'capturar_telefone' in request.POST
            config.capturar_endereco = 'capturar_endereco' in request.POST
            config.titulo_pagina    = request.POST.get('titulo_pagina', '').strip()
            config.descricao_pagina = request.POST.get('descricao_pagina', '').strip()
            config.cor_primaria     = request.POST.get('cor_primaria', '#0d6efd').strip()
            config.save()

            campanha.nome     = request.POST.get('nome', campanha.nome).strip() or campanha.nome
            campanha.data_fim  = request.POST.get('data_fim') or None
            campanha.status    = request.POST.get('status', campanha.status)
            campanha.save()

            messages.success(request, 'Campanha salva com sucesso!')
            return redirect('campanha_detail', pk=campanha.pk)

        return render(request, 'campanhas/campanha_configure_cupom.html', {
            'campanha': campanha,
            'config': config,
            'modo_choices': CampanhaCupomConfig.MODO_CHOICES,
            'status_choices': Campanha.STATUS_CHOICES,
        })

    if campanha.tipo == 'ROLETA':
        config, _ = CampanhaRoletaConfig.objects.get_or_create(campanha=campanha)
        premios = campanha.premios_roleta.all()

        if request.method == 'POST':
            action = request.POST.get('action', 'save_config')

            if action == 'add_premio':
                CampanhaRoletaPremio.objects.create(
                    campanha=campanha,
                    nome=request.POST.get('p_nome', 'Novo Prêmio').strip() or 'Novo Prêmio',
                    descricao=request.POST.get('p_descricao', '').strip(),
                    codigo_resgate=request.POST.get('p_codigo_resgate', '').strip(),
                    peso=int(request.POST.get('p_peso', 10) or 10),
                    quantidade_maxima=int(request.POST.get('p_quantidade_maxima')) if request.POST.get('p_quantidade_maxima') else None,
                    cor=request.POST.get('p_cor', '#f4a261').strip() or '#f4a261',
                    emoji=request.POST.get('p_emoji', '🎁').strip(),
                    eh_perdedor='p_eh_perdedor' in request.POST,
                    ativo=True,
                    ordem=int(request.POST.get('p_ordem', 0) or 0),
                )
                messages.success(request, 'Prêmio adicionado.')
                return redirect('campanha_configure', pk=campanha.pk)

            elif action == 'edit_premio':
                pid = request.POST.get('p_id')
                try:
                    premio = campanha.premios_roleta.get(pk=pid)
                    premio.nome = request.POST.get('p_nome', premio.nome).strip() or premio.nome
                    premio.descricao = request.POST.get('p_descricao', '').strip()
                    premio.codigo_resgate = request.POST.get('p_codigo_resgate', '').strip()
                    premio.peso = int(request.POST.get('p_peso', premio.peso) or premio.peso)
                    qmax = request.POST.get('p_quantidade_maxima')
                    premio.quantidade_maxima = int(qmax) if qmax else None
                    premio.cor = request.POST.get('p_cor', premio.cor).strip() or premio.cor
                    premio.emoji = request.POST.get('p_emoji', '🎁').strip()
                    premio.eh_perdedor = 'p_eh_perdedor' in request.POST
                    premio.ativo = 'p_ativo' in request.POST
                    premio.ordem = int(request.POST.get('p_ordem', premio.ordem) or 0)
                    premio.save()
                    messages.success(request, 'Prêmio atualizado.')
                except CampanhaRoletaPremio.DoesNotExist:
                    messages.error(request, 'Prêmio não encontrado.')
                return redirect('campanha_configure', pk=campanha.pk)

            elif action == 'delete_premio':
                pid = request.POST.get('p_id')
                campanha.premios_roleta.filter(pk=pid).delete()
                messages.success(request, 'Prêmio removido.')
                return redirect('campanha_configure', pk=campanha.pk)

            else:  # save_config
                config.max_jogadas_por_ip_por_dia = int(request.POST.get('max_jogadas_por_ip_por_dia', 1) or 1)
                mjt = request.POST.get('max_jogadas_total_por_ip')
                config.max_jogadas_total_por_ip = int(mjt) if mjt else None
                config.capturar_nome     = 'capturar_nome'     in request.POST
                config.capturar_cpf      = 'capturar_cpf'      in request.POST
                config.capturar_telefone = 'capturar_telefone' in request.POST
                config.capturar_endereco = 'capturar_endereco' in request.POST
                config.titulo_pagina    = request.POST.get('titulo_pagina', '').strip()
                config.descricao_pagina = request.POST.get('descricao_pagina', '').strip()
                config.cor_primaria     = request.POST.get('cor_primaria', '#e63946').strip()
                config.texto_botao_girar = request.POST.get('texto_botao_girar', 'GIRAR!').strip() or 'GIRAR!'
                config.texto_sem_premio = request.POST.get('texto_sem_premio', '').strip() or config.texto_sem_premio
                config.save()

                campanha.nome    = request.POST.get('nome', campanha.nome).strip() or campanha.nome
                campanha.data_fim = request.POST.get('data_fim') or None
                campanha.status  = request.POST.get('status', campanha.status)
                campanha.save()

                messages.success(request, 'Roleta configurada com sucesso!')
                return redirect('campanha_detail', pk=campanha.pk)

        return render(request, 'campanhas/campanha_configure_roleta.html', {
            'campanha': campanha,
            'config': config,
            'premios': premios,
            'status_choices': Campanha.STATUS_CHOICES,
        })

    if campanha.tipo == 'CARTA':
        config, _ = CampanhaCartaConfig.objects.get_or_create(campanha=campanha)
        premios = campanha.premios_roleta.all()

        if request.method == 'POST':
            action = request.POST.get('action', 'save_config')

            if action == 'add_premio':
                CampanhaRoletaPremio.objects.create(
                    campanha=campanha,
                    nome=request.POST.get('p_nome', 'Novo Prêmio').strip() or 'Novo Prêmio',
                    descricao=request.POST.get('p_descricao', '').strip(),
                    codigo_resgate=request.POST.get('p_codigo_resgate', '').strip(),
                    peso=int(request.POST.get('p_peso', 10) or 10),
                    quantidade_maxima=int(request.POST.get('p_quantidade_maxima')) if request.POST.get('p_quantidade_maxima') else None,
                    cor=request.POST.get('p_cor', '#e63946').strip() or '#e63946',
                    emoji=request.POST.get('p_emoji', '🎁').strip(),
                    eh_perdedor='p_eh_perdedor' in request.POST,
                    ativo=True,
                    ordem=int(request.POST.get('p_ordem', 0) or 0),
                )
                messages.success(request, 'Carta adicionada.')
                return redirect('campanha_configure', pk=campanha.pk)

            elif action == 'edit_premio':
                pid = request.POST.get('p_id')
                try:
                    premio = campanha.premios_roleta.get(pk=pid)
                    premio.nome = request.POST.get('p_nome', premio.nome).strip() or premio.nome
                    premio.descricao = request.POST.get('p_descricao', '').strip()
                    premio.codigo_resgate = request.POST.get('p_codigo_resgate', '').strip()
                    premio.peso = int(request.POST.get('p_peso', premio.peso) or premio.peso)
                    qmax = request.POST.get('p_quantidade_maxima')
                    premio.quantidade_maxima = int(qmax) if qmax else None
                    premio.cor = request.POST.get('p_cor', premio.cor).strip() or premio.cor
                    premio.emoji = request.POST.get('p_emoji', '🎁').strip()
                    premio.eh_perdedor = 'p_eh_perdedor' in request.POST
                    premio.ativo = 'p_ativo' in request.POST
                    premio.ordem = int(request.POST.get('p_ordem', premio.ordem) or 0)
                    premio.save()
                    messages.success(request, 'Carta atualizada.')
                except CampanhaRoletaPremio.DoesNotExist:
                    messages.error(request, 'Carta não encontrada.')
                return redirect('campanha_configure', pk=campanha.pk)

            elif action == 'delete_premio':
                pid = request.POST.get('p_id')
                campanha.premios_roleta.filter(pk=pid).delete()
                messages.success(request, 'Carta removida.')
                return redirect('campanha_configure', pk=campanha.pk)

            else:  # save_config
                config.max_jogadas_por_ip_por_dia = int(request.POST.get('max_jogadas_por_ip_por_dia', 1) or 1)
                mjt = request.POST.get('max_jogadas_total_por_ip')
                config.max_jogadas_total_por_ip = int(mjt) if mjt else None
                config.capturar_nome     = 'capturar_nome'     in request.POST
                config.capturar_cpf      = 'capturar_cpf'      in request.POST
                config.capturar_telefone = 'capturar_telefone' in request.POST
                config.capturar_endereco = 'capturar_endereco' in request.POST
                config.titulo_pagina    = request.POST.get('titulo_pagina', '').strip()
                config.descricao_pagina = request.POST.get('descricao_pagina', '').strip()
                config.cor_primaria     = request.POST.get('cor_primaria', '#1a1a2e').strip()
                config.cor_verso_carta  = request.POST.get('cor_verso_carta', '#16213e').strip()
                config.cor_frente_carta = request.POST.get('cor_frente_carta', '#e63946').strip()
                config.texto_verso_carta = request.POST.get('texto_verso_carta', 'Vire a carta!').strip() or 'Vire a carta!'
                config.texto_botao_virar = request.POST.get('texto_botao_virar', 'Virar a Carta!').strip() or 'Virar a Carta!'
                config.texto_sem_premio = request.POST.get('texto_sem_premio', '').strip() or config.texto_sem_premio
                if 'logo' in request.FILES:
                    config.logo = request.FILES['logo']
                config.save()

                campanha.nome    = request.POST.get('nome', campanha.nome).strip() or campanha.nome
                campanha.data_fim = request.POST.get('data_fim') or None
                campanha.status  = request.POST.get('status', campanha.status)
                campanha.save()

                messages.success(request, 'Virar a Carta configurada com sucesso!')
                return redirect('campanha_detail', pk=campanha.pk)

        return render(request, 'campanhas/campanha_configure_carta.html', {
            'campanha': campanha,
            'config': config,
            'premios': premios,
            'status_choices': Campanha.STATUS_CHOICES,
        })

    if campanha.tipo == 'ALERTA':
        config, _ = CampanhaAlertaConfig.objects.get_or_create(campanha=campanha)
        campos = campanha.campos_alerta.filter(ativo=True).order_by('ordem', 'id')

        CAMPO_TIPO_CHOICES = CampanhaAlertaCampo.TIPO_CHOICES

        SUGESTOES = [
            {'rotulo': 'Tipo de imóvel', 'tipo': 'SELECT', 'obrigatorio': True,
             'opcoes': 'Apartamento\nCasa\nSala Comercial\nTerreno\nGalpão / Industrial\nSitio / Chácara\nOutro'},
            {'rotulo': 'Finalidade', 'tipo': 'SELECT', 'obrigatorio': True,
             'opcoes': 'Compra\nLocação'},
            {'rotulo': 'Número de dormitórios', 'tipo': 'SELECT', 'obrigatorio': False,
             'opcoes': '1\n2\n3\n4 ou mais'},
            {'rotulo': 'Orçamento máximo (R$)', 'tipo': 'MOEDA', 'obrigatorio': False, 'opcoes': ''},
            {'rotulo': 'Bairros ou regiões de interesse', 'tipo': 'TEXTO', 'obrigatorio': False, 'opcoes': ''},
            {'rotulo': 'Prazo para fechar negócio', 'tipo': 'SELECT', 'obrigatorio': False,
             'opcoes': 'Imediato\nAté 3 meses\nAté 6 meses\nAté 1 ano\nSem prazo definido'},
            {'rotulo': 'Observações ou características importantes', 'tipo': 'TEXTO',
             'obrigatorio': False, 'opcoes': ''},
        ]

        if request.method == 'POST':
            action = request.POST.get('action', 'save_config')

            if action == 'add_campo':
                CampanhaAlertaCampo.objects.create(
                    campanha=campanha,
                    tipo=request.POST.get('c_tipo', 'TEXTO'),
                    rotulo=request.POST.get('c_rotulo', '').strip() or 'Novo Campo',
                    placeholder=request.POST.get('c_placeholder', '').strip(),
                    opcoes=request.POST.get('c_opcoes', '').strip(),
                    obrigatorio='c_obrigatorio' in request.POST,
                    ordem=int(request.POST.get('c_ordem', 0) or 0),
                )
                messages.success(request, 'Campo adicionado.')
                return redirect('campanha_configure', pk=campanha.pk)

            elif action == 'edit_campo':
                cid = request.POST.get('c_id')
                try:
                    campo = campanha.campos_alerta.get(pk=cid)
                    campo.tipo = request.POST.get('c_tipo', campo.tipo)
                    campo.rotulo = request.POST.get('c_rotulo', campo.rotulo).strip() or campo.rotulo
                    campo.placeholder = request.POST.get('c_placeholder', '').strip()
                    campo.opcoes = request.POST.get('c_opcoes', '').strip()
                    campo.obrigatorio = 'c_obrigatorio' in request.POST
                    campo.ativo = 'c_ativo' in request.POST
                    campo.ordem = int(request.POST.get('c_ordem', campo.ordem) or 0)
                    campo.save()
                    messages.success(request, 'Campo atualizado.')
                except CampanhaAlertaCampo.DoesNotExist:
                    messages.error(request, 'Campo não encontrado.')
                return redirect('campanha_configure', pk=campanha.pk)

            elif action == 'delete_campo':
                campanha.campos_alerta.filter(pk=request.POST.get('c_id')).delete()
                messages.success(request, 'Campo removido.')
                return redirect('campanha_configure', pk=campanha.pk)

            elif action == 'add_sugestoes':
                # Adiciona todos os campos sugeridos de uma vez
                ordem_atual = campanha.campos_alerta.count() * 10
                for sg in SUGESTOES:
                    CampanhaAlertaCampo.objects.create(
                        campanha=campanha,
                        tipo=sg['tipo'],
                        rotulo=sg['rotulo'],
                        opcoes=sg.get('opcoes', ''),
                        obrigatorio=sg.get('obrigatorio', False),
                        ordem=ordem_atual,
                    )
                    ordem_atual += 10
                messages.success(request, 'Campos sugeridos de imobiliária adicionados.')
                return redirect('campanha_configure', pk=campanha.pk)

            else:  # save_config
                config.titulo_pagina      = request.POST.get('titulo_pagina', '').strip()
                config.subtitulo_pagina   = request.POST.get('subtitulo_pagina', '').strip()
                config.descricao_pagina   = request.POST.get('descricao_pagina', '').strip()
                config.mensagem_sucesso   = request.POST.get('mensagem_sucesso', config.mensagem_sucesso).strip() or config.mensagem_sucesso
                config.whatsapp_contato   = request.POST.get('whatsapp_contato', '').strip()
                config.cor_primaria       = request.POST.get('cor_primaria', '#1a1a2e').strip()
                config.cor_destaque       = request.POST.get('cor_destaque', '#e63946').strip()
                config.capturar_nome      = 'capturar_nome' in request.POST
                config.capturar_telefone  = 'capturar_telefone' in request.POST
                config.capturar_email     = 'capturar_email' in request.POST
                if 'logo' in request.FILES:
                    config.logo = request.FILES['logo']
                config.save()

                campanha.nome    = request.POST.get('nome', campanha.nome).strip() or campanha.nome
                campanha.data_fim = request.POST.get('data_fim') or None
                campanha.status  = request.POST.get('status', campanha.status)
                campanha.save()

                messages.success(request, 'Alerta Inteligente configurado com sucesso!')
                return redirect('campanha_detail', pk=campanha.pk)

        return render(request, 'campanhas/campanha_configure_alerta.html', {
            'campanha': campanha,
            'config': config,
            'campos': campos,
            'campo_tipo_choices': CAMPO_TIPO_CHOICES,
            'sugestoes': SUGESTOES,
            'status_choices': Campanha.STATUS_CHOICES,
        })

    if campanha.tipo == 'SORTEIO':
        return _campanha_configure_sorteio(request, campanha)

    messages.warning(request, 'Tipo de campanha ainda não suportado.')
    return redirect('campanha_list')


def _campanha_configure_sorteio(request, campanha):
    """Handler interno para configurar campanha do tipo SORTEIO."""
    config, _ = CampanhaSorteioConfig.objects.get_or_create(campanha=campanha)
    if request.method == 'POST':
        config.titulo_pagina    = request.POST.get('titulo_pagina', '').strip()
        config.descricao_pagina = request.POST.get('descricao_pagina', '').strip()
        config.mensagem_sucesso = request.POST.get('mensagem_sucesso', config.mensagem_sucesso).strip() or config.mensagem_sucesso
        config.capturar_telefone  = 'capturar_telefone' in request.POST
        config.capturar_endereco  = 'capturar_endereco' in request.POST
        config.bloquear_duplicados_cpf = 'bloquear_duplicados_cpf' in request.POST
        config.bloquear_duplicados_ip  = 'bloquear_duplicados_ip' in request.POST
        config.cor_primaria     = request.POST.get('cor_primaria', '#6366f1').strip()
        data_sorteio_raw = request.POST.get('data_sorteio', '').strip()
        if data_sorteio_raw:
            from django.utils.dateparse import parse_datetime
            config.data_sorteio = parse_datetime(data_sorteio_raw)
        else:
            config.data_sorteio = None
        if 'foto_item' in request.FILES:
            config.foto_item = request.FILES['foto_item']
        config.save()

        campanha.nome    = request.POST.get('nome', campanha.nome).strip() or campanha.nome
        campanha.data_fim = request.POST.get('data_fim') or None
        campanha.status  = request.POST.get('status', campanha.status)
        campanha.save()
        messages.success(request, 'Sorteio configurado com sucesso!')
        return redirect('campanha_detail', pk=campanha.pk)

    return render(request, 'campanhas/campanha_configure_sorteio.html', {
        'campanha': campanha,
        'config': config,
        'status_choices': Campanha.STATUS_CHOICES,
    })


@login_required
def campanha_detail_view(request, pk):
    campanha = get_object_or_404(Campanha, pk=pk)
    user = request.user
    if not user.is_owner() and campanha.franqueado != user:
        messages.error(request, 'Sem permissão.')
        return redirect('campanha_list')

    leads_count = campanha.leads.count()
    config = getattr(campanha, 'config_cupom', None)
    config_roleta = getattr(campanha, 'config_roleta', None)
    config_carta = getattr(campanha, 'config_carta', None)
    config_alerta = getattr(campanha, 'config_alerta', None)
    config_sorteio = getattr(campanha, 'config_sorteio', None)
    jogadas_count = campanha.jogadas.count() if campanha.tipo in ('ROLETA', 'CARTA') else 0
    ganhadores_count = campanha.jogadas.filter(ganhou=True).count() if campanha.tipo in ('ROLETA', 'CARTA') else 0
    alerta_leads_count = campanha.leads_alerta.count() if campanha.tipo == 'ALERTA' else 0
    sorteio_participantes_count = campanha.participantes_sorteio.count() if campanha.tipo == 'SORTEIO' else 0
    sorteio_ativos_count = campanha.participantes_sorteio.filter(ativo_sorteio=True).count() if campanha.tipo == 'SORTEIO' else 0
    return render(request, 'campanhas/campanha_detail.html', {
        'campanha': campanha,
        'config': config,
        'config_roleta': config_roleta,
        'config_carta': config_carta,
        'config_alerta': config_alerta,
        'config_sorteio': config_sorteio,
        'leads_count': leads_count,
        'jogadas_count': jogadas_count,
        'ganhadores_count': ganhadores_count,
        'alerta_leads_count': alerta_leads_count,
        'sorteio_participantes_count': sorteio_participantes_count,
        'sorteio_ativos_count': sorteio_ativos_count,
    })


@login_required
def campanha_delete_view(request, pk):
    campanha = get_object_or_404(Campanha, pk=pk)
    user = request.user
    if not user.is_owner() and campanha.franqueado != user:
        return JsonResponse({'success': False, 'error': 'Sem permissão.'}, status=403)
    campanha.delete()
    messages.success(request, 'Campanha excluída.')
    return JsonResponse({'success': True})


@login_required
def campanha_toggle_status_view(request, pk):
    """Ativa / Encerra a campanha via POST."""
    if request.method != 'POST':
        return JsonResponse({'error': 'Method not allowed'}, status=405)
    campanha = get_object_or_404(Campanha, pk=pk)
    user = request.user
    if not user.is_owner() and campanha.franqueado != user:
        return JsonResponse({'success': False, 'error': 'Sem permissão.'}, status=403)
    new_status = request.POST.get('status', 'ATIVA')
    if new_status in ('ATIVA', 'ENCERRADA', 'RASCUNHO'):
        campanha.status = new_status
        campanha.save()
    return JsonResponse({'success': True, 'status': campanha.status})


@login_required
def campanha_leads_view(request, pk):
    campanha = get_object_or_404(Campanha, pk=pk)
    user = request.user
    if not user.is_owner() and campanha.franqueado != user:
        messages.error(request, 'Sem permissão.')
        return redirect('campanha_list')

    # Campanhas do tipo ALERTA têm seus próprios leads
    if campanha.tipo == 'ALERTA':
        return redirect('campanha_alerta_leads', pk=pk)

    # Campanhas do tipo SORTEIO vão para a tela do sorteio
    if campanha.tipo == 'SORTEIO':
        return redirect('campanha_sorteio_draw', pk=pk)

    leads = campanha.leads.order_by('-criado_em')

    # CSV export
    if request.GET.get('export') == 'csv':
        import csv
        from django.http import HttpResponse
        response = HttpResponse(content_type='text/csv; charset=utf-8')
        fname = f'leads_{campanha.pk}.csv'
        response['Content-Disposition'] = f'attachment; filename="{fname}"'
        writer = csv.writer(response)
        writer.writerow(['#', 'Nome', 'CPF', 'Telefone', 'Endereço', 'Código Cupom', 'Data'])
        for i, lead in enumerate(leads, 1):
            writer.writerow([i, lead.nome, lead.cpf, lead.telefone,
                             lead.endereco, lead.codigo_cupom,
                             lead.criado_em.strftime('%d/%m/%Y %H:%M')])
        return response

    return render(request, 'campanhas/campanha_leads.html', {
        'campanha': campanha,
        'leads': leads,
    })


# ── LANDING PAGE (pública) ────────────────────────────────────────────────────

def campanha_landing_view(request, token):
    """Página pública da campanha — sem login."""
    from django.utils.crypto import get_random_string
    campanha = get_object_or_404(Campanha, token=token)

    # Verificar se está ativa
    if campanha.expirada and campanha.status == 'ATIVA':
        campanha.status = 'ENCERRADA'
        campanha.save(update_fields=['status'])

    encerrada = not campanha.is_ativa
    config = getattr(campanha, 'config_cupom', None)

    # Dispatch por tipo
    if campanha.tipo == 'ROLETA':
        return _campanha_roleta_landing(request, campanha, encerrada)

    if campanha.tipo == 'CARTA':
        return _campanha_carta_landing(request, campanha, encerrada)

    if campanha.tipo == 'ALERTA':
        return _campanha_alerta_landing(request, campanha, encerrada)

    if campanha.tipo == 'SORTEIO':
        return _campanha_sorteio_landing(request, campanha, encerrada)
        ip = request.META.get('HTTP_X_FORWARDED_FOR', request.META.get('REMOTE_ADDR', '')).split(',')[0].strip() or None

        # ── Deduplicação por IP (modos com lead) ──────────────────────────────
        if config.modo_codigo in ('CODIGO_UNICO', 'CODIGO_POR_CLIENTE') and ip:
            lead_existente = campanha.leads.filter(ip=ip).first()
            if lead_existente:
                return render(request, 'campanhas/campanha_resgate_sucesso.html', {
                    'campanha': campanha,
                    'config': config,
                    'codigo': lead_existente.codigo_cupom,
                    'ja_resgatado': True,
                })

        nome     = request.POST.get('nome', '').strip()
        cpf      = request.POST.get('cpf', '').strip()
        telefone = request.POST.get('telefone', '').strip()
        endereco = request.POST.get('endereco', '').strip()

        # Validação de CPF
        if (config and config.capturar_cpf) and cpf:
            if not _validar_cpf(cpf):
                return render(request, 'campanhas/campanha_landing.html', {
                    'campanha': campanha,
                    'config': config,
                    'erro': 'CPF inválido. Verifique os dígitos e tente novamente.',
                    'form_nome': nome, 'form_cpf': cpf,
                    'form_telefone': telefone, 'form_endereco': endereco,
                })

        # Determinar o código a entregar
        if config.modo_codigo == 'CODIGO_UNICO':
            codigo = config.codigo_unico
        elif config.modo_codigo == 'CODIGO_POR_CLIENTE':
            prefix = (config.prefixo_codigo.rstrip('-') + '-') if config.prefixo_codigo else ''
            codigo = prefix + get_random_string(6).upper()
        else:  # SEM_LEAD
            codigo = config.codigo_unico

        # Salvar lead (sempre salva quando há captura ou código individual)
        if config.modo_codigo != 'SEM_LEAD' or config.captura_algum_dado:
            CampanhaLead.objects.create(
                campanha=campanha,
                nome=nome,
                cpf=cpf,
                telefone=telefone,
                endereco=endereco,
                codigo_cupom=codigo,
                ip=ip,
            )

        return render(request, 'campanhas/campanha_resgate_sucesso.html', {
            'campanha': campanha,
            'config': config,
            'codigo': codigo,
            'ja_resgatado': False,
        })

    # GET: mostrar o formulário (ou código direto no modo SEM_LEAD)
    exibir_codigo_direto = (
        config and config.modo_codigo == 'SEM_LEAD' and not encerrada
    )
    if exibir_codigo_direto and config:
        codigo_direto = config.codigo_unico
    else:
        codigo_direto = None

    return render(request, 'campanhas/campanha_landing.html', {
        'campanha': campanha,
        'config': config,
        'encerrada': encerrada,
        'exibir_codigo_direto': exibir_codigo_direto,
        'codigo_direto': codigo_direto,
    })


# ─────────────────────────────────────────────────────────────────────────────
#  ROLETA — views públicas e autenticadas
# ─────────────────────────────────────────────────────────────────────────────

def _campanha_roleta_landing(request, campanha, encerrada):
    """Renderiza a landing page da roleta (READ-ONLY — giro é via AJAX)."""
    config = getattr(campanha, 'config_roleta', None)
    premios = list(campanha.premios_roleta.filter(ativo=True).order_by('ordem', 'id'))
    premios_json = [
        {
            'id': p.pk,
            'nome': p.nome,
            'emoji': p.emoji,
            'cor': p.cor,
            'eh_perdedor': p.eh_perdedor,
        }
        for p in premios
    ]
    import json
    return render(request, 'campanhas/campanha_landing_roleta.html', {
        'campanha': campanha,
        'config': config,
        'premios': premios,
        'premios_json': json.dumps(premios_json),
        'encerrada': encerrada,
    })


@csrf_exempt
def campanha_spin_view(request, token):
    """AJAX POST — realiza uma jogada na roleta. Retorna JSON."""
    import random
    from django.utils import timezone

    if request.method != 'POST':
        return JsonResponse({'success': False, 'motivo': 'metodo_invalido'}, status=405)

    campanha = get_object_or_404(Campanha, token=token, tipo='ROLETA')

    if not campanha.is_ativa:
        return JsonResponse({'success': False, 'motivo': 'encerrada'})

    config = getattr(campanha, 'config_roleta', None)
    if not config:
        return JsonResponse({'success': False, 'motivo': 'sem_configuracao'})

    ip = request.META.get('HTTP_X_FORWARDED_FOR', request.META.get('REMOTE_ADDR', '')).split(',')[0].strip() or None

    # ── Verificar limite diário ───────────────────────────────────────────
    if ip and config.max_jogadas_por_ip_por_dia > 0:
        hoje = timezone.localdate()
        jogadas_hoje = campanha.jogadas.filter(ip=ip, criado_em__date=hoje).count()
        if jogadas_hoje >= config.max_jogadas_por_ip_por_dia:
            return JsonResponse({'success': False, 'motivo': 'limite_diario'})

    # ── Verificar limite total ────────────────────────────────────────────
    if ip and config.max_jogadas_total_por_ip:
        total = campanha.jogadas.filter(ip=ip).count()
        if total >= config.max_jogadas_total_por_ip:
            return JsonResponse({'success': False, 'motivo': 'limite_total'})

    # ── Sortear prêmio ────────────────────────────────────────────────────
    premios_ativos = list(campanha.premios_roleta.filter(ativo=True).order_by('ordem', 'id'))
    if not premios_ativos:
        return JsonResponse({'success': False, 'motivo': 'sem_premios'})

    disponiveis = [p for p in premios_ativos if not p.esgotado]
    if not disponiveis:
        return JsonResponse({'success': False, 'motivo': 'todos_esgotados'})

    pesos = [p.peso for p in disponiveis]
    winner = random.choices(disponiveis, weights=pesos, k=1)[0]

    # Segmento index na lista completa (para animar a roleta)
    try:
        segment_index = premios_ativos.index(winner)
    except ValueError:
        segment_index = 0

    # ── Criar jogada ──────────────────────────────────────────────────────
    jogada = CampanhaJogada.objects.create(
        campanha=campanha,
        ip=ip,
        premio=winner,
        ganhou=not winner.eh_perdedor,
    )

    precisa_lead = (not winner.eh_perdedor) and config.captura_algum_dado

    return JsonResponse({
        'success': True,
        'segment_index': segment_index,
        'jogada_id': jogada.pk,
        'ganhou': jogada.ganhou,
        'eh_perdedor': winner.eh_perdedor,
        'nome_premio': winner.nome,
        'emoji': winner.emoji,
        'codigo_resgate': winner.codigo_resgate if not winner.eh_perdedor else '',
        'precisa_lead': precisa_lead,
        'texto_sem_premio': config.texto_sem_premio,
    })


def _validar_cpf(cpf: str) -> bool:
    """Valida CPF brasileiro. Aceita '000.000.000-00' ou '00000000000'."""
    import re
    cpf = re.sub(r'\D', '', cpf)
    if len(cpf) != 11 or cpf == cpf[0] * 11:
        return False
    for i in range(9, 11):
        soma = sum(int(cpf[j]) * (i + 1 - j) for j in range(i))
        if (soma * 10 % 11) % 10 != int(cpf[i]):
            return False
    return True


@csrf_exempt
def campanha_roleta_lead_view(request, token, jogada_pk):
    """AJAX/POST público — salva os dados de lead após o ganhador preencher o formulário."""
    if request.method != 'POST':
        return JsonResponse({'success': False}, status=405)

    campanha = get_object_or_404(Campanha, token=token, tipo='ROLETA')
    jogada = get_object_or_404(CampanhaJogada, pk=jogada_pk, campanha=campanha)

    if jogada.lead_salvo:
        return JsonResponse({'success': True, 'ja_salvo': True})

    config = getattr(campanha, 'config_roleta', None)

    jogada.nome     = request.POST.get('nome', '').strip() if (config and config.capturar_nome)     else ''
    jogada.cpf      = request.POST.get('cpf',  '').strip() if (config and config.capturar_cpf)      else ''
    jogada.telefone = request.POST.get('telefone', '').strip() if (config and config.capturar_telefone) else ''
    jogada.endereco = request.POST.get('endereco', '').strip() if (config and config.capturar_endereco) else ''

    # Validação de CPF
    if (config and config.capturar_cpf) and jogada.cpf:
        if not _validar_cpf(jogada.cpf):
            return JsonResponse({'success': False, 'error': 'CPF inválido. Verifique os dígitos e tente novamente.'}, status=400)

    jogada.lead_salvo = True
    jogada.save(update_fields=['nome', 'cpf', 'telefone', 'endereco', 'lead_salvo'])

    # Também registrar como CampanhaLead para aparecer na aba "Leads"
    CampanhaLead.objects.create(
        campanha=campanha,
        nome=jogada.nome,
        cpf=jogada.cpf,
        telefone=jogada.telefone,
        endereco=jogada.endereco,
        codigo_cupom=jogada.premio.codigo_resgate if jogada.premio else '',
        ip=jogada.ip,
    )

    return JsonResponse({'success': True, 'ja_salvo': False})


@login_required
def campanha_jogadas_view(request, pk):
    """Painel autenticado: lista todas as jogadas da roleta com filtros e exportação CSV."""
    import csv

    campanha = get_object_or_404(Campanha, pk=pk)
    user = request.user
    if not user.is_owner() and campanha.franqueado != user:
        messages.error(request, 'Sem permissão.')
        return redirect('campanha_list')

    jogadas = campanha.jogadas.select_related('premio').order_by('-criado_em')

    # Exportação CSV
    if request.GET.get('export') == 'csv':
        from django.http import HttpResponse
        import io
        output = io.StringIO()
        output.write('\ufeff')  # BOM para Excel abrir sem problemas de encoding
        writer = csv.writer(output)
        writer.writerow(['#', 'Data', 'IP', 'Premio', 'Ganhou', 'Nome', 'CPF', 'Telefone', 'Endereco'])
        for j in jogadas:
            writer.writerow([
                j.pk,
                j.criado_em.strftime('%d/%m/%Y %H:%M'),
                j.ip or '',
                j.premio.nome if j.premio else '',
                'Sim' if j.ganhou else 'Nao',
                j.nome or '',
                j.cpf or '',
                j.telefone or '',
                j.endereco or '',
            ])
        response = HttpResponse(output.getvalue(), content_type='text/csv')
        response['Content-Disposition'] = f'attachment; filename="jogadas_{campanha.pk}.csv"'
        return response

    return render(request, 'campanhas/campanha_jogadas.html', {
        'campanha': campanha,
        'jogadas': jogadas,
        'total': jogadas.count(),
        'ganhadores': jogadas.filter(ganhou=True).count(),
    })


# ─────────────────────────────────────────────────────────────────────────────
#  VIRAR A CARTA — views públicas e autenticadas
# ─────────────────────────────────────────────────────────────────────────────

def _campanha_carta_landing(request, campanha, encerrada):
    """Renderiza a landing page do jogo Virar a Carta (flip é via AJAX)."""
    import json
    config = getattr(campanha, 'config_carta', None)
    premios = list(campanha.premios_roleta.filter(ativo=True).order_by('ordem', 'id'))
    premios_json = json.dumps([
        {'id': p.pk, 'nome': p.nome, 'emoji': p.emoji, 'cor': p.cor, 'eh_perdedor': p.eh_perdedor}
        for p in premios
    ])
    logo_url = None
    if config and config.logo:
        try:
            logo_url = config.logo.url
        except Exception:
            pass
    return render(request, 'campanhas/campanha_landing_carta.html', {
        'campanha': campanha,
        'config': config,
        'premios': premios,
        'premios_json': premios_json,
        'encerrada': encerrada,
        'logo_url': logo_url,
    })


# ─────────────────────────────────────────────────────────────────────────────
#  ALERTA INTELIGENTE — landing pública e leads
# ─────────────────────────────────────────────────────────────────────────────

def _campanha_alerta_landing(request, campanha, encerrada):
    """Renderiza e processa o formulário público do Alerta Inteligente."""
    config = getattr(campanha, 'config_alerta', None)
    campos = list(campanha.campos_alerta.filter(ativo=True).order_by('ordem', 'id'))

    if request.method == 'POST' and not encerrada and config:
        ip = request.META.get('HTTP_X_FORWARDED_FOR',
                              request.META.get('REMOTE_ADDR', '')).split(',')[0].strip() or None

        nome     = request.POST.get('nome', '').strip()
        telefone = request.POST.get('telefone', '').strip()
        email    = request.POST.get('email', '').strip()

        # Coletar respostas dinâmicas
        respostas = {}
        for campo in campos:
            chave = f'campo_{campo.pk}'
            if campo.tipo == 'MULTISELECT':
                respostas[str(campo.pk)] = request.POST.getlist(chave)
            else:
                respostas[str(campo.pk)] = request.POST.get(chave, '').strip()

        CampanhaAlertaLead.objects.create(
            campanha=campanha,
            nome=nome,
            telefone=telefone,
            email=email,
            respostas=respostas,
            ip=ip,
        )

        # Montar URL de WhatsApp com resumo (se configurado)
        wa_url = None
        if config.whatsapp_contato:
            numero = ''.join(filter(str.isdigit, config.whatsapp_contato))
            import urllib.parse
            msg_parts = [f'Olá! Me cadastrei no alerta "{campanha.nome}".']
            if nome:
                msg_parts.append(f'Nome: {nome}')
            for campo in campos:
                valor = respostas.get(str(campo.pk), '')
                if isinstance(valor, list):
                    valor = ', '.join(valor)
                if valor:
                    msg_parts.append(f'{campo.rotulo}: {valor}')
            wa_url = f'https://wa.me/{numero}?text={urllib.parse.quote(chr(10).join(msg_parts))}'

        return render(request, 'campanhas/campanha_alerta_sucesso.html', {
            'campanha': campanha,
            'config': config,
            'wa_url': wa_url,
        })

    return render(request, 'campanhas/campanha_landing_alerta.html', {
        'campanha': campanha,
        'config': config,
        'campos': campos,
        'encerrada': encerrada,
    })


# ─────────────────────────────────────────────────────────────────────────────
#  SORTEIO — views públicas e autenticadas
# ─────────────────────────────────────────────────────────────────────────────

def _campanha_sorteio_landing(request, campanha, encerrada):
    """Formulário público de inscrição no sorteio."""
    config = getattr(campanha, 'config_sorteio', None)

    if request.method == 'POST' and not encerrada and config:
        ip = request.META.get('HTTP_X_FORWARDED_FOR',
                              request.META.get('REMOTE_ADDR', '')).split(',')[0].strip() or None

        nome     = request.POST.get('nome', '').strip()
        cpf_raw  = request.POST.get('cpf', '').strip()
        cpf      = ''.join(filter(str.isdigit, cpf_raw))  # normaliza para só dígitos
        telefone = request.POST.get('telefone', '').strip()
        endereco = request.POST.get('endereco', '').strip()

        erro = None

        # Validação CPF
        if not _validar_cpf(cpf):
            erro = 'CPF inválido. Verifique os dígitos e tente novamente.'

        # Bloquear duplicado por CPF
        if not erro and config.bloquear_duplicados_cpf:
            if campanha.participantes_sorteio.filter(cpf=cpf).exists():
                erro = 'Este CPF já está inscrito neste sorteio.'

        # Bloquear duplicado por IP
        if not erro and config.bloquear_duplicados_ip and ip:
            if campanha.participantes_sorteio.filter(ip=ip).exists():
                erro = 'Você já está inscrito neste sorteio.'

        if erro:
            return render(request, 'campanhas/campanha_landing_sorteio.html', {
                'campanha': campanha, 'config': config, 'encerrada': encerrada,
                'erro': erro,
                'form_nome': nome, 'form_cpf': cpf_raw,
                'form_telefone': telefone, 'form_endereco': endereco,
            })

        CampanhaParticipanteSorteio.objects.create(
            campanha=campanha,
            nome=nome,
            cpf=cpf,
            telefone=telefone,
            endereco=endereco,
            ip=ip,
        )

        return render(request, 'campanhas/campanha_sorteio_sucesso.html', {
            'campanha': campanha,
            'config': config,
        })

    return render(request, 'campanhas/campanha_landing_sorteio.html', {
        'campanha': campanha,
        'config': config,
        'encerrada': encerrada,
    })


@login_required
def campanha_sorteio_draw_view(request, pk):
    """Tela do sorteio animado — apenas autenticados (franqueado/owner)."""
    campanha = get_object_or_404(Campanha, pk=pk, tipo='SORTEIO')
    user = request.user
    if not user.is_owner() and campanha.franqueado != user:
        messages.error(request, 'Sem permissão.')
        return redirect('campanha_list')

    config = getattr(campanha, 'config_sorteio', None)
    participantes = campanha.participantes_sorteio.all().order_by('nome')
    total = campanha.participantes_sorteio.filter(ativo_sorteio=True).count()
    inativos = campanha.participantes_sorteio.filter(ativo_sorteio=False).count()

    return render(request, 'campanhas/campanha_sorteio_draw.html', {
        'campanha': campanha,
        'config': config,
        'participantes': participantes,
        'total': total,
        'inativos': inativos,
    })


@login_required
def campanha_sorteio_participantes_view(request, pk):
    """JSON: lista participantes ativos para o sorteador."""
    campanha = get_object_or_404(Campanha, pk=pk, tipo='SORTEIO')
    user = request.user
    if not user.is_owner() and campanha.franqueado != user:
        return JsonResponse({'error': 'Sem permissão.'}, status=403)

    qs = campanha.participantes_sorteio.filter(ativo_sorteio=True).values('id', 'nome', 'cpf', 'telefone')
    return JsonResponse({'participantes': list(qs)})


@login_required
def campanha_sorteio_toggle_participante_view(request, pk, participante_pk):
    """Toggle ativo_sorteio de um participante."""
    if request.method != 'POST':
        return JsonResponse({'error': 'Method not allowed.'}, status=405)
    campanha = get_object_or_404(Campanha, pk=pk, tipo='SORTEIO')
    user = request.user
    if not user.is_owner() and campanha.franqueado != user:
        return JsonResponse({'error': 'Sem permissão.'}, status=403)

    participante = get_object_or_404(CampanhaParticipanteSorteio, pk=participante_pk, campanha=campanha)
    participante.ativo_sorteio = not participante.ativo_sorteio
    participante.save(update_fields=['ativo_sorteio'])
    return JsonResponse({'success': True, 'ativo_sorteio': participante.ativo_sorteio})


@login_required
def campanha_alerta_leads_view(request, pk):
    """Listagem e export CSV dos leads do Alerta Inteligente."""
    campanha = get_object_or_404(Campanha, pk=pk, tipo='ALERTA')
    user = request.user
    if not user.is_owner() and campanha.franqueado != user:
        messages.error(request, 'Sem permissão.')
        return redirect('campanha_list')

    campos = list(campanha.campos_alerta.filter(ativo=True).order_by('ordem', 'id'))
    leads  = campanha.leads_alerta.order_by('-criado_em')

    if request.GET.get('export') == 'csv':
        import csv
        from django.http import HttpResponse
        response = HttpResponse(content_type='text/csv; charset=utf-8')
        fname = f'alerta_leads_{campanha.pk}.csv'
        response['Content-Disposition'] = f'attachment; filename="{fname}"'
        writer = csv.writer(response)
        cabecalho = ['#', 'Nome', 'Telefone', 'E-mail'] + [c.rotulo for c in campos] + ['Data']
        writer.writerow(cabecalho)
        for i, lead in enumerate(leads, 1):
            linha = [i, lead.nome, lead.telefone, lead.email]
            for campo in campos:
                v = lead.respostas.get(str(campo.pk), '')
                if isinstance(v, list):
                    v = ', '.join(v)
                linha.append(v)
            linha.append(lead.criado_em.strftime('%d/%m/%Y %H:%M'))
            writer.writerow(linha)
        return response

    return render(request, 'campanhas/campanha_alerta_leads.html', {
        'campanha': campanha,
        'campos': campos,
        'leads': leads,
    })


@csrf_exempt
def campanha_carta_flip_view(request, token):
    """AJAX POST — realiza uma jogada no Virar a Carta. Retorna JSON."""
    import random
    from django.utils import timezone

    if request.method != 'POST':
        return JsonResponse({'success': False, 'motivo': 'metodo_invalido'}, status=405)

    campanha = get_object_or_404(Campanha, token=token, tipo='CARTA')

    if not campanha.is_ativa:
        return JsonResponse({'success': False, 'motivo': 'encerrada'})

    config = getattr(campanha, 'config_carta', None)
    if not config:
        return JsonResponse({'success': False, 'motivo': 'sem_configuracao'})

    ip = request.META.get('HTTP_X_FORWARDED_FOR', request.META.get('REMOTE_ADDR', '')).split(',')[0].strip() or None

    # ── Verificar limite diário ───────────────────────────────────────────
    if ip and config.max_jogadas_por_ip_por_dia > 0:
        hoje = timezone.localdate()
        jogadas_hoje = campanha.jogadas.filter(ip=ip, criado_em__date=hoje).count()
        if jogadas_hoje >= config.max_jogadas_por_ip_por_dia:
            return JsonResponse({'success': False, 'motivo': 'limite_diario'})

    # ── Verificar limite total ────────────────────────────────────────────
    if ip and config.max_jogadas_total_por_ip:
        total = campanha.jogadas.filter(ip=ip).count()
        if total >= config.max_jogadas_total_por_ip:
            return JsonResponse({'success': False, 'motivo': 'limite_total'})

    # ── Sortear prêmio ────────────────────────────────────────────────────
    premios_ativos = list(campanha.premios_roleta.filter(ativo=True).order_by('ordem', 'id'))
    if not premios_ativos:
        return JsonResponse({'success': False, 'motivo': 'sem_premios'})

    disponiveis = [p for p in premios_ativos if not p.esgotado]
    if not disponiveis:
        return JsonResponse({'success': False, 'motivo': 'todos_esgotados'})

    pesos = [p.peso for p in disponiveis]
    winner = random.choices(disponiveis, weights=pesos, k=1)[0]

    # ── Criar jogada ──────────────────────────────────────────────────────
    jogada = CampanhaJogada.objects.create(
        campanha=campanha,
        ip=ip,
        premio=winner,
        ganhou=not winner.eh_perdedor,
    )

    precisa_lead = (not winner.eh_perdedor) and config.captura_algum_dado

    return JsonResponse({
        'success': True,
        'jogada_id': jogada.pk,
        'ganhou': jogada.ganhou,
        'eh_perdedor': winner.eh_perdedor,
        'nome_premio': winner.nome,
        'emoji': winner.emoji,
        'cor': winner.cor,
        'codigo_resgate': winner.codigo_resgate if not winner.eh_perdedor else '',
        'precisa_lead': precisa_lead,
        'texto_sem_premio': config.texto_sem_premio,
    })


@csrf_exempt
def campanha_carta_lead_view(request, token, jogada_pk):
    """AJAX/POST público — salva os dados de lead após virar a carta e ganhar."""
    if request.method != 'POST':
        return JsonResponse({'success': False}, status=405)

    campanha = get_object_or_404(Campanha, token=token, tipo='CARTA')
    jogada = get_object_or_404(CampanhaJogada, pk=jogada_pk, campanha=campanha)

    if jogada.lead_salvo:
        return JsonResponse({'success': True, 'ja_salvo': True})

    config = getattr(campanha, 'config_carta', None)

    jogada.nome     = request.POST.get('nome', '').strip() if (config and config.capturar_nome)     else ''
    jogada.cpf      = request.POST.get('cpf',  '').strip() if (config and config.capturar_cpf)      else ''
    jogada.telefone = request.POST.get('telefone', '').strip() if (config and config.capturar_telefone) else ''
    jogada.endereco = request.POST.get('endereco', '').strip() if (config and config.capturar_endereco) else ''

    # Validação de CPF
    if (config and config.capturar_cpf) and jogada.cpf:
        if not _validar_cpf(jogada.cpf):
            return JsonResponse({'success': False, 'error': 'CPF inválido. Verifique os dígitos e tente novamente.'}, status=400)

    jogada.lead_salvo = True
    jogada.save(update_fields=['nome', 'cpf', 'telefone', 'endereco', 'lead_salvo'])

    # Registrar como CampanhaLead para aparecer na aba "Leads"
    CampanhaLead.objects.create(
        campanha=campanha,
        nome=jogada.nome,
        cpf=jogada.cpf,
        telefone=jogada.telefone,
        endereco=jogada.endereco,
        codigo_cupom=jogada.premio.codigo_resgate if jogada.premio else '',
        ip=jogada.ip,
    )

    return JsonResponse({'success': True, 'ja_salvo': False})


@login_required
def treinamento_franqueado_view(request):
    """Guia de vendas e treinamento para franqueados"""
    return render(request, 'treinamento/guia_franqueado.html', {})


@login_required
def landing_leads_view(request):
    """Lista os leads capturados pela landing page pública (apenas OWNER)."""
    from .models import LandingLead
    if not request.user.is_owner:
        from django.http import HttpResponseForbidden
        return HttpResponseForbidden()

    import csv
    # Export CSV
    if request.GET.get('export') == 'csv':
        response = HttpResponse(content_type='text/csv; charset=utf-8-sig')
        response['Content-Disposition'] = 'attachment; filename="leads_landing.csv"'
        writer = csv.writer(response)
        writer.writerow(['#', 'Nome', 'WhatsApp', 'E-mail', 'Cidade', 'Segmento', 'Mensagem', 'Data'])
        for i, lead in enumerate(LandingLead.objects.all(), 1):
            writer.writerow([i, lead.nome, lead.whatsapp, lead.email, lead.cidade,
                             lead.get_segmento_display(), lead.mensagem,
                             lead.criado_em.strftime('%d/%m/%Y %H:%M')])
        return response

    leads = LandingLead.objects.all()
    total = leads.count()
    return render(request, 'landing/leads.html', {'leads': leads, 'total': total})


# ════════════════════════════════════════════════════════════════════════════
#  AGENTES DE IA
# ════════════════════════════════════════════════════════════════════════════

def _get_agente_or_403(request, pk):
    """Helper: retorna o agente garantindo que o user tem acesso."""
    agente = get_object_or_404(AgenteIA, pk=pk)
    user = request.user
    if not user.is_owner() and agente.franqueado != user:
        return None, redirect('agente_list')
    return agente, None


@login_required
def agente_list_view(request):
    user = request.user
    if user.is_owner():
        agentes = AgenteIA.objects.select_related('franqueado').all()
    else:
        agentes = AgenteIA.objects.filter(franqueado=user)
    return render(request, 'agentes/agente_list.html', {'agentes': agentes})


@login_required
def agente_create_view(request):
    def _initial_from_post(post):
        return {
            'nome': post.get('nome', ''),
            'nome_empresa': post.get('nome_empresa', ''),
            'descricao_curta': post.get('descricao_curta', ''),
            'modelo_ia': post.get('modelo_ia', 'gpt-4o-mini'),
            'temperatura': post.get('temperatura', '0.7'),
            'max_tokens': post.get('max_tokens', '500'),
            'prompt_sistema': post.get('prompt_sistema', ''),
            'restricoes': post.get('restricoes', ''),
            'mensagem_boas_vindas': post.get('mensagem_boas_vindas', 'Olá! Como posso te ajudar hoje? 😊'),
            'cor_primaria': post.get('cor_primaria', '#6366f1'),
            'placeholder_input': post.get('placeholder_input', 'Digite sua mensagem…'),
            'mensagem_coleta': post.get('mensagem_coleta', 'Para continuar, preencha seus dados:'),
            'limite_mensagens': post.get('limite_mensagens', '0'),
            'whatsapp_escalada': post.get('whatsapp_escalada', ''),
            'mensagem_escalada': post.get('mensagem_escalada', 'Preciso falar com um atendente'),
            'coleta_contato': 'coleta_contato' in post,
            'coleta_nome': 'coleta_nome' in post,
            'coleta_telefone': 'coleta_telefone' in post,
            'coleta_email': 'coleta_email' in post,
            'ativo': 'ativo' in post,
        }

    if request.method == 'POST':
        nome      = request.POST.get('nome', '').strip()
        modelo_ia = request.POST.get('modelo_ia', 'gpt-4o-mini')
        api_key   = request.POST.get('api_key', '').strip()

        if not nome or not api_key:
            messages.error(request, 'Nome e chave de API são obrigatórios.')
            franchisees = User.objects.filter(role='FRANCHISEE') if request.user.is_owner() else None
            return render(request, 'agentes/agente_configure.html', {
                'agente': None,
                'modelo_choices': AgenteIA.MODELO_CHOICES,
                'franchisees': franchisees,
                'initial': _initial_from_post(request.POST),
            })

        franqueado = request.user
        if request.user.is_owner():
            fk = request.POST.get('franqueado_id')
            if fk:
                franqueado = get_object_or_404(User, pk=fk, role='FRANCHISEE')

        agente = AgenteIA.objects.create(
            franqueado=franqueado,
            nome=nome,
            modelo_ia=modelo_ia,
            api_key=api_key,
            descricao_curta=request.POST.get('descricao_curta', '').strip(),
            nome_empresa=request.POST.get('nome_empresa', '').strip(),
            temperatura=float(request.POST.get('temperatura', 0.7)),
            max_tokens=int(request.POST.get('max_tokens', 500)),
            prompt_sistema=request.POST.get('prompt_sistema', '').strip(),
            restricoes=request.POST.get('restricoes', '').strip(),
            mensagem_boas_vindas=request.POST.get('mensagem_boas_vindas', 'Olá! Como posso te ajudar hoje? 😊').strip(),
            cor_primaria=request.POST.get('cor_primaria', '#6366f1').strip(),
            placeholder_input=request.POST.get('placeholder_input', 'Digite sua mensagem…').strip(),
            coleta_contato='coleta_contato' in request.POST,
            coleta_nome='coleta_nome' in request.POST,
            coleta_telefone='coleta_telefone' in request.POST,
            coleta_email='coleta_email' in request.POST,
            mensagem_coleta=request.POST.get('mensagem_coleta', '').strip(),
            limite_mensagens=int(request.POST.get('limite_mensagens', 0)),
            whatsapp_escalada=request.POST.get('whatsapp_escalada', '').strip(),
            mensagem_escalada=request.POST.get('mensagem_escalada', 'Preciso falar com um atendente').strip(),
            ativo='ativo' in request.POST,
        )
        if 'avatar' in request.FILES:
            agente.avatar = request.FILES['avatar']
            agente.save()
        if 'base_conhecimento' in request.FILES:
            f_base = request.FILES['base_conhecimento']
            _ext_ok = f_base.name.lower().endswith(('.txt', '.csv', '.tsv'))
            _size_ok = f_base.size <= 500 * 1024  # 500 KB
            if not _ext_ok or not _size_ok:
                messages.error(request, 'Base de conhecimento: use .txt/.csv/.tsv com até 500 KB.')
                return redirect('agente_configure', pk=agente.pk)
            agente.base_conhecimento = f_base
            agente.save(update_fields=['base_conhecimento'])

        messages.success(request, f'Agente "{agente.nome}" criado com sucesso!')
        return redirect('agente_configure', pk=agente.pk)

    franchisees = User.objects.filter(role='FRANCHISEE') if request.user.is_owner() else None
    return render(request, 'agentes/agente_configure.html', {
        'agente': None,
        'modelo_choices': AgenteIA.MODELO_CHOICES,
        'franchisees': franchisees,
        'initial': _initial_from_post({}),
    })


@login_required
def agente_configure_view(request, pk):
    agente, redir = _get_agente_or_403(request, pk)
    if redir:
        return redir

    if request.method == 'POST':
        agente.nome               = request.POST.get('nome', agente.nome).strip()
        agente.descricao_curta    = request.POST.get('descricao_curta', '').strip()
        agente.nome_empresa       = request.POST.get('nome_empresa', '').strip()
        agente.modelo_ia          = request.POST.get('modelo_ia', agente.modelo_ia)
        novo_key                  = request.POST.get('api_key', '').strip()
        if novo_key:
            agente.api_key        = novo_key
        agente.temperatura        = float(request.POST.get('temperatura', agente.temperatura))
        agente.max_tokens         = int(request.POST.get('max_tokens', agente.max_tokens))
        agente.prompt_sistema     = request.POST.get('prompt_sistema', '').strip()
        agente.restricoes         = request.POST.get('restricoes', '').strip()
        agente.mensagem_boas_vindas = request.POST.get('mensagem_boas_vindas', '').strip()
        agente.cor_primaria       = request.POST.get('cor_primaria', '#6366f1').strip()
        agente.placeholder_input  = request.POST.get('placeholder_input', '').strip()
        agente.coleta_contato     = 'coleta_contato' in request.POST
        agente.coleta_nome        = 'coleta_nome' in request.POST
        agente.coleta_telefone    = 'coleta_telefone' in request.POST
        agente.coleta_email       = 'coleta_email' in request.POST
        agente.mensagem_coleta    = request.POST.get('mensagem_coleta', '').strip()
        agente.limite_mensagens   = int(request.POST.get('limite_mensagens', 0))
        agente.whatsapp_escalada  = request.POST.get('whatsapp_escalada', '').strip()
        agente.mensagem_escalada  = request.POST.get('mensagem_escalada', '').strip()
        agente.habilitar_qualificacao = 'habilitar_qualificacao' in request.POST
        agente.ativo              = 'ativo' in request.POST
        if 'avatar' in request.FILES:
            agente.avatar = request.FILES['avatar']
        if 'base_conhecimento' in request.FILES:
            f_base = request.FILES['base_conhecimento']
            _ext_ok = f_base.name.lower().endswith(('.txt', '.csv', '.tsv'))
            _size_ok = f_base.size <= 500 * 1024  # 500 KB
            if not _ext_ok or not _size_ok:
                messages.error(request, 'Base de conhecimento: use .txt/.csv/.tsv com até 500 KB.')
                return redirect('agente_configure', pk=agente.pk)
            agente.base_conhecimento = f_base
        elif request.POST.get('limpar_base_conhecimento') == '1' and agente.base_conhecimento:
            agente.base_conhecimento.delete(save=False)
            agente.base_conhecimento = None
        agente.save()
        messages.success(request, 'Agente atualizado com sucesso!')
        return redirect('agente_configure', pk=agente.pk)

    return render(request, 'agentes/agente_configure.html', {
        'agente': agente,
        'modelo_choices': AgenteIA.MODELO_CHOICES,
        'embed_url': request.build_absolute_uri(
            reverse('agente_chat', kwargs={'public_id': agente.public_id}) + '?embed=1'
        ),
        'acoes': agente.acoes.order_by('ordem', 'nome'),
    })


@login_required
def agente_delete_view(request, pk):
    agente, redir = _get_agente_or_403(request, pk)
    if redir:
        return redir
    if request.method == 'POST':
        nome = agente.nome
        agente.delete()
        messages.success(request, f'Agente "{nome}" removido.')
        return redirect('agente_list')
    return render(request, 'agentes/agente_confirm_delete.html', {'agente': agente})


@login_required
def agente_historico_view(request, pk):
    agente, redir = _get_agente_or_403(request, pk)
    if redir:
        return redir
    conversas = agente.conversas.prefetch_related('mensagens', 'capturas__acao').order_by('-criado_em')
    return render(request, 'agentes/agente_historico.html', {
        'agente': agente,
        'conversas': conversas,
    })


# ── AgenteIAAcao CRUD ─────────────────────────────────────────────────────────

@login_required
def agente_acao_create_view(request, agente_pk):
    """Cria uma nova ação para o agente."""
    import json as _json
    agente, redir = _get_agente_or_403(request, agente_pk)
    if redir:
        return redir
    if request.method == 'POST':
        from .models import AgenteIAAcao
        acao = AgenteIAAcao(agente=agente)
        _tipos_validos  = ('http', 'capturar')
        _metodos_validos = ('GET', 'POST', 'PUT', 'DELETE')
        acao.tipo            = request.POST.get('tipo', 'http') if request.POST.get('tipo') in _tipos_validos else 'http'
        acao.nome            = request.POST.get('nome', '').strip()
        acao.descricao       = request.POST.get('descricao', '').strip()
        acao.url             = request.POST.get('url', '').strip()
        acao.metodo          = request.POST.get('metodo', 'GET') if request.POST.get('metodo') in _metodos_validos else 'GET'
        acao.headers_json    = request.POST.get('headers_json', '{}').strip() or '{}'
        acao.corpo_template  = request.POST.get('corpo_template', '').strip()
        acao.parametros_json = request.POST.get('parametros_json', '[]').strip() or '[]'
        acao.exibir_botao    = 'exibir_botao' in request.POST
        acao.texto_botao     = request.POST.get('texto_botao', '').strip()
        acao.ativo           = 'ativo' in request.POST
        try:
            acao.ordem = int(request.POST.get('ordem', 0))
        except ValueError:
            acao.ordem = 0
        # Valida JSON simples
        try:
            _json.loads(acao.headers_json)
            _json.loads(acao.parametros_json)
        except Exception as e:
            messages.error(request, f'JSON inválido: {e}')
            return redirect('agente_configure', pk=agente.pk)
        acao.save()
        messages.success(request, f'Ação "{acao.nome}" criada.')
    return redirect('agente_configure', pk=agente.pk)


@login_required
def agente_acao_edit_view(request, pk):
    """Edita uma ação existente."""
    import json as _json
    from .models import AgenteIAAcao
    acao = get_object_or_404(AgenteIAAcao, pk=pk)
    agente, redir = _get_agente_or_403(request, acao.agente_id)
    if redir:
        return redir
    if request.method == 'POST':
        _tipos_validos  = ('http', 'capturar')
        _metodos_validos = ('GET', 'POST', 'PUT', 'DELETE')
        acao.nome            = request.POST.get('nome', acao.nome).strip()
        acao.tipo            = request.POST.get('tipo', acao.tipo) if request.POST.get('tipo') in _tipos_validos else acao.tipo
        acao.descricao       = request.POST.get('descricao', '').strip()
        acao.url             = request.POST.get('url', '').strip()
        acao.metodo          = request.POST.get('metodo', acao.metodo) if request.POST.get('metodo') in _metodos_validos else acao.metodo
        acao.headers_json    = request.POST.get('headers_json', '{}').strip() or '{}'
        acao.corpo_template  = request.POST.get('corpo_template', '').strip()
        acao.parametros_json = request.POST.get('parametros_json', '[]').strip() or '[]'
        acao.exibir_botao    = 'exibir_botao' in request.POST
        acao.texto_botao     = request.POST.get('texto_botao', '').strip()
        acao.ativo           = 'ativo' in request.POST
        try:
            acao.ordem = int(request.POST.get('ordem', 0))
        except ValueError:
            acao.ordem = 0
        try:
            _json.loads(acao.headers_json)
            _json.loads(acao.parametros_json)
        except Exception as e:
            messages.error(request, f'JSON inválido: {e}')
            return redirect('agente_configure', pk=acao.agente_id)
        acao.save()
        messages.success(request, f'Ação "{acao.nome}" atualizada.')
    return redirect('agente_configure', pk=acao.agente_id)


@login_required
def agente_acao_delete_view(request, pk):
    """Remove uma ação."""
    from .models import AgenteIAAcao
    acao = get_object_or_404(AgenteIAAcao, pk=pk)
    agente_pk = acao.agente_id
    agente, redir = _get_agente_or_403(request, agente_pk)
    if redir:
        return redir
    if request.method == 'POST':
        nome = acao.nome
        acao.delete()
        messages.success(request, f'Ação "{nome}" removida.')
    return redirect('agente_configure', pk=agente_pk)


# ── Capturas (pedidos) ─────────────────────────────────────────────────────────

@login_required
def agente_capturas_view(request, pk):
    """Lista todos os pedidos/capturas de um agente."""
    from .models import AgenteIACaptura
    agente, redir = _get_agente_or_403(request, pk)
    if redir:
        return redir

    status_filtro = request.GET.get('status', '')
    capturas = agente.capturas.select_related('conversa', 'acao').order_by('-criado_em')
    if status_filtro:
        capturas = capturas.filter(status=status_filtro)

    return render(request, 'agentes/agente_capturas.html', {
        'agente': agente,
        'capturas': capturas,
        'status_filtro': status_filtro,
    })


@login_required
def agente_captura_status_view(request, pk):
    """POST — atualiza o status de uma captura."""
    from .models import AgenteIACaptura
    captura = get_object_or_404(AgenteIACaptura, pk=pk)
    agente, redir = _get_agente_or_403(request, captura.agente_id)
    if redir:
        return redir
    if request.method == 'POST':
        novo = request.POST.get('status', '')
        if novo in ('pendente', 'confirmado', 'cancelado'):
            captura.status = novo
            captura.save(update_fields=['status', 'atualizado_em'])
            messages.success(request, f'Status atualizado para "{captura.get_status_display()}".')
    # Redireciona de volta para onde veio (historico ou capturas)
    # Valida next_url contra open redirect — só permite URLs relativas do próprio site
    next_url = request.POST.get('next', '')
    if next_url and next_url.startswith('/'):
        return redirect(next_url)
    return redirect('agente_capturas', pk=captura.agente_id)


# ── Helpers ──────────────────────────────────────────────────────────────────

def _executar_captura(acao, argumentos: dict, conversa) -> str:
    """Salva dados capturados pela IA no banco e (opcionalmente) dispara webhook."""
    from .models import AgenteIACaptura
    import requests as _req

    captura = AgenteIACaptura.objects.create(
        conversa=conversa,
        agente=acao.agente,
        acao=acao,
        dados=argumentos,
        status='pendente',
    )

    # Webhook opcional
    if acao.url:
        try:
            headers = acao.headers()
            if 'Content-Type' not in headers:
                headers['Content-Type'] = 'application/json'
            resp = _req.post(acao.url, json=argumentos, headers=headers, timeout=10)
            captura.webhook_enviado = True
            captura.webhook_resposta = resp.text[:1000]
            captura.save(update_fields=['webhook_enviado', 'webhook_resposta'])
            return f'Pedido registrado com sucesso (#{captura.pk}). Webhook enviado ({resp.status_code}).'
        except Exception as e:
            captura.webhook_resposta = str(e)[:500]
            captura.save(update_fields=['webhook_resposta'])
            return f'Pedido registrado (#{captura.pk}), mas falha ao enviar webhook: {str(e)[:200]}'
    return f'Pedido registrado com sucesso (#{captura.pk}).'


def _executar_acao(acao, argumentos: dict, conversa=None) -> str:
    """Executa uma AgenteIAAcao e retorna texto com o resultado."""
    if acao.tipo == 'capturar':
        return _executar_captura(acao, argumentos, conversa)

    # tipo == 'http' — chamada HTTP externa
    import json as _json
    import requests as _req

    # Substitui placeholders na URL
    url = acao.url
    for k, v in argumentos.items():
        url = url.replace(f'{{{k}}}', str(v))

    headers = acao.headers()

    # Monta corpo se houver template
    corpo = None
    if acao.corpo_template:
        corpo_str = acao.corpo_template
        for k, v in argumentos.items():
            corpo_str = corpo_str.replace(f'{{{k}}}', str(v))
        try:
            corpo = _json.loads(corpo_str)
        except Exception:
            corpo = corpo_str  # envia como string se não for JSON válido

    try:
        resp = _req.request(
            acao.metodo, url, headers=headers,
            json=corpo if isinstance(corpo, dict) else None,
            data=corpo if isinstance(corpo, str) else None,
            timeout=10,
        )
        return resp.text[:3000]
    except Exception as e:
        return f'Erro ao executar ação: {str(e)}'


def _construir_tools_openai(acoes):
    """Constrói lista de tools no formato OpenAI a partir de AgenteIAAcao."""
    tools = []
    for acao in acoes:
        params = acao.parametros()
        properties = {}
        required = []
        for p in params:
            prop = {'type': p.get('tipo', 'string')}
            if p.get('descricao'):
                prop['description'] = p['descricao']
            properties[p['nome']] = prop
            if p.get('obrigatorio'):
                required.append(p['nome'])
        tools.append({
            'type': 'function',
            'function': {
                'name': acao.nome,
                'description': acao.descricao,
                'parameters': {
                    'type': 'object',
                    'properties': properties,
                    'required': required,
                },
            },
        })
    return tools


def _construir_tools_anthropic(acoes):
    """Constrói lista de tools no formato Anthropic a partir de AgenteIAAcao."""
    tools = []
    for acao in acoes:
        params = acao.parametros()
        properties = {}
        required = []
        for p in params:
            prop = {'type': p.get('tipo', 'string')}
            if p.get('descricao'):
                prop['description'] = p['descricao']
            properties[p['nome']] = prop
            if p.get('obrigatorio'):
                required.append(p['nome'])
        tools.append({
            'name': acao.nome,
            'description': acao.descricao,
            'input_schema': {
                'type': 'object',
                'properties': properties,
                'required': required,
            },
        })
    return tools


def _construir_tools_gemini(acoes):
    """Constrói FunctionDeclarations para Gemini a partir de AgenteIAAcao."""
    try:
        import google.generativeai as genai
        from google.generativeai.types import FunctionDeclaration, Tool
    except ImportError:
        return None

    declarations = []
    for acao in acoes:
        params = acao.parametros()
        properties = {}
        required = []
        for p in params:
            prop_schema = {'type': p.get('tipo', 'string').upper()}
            if p.get('descricao'):
                prop_schema['description'] = p['descricao']
            properties[p['nome']] = prop_schema
            if p.get('obrigatorio'):
                required.append(p['nome'])
        declarations.append(FunctionDeclaration(
            name=acao.nome,
            description=acao.descricao,
            parameters={'type': 'OBJECT', 'properties': properties, 'required': required},
        ))
    return [Tool(function_declarations=declarations)] if declarations else None


def _chamar_ia(agente, historico_msgs, conversa=None):
    """
    Chama o provedor de IA e retorna o texto da resposta.
    `historico_msgs` é lista de dicts [{role, content}, …].
    Suporta function calling via AgenteIAAcao quando o agente tem ações ativas.
    Raises: Exception com mensagem amigável em caso de erro.
    """
    import json as _json

    sistema_completo = agente.prompt_sistema
    if agente.base_conhecimento:
        try:
            with agente.base_conhecimento.open('r') as f:
                conteudo_base = f.read(60000)
            sistema_completo += f'\n\nBASE DE CONHECIMENTO (use estas informações para responder):\n{conteudo_base}'
        except Exception:
            pass
    if agente.restricoes:
        sistema_completo += f'\n\nRESTRIÇÕES IMPORTANTES:\n{agente.restricoes}'
    if agente.nome_empresa:
        sistema_completo = f'Você é um assistente da empresa "{agente.nome_empresa}".\n\n' + sistema_completo

    if agente.max_tokens <= 150:
        brevidade = 'IMPORTANTE: Responda em no máximo 1-2 frases curtas e sempre completas. Nunca corte uma frase no meio.'
    elif agente.max_tokens <= 300:
        brevidade = 'IMPORTANTE: Seja conciso. Responda em no máximo 3-4 frases e sempre termine com uma frase completa. Nunca deixe a resposta cortada no meio.'
    elif agente.max_tokens <= 600:
        brevidade = 'Seja direto e objetivo. Sempre termine suas respostas com uma frase completa.'
    else:
        brevidade = None
    if brevidade:
        sistema_completo += f'\n\n{brevidade}'

    # Injeta contexto de data/hora atual para que a IA saiba o dia exato
    import datetime as _dt
    _agora = _dt.datetime.now()
    _dias_pt = ['Segunda-feira', 'Terça-feira', 'Quarta-feira', 'Quinta-feira',
                'Sexta-feira', 'Sábado', 'Domingo']
    _meses_pt = ['janeiro', 'fevereiro', 'março', 'abril', 'maio', 'junho',
                 'julho', 'agosto', 'setembro', 'outubro', 'novembro', 'dezembro']
    _ctx_data = (
        f'CONTEXTO ATUAL: {_dias_pt[_agora.weekday()]}, '
        f'{_agora.day} de {_meses_pt[_agora.month - 1]} de {_agora.year}, '
        f'{_agora.strftime("%H:%M")}h.'
    )
    sistema_completo = _ctx_data + '\n\n' + sistema_completo

    # Injeta dados do visitante coletados no início da conversa
    if conversa:
        dados_visitante = []
        if getattr(conversa, 'nome_visitante', None):
            dados_visitante.append(f'Nome: {conversa.nome_visitante}')
        if getattr(conversa, 'telefone_visitante', None):
            dados_visitante.append(f'Telefone: {conversa.telefone_visitante}')
        if getattr(conversa, 'email_visitante', None):
            dados_visitante.append(f'E-mail: {conversa.email_visitante}')
        if dados_visitante:
            sistema_completo += (
                '\n\nDADOS DO CLIENTE NESTA CONVERSA (use para personalizar o atendimento):\n'
                + '\n'.join(dados_visitante)
            )

    # Injeta capturas já realizadas nesta conversa para evitar re-perguntar
    if conversa:
        try:
            from .models import AgenteIACaptura
            caps = AgenteIACaptura.objects.filter(conversa=conversa).order_by('criado_em')
            if caps.exists():
                import json as _json_caps
                linhas = []
                for cap in caps:
                    acao_nome = cap.acao.nome if cap.acao else 'captura'
                    dados_str = ', '.join(f'{k}: {v}' for k, v in cap.dados.items()) if cap.dados else '(sem dados)'
                    linhas.append(f'- #{cap.pk} [{acao_nome}] status={cap.get_status_display()} | {dados_str}')
                sistema_completo += (
                    '\n\nCAPTURAS JÁ REALIZADAS NESTA CONVERSA (NÃO PEÇA ESSES DADOS NOVAMENTE):\n'
                    + '\n'.join(linhas)
                )
        except Exception:
            pass

    # Carrega ações ativas (se o agente for real, não _FakeAgente)
    acoes_ativas = []
    if hasattr(agente, 'acoes'):
        try:
            acoes_ativas = list(agente.acoes.filter(ativo=True).order_by('ordem', 'nome'))
        except Exception:
            acoes_ativas = []

    provedor = agente.provedor

    # ── OpenAI ────────────────────────────────────────────────────────────────
    if provedor == 'openai':
        try:
            import openai as _oai
        except ImportError:
            raise Exception('Pacote openai não instalado no servidor.')
        client = _oai.OpenAI(api_key=agente.api_key)

        messages = [{'role': 'system', 'content': sistema_completo}] + historico_msgs

        if acoes_ativas:
            tools = _construir_tools_openai(acoes_ativas)
            acoes_map = {a.nome: a for a in acoes_ativas}

            # Loop agentico: executa até 5 rodadas de tool_calls
            for _ in range(5):
                resp = client.chat.completions.create(
                    model=agente.modelo_ia,
                    temperature=agente.temperatura,
                    max_tokens=agente.max_tokens,
                    messages=messages,
                    tools=tools,
                    tool_choice='auto',
                )
                choice = resp.choices[0]
                if choice.finish_reason == 'tool_calls' and choice.message.tool_calls:
                    # Adiciona mensagem do assistente com tool_calls
                    messages.append(choice.message)
                    for tc in choice.message.tool_calls:
                        try:
                            args = _json.loads(tc.function.arguments)
                        except Exception:
                            args = {}
                        acao = acoes_map.get(tc.function.name)
                        resultado = _executar_acao(acao, args, conversa) if acao else 'Ação não encontrada.'
                        messages.append({
                            'role': 'tool',
                            'tool_call_id': tc.id,
                            'content': resultado,
                        })
                else:
                    # Resposta final
                    return (choice.message.content or '').strip()
            # Fallback após 5 rodadas
            return (resp.choices[0].message.content or '').strip()
        else:
            resp = client.chat.completions.create(
                model=agente.modelo_ia,
                temperature=agente.temperatura,
                max_tokens=agente.max_tokens,
                messages=messages,
            )
            return resp.choices[0].message.content.strip()

    # ── Anthropic ─────────────────────────────────────────────────────────────
    elif provedor == 'anthropic':
        try:
            import anthropic as _ant
        except ImportError:
            raise Exception('Pacote anthropic não instalado no servidor.')
        client = _ant.Anthropic(api_key=agente.api_key)

        if acoes_ativas:
            tools = _construir_tools_anthropic(acoes_ativas)
            acoes_map = {a.nome: a for a in acoes_ativas}
            messages = list(historico_msgs)

            for _ in range(5):
                resp = client.messages.create(
                    model=agente.modelo_ia,
                    max_tokens=agente.max_tokens,
                    system=sistema_completo,
                    tools=tools,
                    messages=messages,
                )
                if resp.stop_reason == 'tool_use':
                    # Mensagem do assistente
                    messages.append({'role': 'assistant', 'content': resp.content})
                    tool_results = []
                    for blk in resp.content:
                        if blk.type == 'tool_use':
                            acao = acoes_map.get(blk.name)
                            resultado = _executar_acao(acao, blk.input, conversa) if acao else 'Ação não encontrada.'
                            tool_results.append({
                                'type': 'tool_result',
                                'tool_use_id': blk.id,
                                'content': resultado,
                            })
                    messages.append({'role': 'user', 'content': tool_results})
                else:
                    # Resposta final — extrai texto
                    for blk in resp.content:
                        if hasattr(blk, 'text'):
                            return blk.text.strip()
                    return ''
            return ''
        else:
            resp = client.messages.create(
                model=agente.modelo_ia,
                max_tokens=agente.max_tokens,
                system=sistema_completo,
                messages=historico_msgs,
            )
            return resp.content[0].text.strip()

    # ── Google Gemini ─────────────────────────────────────────────────────────
    elif provedor == 'google':
        try:
            import google.generativeai as genai
        except ImportError:
            raise Exception('Pacote google-generativeai não instalado no servidor.')
        genai.configure(api_key=agente.api_key)

        gemini_tools = _construir_tools_gemini(acoes_ativas) if acoes_ativas else None
        acoes_map = {a.nome: a for a in acoes_ativas}

        model = genai.GenerativeModel(
            model_name=agente.modelo_ia,
            system_instruction=sistema_completo,
            tools=gemini_tools,
        )
        history_gemini = [
            {'role': m['role'] if m['role'] != 'assistant' else 'model', 'parts': [m['content']]}
            for m in historico_msgs[:-1]
        ]
        chat = model.start_chat(history=history_gemini)

        if acoes_ativas and gemini_tools:
            resp = chat.send_message(historico_msgs[-1]['content'])
            for _ in range(5):
                # Verifica se há function_call
                part = resp.candidates[0].content.parts[0] if resp.candidates else None
                if part and hasattr(part, 'function_call') and part.function_call.name:
                    fc = part.function_call
                    acao = acoes_map.get(fc.name)
                    args = dict(fc.args) if fc.args else {}
                    resultado = _executar_acao(acao, args, conversa) if acao else 'Ação não encontrada.'
                    import google.generativeai.protos as _gp
                    resp = chat.send_message(_gp.Content(parts=[
                        _gp.Part(function_response=_gp.FunctionResponse(
                            name=fc.name,
                            response={'result': resultado},
                        ))
                    ]))
                else:
                    break
            return resp.text.strip()
        else:
            resp = chat.send_message(historico_msgs[-1]['content'])
            return resp.text.strip()

    raise Exception(f'Provedor desconhecido: {provedor}')


# ── Public chat views ─────────────────────────────────────────────────────────

def _chat_rate_limit(ip, key_suffix, limit, window_seconds):
    """Retorna True se o IP excedeu o limite. Usa Django cache."""
    from django.core.cache import cache
    key = f'rl:chat:{key_suffix}:{ip}'
    count = cache.get(key, 0)
    if count >= limit:
        return True
    cache.set(key, count + 1, window_seconds)
    return False


@xframe_options_exempt
def agente_chat_view(request, public_id):
    """Página pública de chat."""
    agente = get_object_or_404(AgenteIA, public_id=public_id, ativo=True)
    embed = request.GET.get('embed') == '1'
    chat_url = request.build_absolute_uri()
    embed_url = request.build_absolute_uri(
        reverse('agente_chat', kwargs={'public_id': agente.public_id}) + '?embed=1'
    )
    acoes_botoes = agente.acoes.filter(ativo=True, exibir_botao=True).order_by('ordem', 'nome')
    return render(request, 'agentes/agente_chat.html', {
        'agente': agente,
        'embed': embed,
        'embed_url': embed_url,
        'acoes_botoes': acoes_botoes,
    })


@csrf_exempt
def agente_chat_iniciar_view(request, public_id):
    """POST — cria a sessão e devolve session_id + boas-vindas."""
    if request.method != 'POST':
        return JsonResponse({'error': 'Method not allowed.'}, status=405)

    ip = request.META.get('HTTP_X_FORWARDED_FOR', request.META.get('REMOTE_ADDR', '')).split(',')[0].strip()

    # Rate limit: máx 10 sessões por IP por hora
    if _chat_rate_limit(ip, 'iniciar', 10, 3600):
        return JsonResponse({'error': 'Muitas tentativas. Tente novamente em alguns minutos.'}, status=429)

    agente = get_object_or_404(AgenteIA, public_id=public_id, ativo=True)

    import json as _json
    try:
        body = _json.loads(request.body)
    except Exception:
        body = {}

    conversa = AgenteIAConversa.objects.create(
        agente=agente,
        nome_visitante=body.get('nome', ''),
        telefone_visitante=body.get('telefone', ''),
        email_visitante=body.get('email', ''),
        ip=ip or None,
    )

    return JsonResponse({
        'ok': True,
        'session_id': str(conversa.session_id),
        'mensagem_boas_vindas': agente.mensagem_boas_vindas,
    })


@csrf_exempt
def agente_chat_enviar_view(request, public_id):
    """POST — envia mensagem do usuário, retorna resposta da IA."""
    if request.method != 'POST':
        return JsonResponse({'error': 'Method not allowed.'}, status=405)

    ip = request.META.get('HTTP_X_FORWARDED_FOR', request.META.get('REMOTE_ADDR', '')).split(',')[0].strip()

    # Rate limit: máx 40 mensagens por IP por minuto
    if _chat_rate_limit(ip, 'enviar', 40, 60):
        return JsonResponse({'error': 'Muitas mensagens. Aguarde um momento.'}, status=429)

    agente = get_object_or_404(AgenteIA, public_id=public_id, ativo=True)

    import json as _json
    try:
        body = _json.loads(request.body)
    except Exception:
        return JsonResponse({'error': 'JSON inválido.'}, status=400)

    session_id = body.get('session_id', '')
    mensagem   = body.get('mensagem', '').strip()

    if not session_id or not mensagem:
        return JsonResponse({'error': 'session_id e mensagem são obrigatórios.'}, status=400)

    # Limita tamanho da mensagem para evitar abuso de tokens da API
    if len(mensagem) > 4000:
        return JsonResponse({'error': 'Mensagem muito longa (máx. 4000 caracteres).'}, status=400)

    conversa = get_object_or_404(AgenteIAConversa, session_id=session_id, agente=agente)

    # Limite de mensagens por sessão
    if agente.limite_mensagens and conversa.total_mensagens >= agente.limite_mensagens:
        return JsonResponse({'error': 'Limite de mensagens atingido nesta sessão.', 'limite': True}, status=429)

    # Salva mensagem do usuário
    AgenteIAMensagem.objects.create(conversa=conversa, role='user', conteudo=mensagem)

    # Se um humano assumiu o atendimento, não chama a IA
    if conversa.modo == 'humano':
        return JsonResponse({'ok': True, 'aguardando': True})

    # Monta histórico para a IA (últimas 20 mensagens para não explodir tokens)
    msgs_db = conversa.mensagens.order_by('criado_em')
    historico = [{'role': m.role if m.role != 'assistant' else 'assistant', 'content': m.conteudo}
                 for m in msgs_db]
    # Para OpenAI/Anthropic "assistant" é válido; já para Gemini é tratado em _chamar_ia

    try:
        resposta_texto = _chamar_ia(agente, historico[-40:], conversa=conversa)  # janela de 40 msgs
    except Exception as e:
        import logging as _log
        _log.getLogger(__name__).error('Erro _chamar_ia agente=%s: %s', agente.pk, e)
        return JsonResponse({'error': 'Erro ao consultar a IA. Tente novamente em instantes.'}, status=500)

    # Salva resposta
    AgenteIAMensagem.objects.create(conversa=conversa, role='assistant', conteudo=resposta_texto)

    # Atualiza contador
    conversa.total_mensagens = conversa.mensagens.filter(role='user').count()
    conversa.save(update_fields=['total_mensagens'])

    # Classificação silenciosa de lead a cada 3 mensagens do usuário
    if agente.habilitar_qualificacao and conversa.total_mensagens % 3 == 0:
        try:
            score, motivo = _classificar_lead(agente, historico[-40:] + [{'role': 'assistant', 'content': resposta_texto}])
            conversa.qualificacao = score
            conversa.qualificacao_motivo = motivo
            conversa.qualificacao_em = timezone.now()
            conversa.save(update_fields=['qualificacao', 'qualificacao_motivo', 'qualificacao_em'])
        except Exception:
            pass  # nunca deixa classificação quebrar o chat

    return JsonResponse({'ok': True, 'resposta': resposta_texto})


def _classificar_lead(agente, historico_msgs):
    """
    Chama a IA em modo silencioso para classificar o lead.
    Retorna tupla (score, motivo) onde score in ['hot','warm','cold'].
    """
    import json as _json

    prompt_classificador = (
        "Você é um classificador de leads. Analise a conversa abaixo e responda SOMENTE com JSON válido, "
        "sem nenhum texto extra. Formato: {\"score\": \"hot\"|\"warm\"|\"cold\", \"motivo\": \"uma frase curta\"}. "
        "Critérios: hot = demonstrou intenção clara de compra/contratação, warm = interesse mas sem decisão, "
        "cold = apenas curiosidade ou sem engajamento real."
    )

    # Cria agente temporário apenas para a chamada — reutiliza _chamar_ia
    class _FakeAgente:
        prompt_sistema = prompt_classificador
        restricoes = ''
        nome_empresa = ''
        provedor = agente.provedor
        modelo_ia = agente.modelo_ia
        api_key = agente.api_key
        max_tokens = 120
        temperatura = 0.1
        base_conhecimento = None
        mensagem_boas_vindas = ''

    # Filtra apenas user/assistant para o classificador
    hist_limpo = [m for m in historico_msgs if m['role'] in ('user', 'assistant')]
    raw = _chamar_ia(_FakeAgente(), hist_limpo[-20:])

    # Extrai JSON da resposta (pode vir com markdown etc)
    import re as _re
    match = _re.search(r'\{.*?\}', raw, _re.DOTALL)
    if match:
        data = _json.loads(match.group())
        score = data.get('score', 'cold')
        if score not in ('hot', 'warm', 'cold'):
            score = 'cold'
        return score, data.get('motivo', '')[:300]
    return 'cold', ''


@csrf_exempt
def agente_chat_poll_view(request, public_id):
    """GET — retorna mensagens após `after_id` para suporte a human takeover."""
    agente = get_object_or_404(AgenteIA, public_id=public_id, ativo=True)
    session_id = request.GET.get('session_id', '')
    after_id   = request.GET.get('after', '0')
    try:
        after_id = int(after_id)
    except ValueError:
        after_id = 0

    conversa = get_object_or_404(AgenteIAConversa, session_id=session_id, agente=agente)
    msgs = conversa.mensagens.filter(id__gt=after_id).order_by('criado_em')

    return JsonResponse({
        'modo': conversa.modo,
        'messages': [
            {'id': m.id, 'role': m.role, 'conteudo': m.conteudo,
             'hora': m.criado_em.strftime('%H:%M')}
            for m in msgs
        ]
    })


@login_required
def agente_conversa_assumir_view(request, pk):
    """POST — atendente assume a conversa, IA para de responder."""
    from django.shortcuts import get_object_or_404
    conversa = get_object_or_404(AgenteIAConversa, pk=pk)
    # Verifica acesso: dono do agente ou owner
    if not (request.user.is_owner() or conversa.agente.franqueado == request.user):
        return JsonResponse({'error': 'Sem permissão.'}, status=403)
    conversa.modo = 'humano'
    conversa.assumido_por = request.user
    conversa.assumido_em = timezone.now()
    conversa.save(update_fields=['modo', 'assumido_por', 'assumido_em'])
    # Mensagem de sistema avisando o visitante
    AgenteIAMensagem.objects.create(
        conversa=conversa,
        role='human',
        conteudo=f'💬 {request.user.get_full_name() or request.user.username} assumiu o atendimento.',
    )
    messages.success(request, 'Você assumiu o atendimento desta conversa.')
    return redirect(request.META.get('HTTP_REFERER', 'agente_historico') )


@login_required
def agente_conversa_responder_view(request, pk):
    """POST — atendente envia mensagem manual para o visitante."""
    conversa = get_object_or_404(AgenteIAConversa, pk=pk)
    if not (request.user.is_owner() or conversa.agente.franqueado == request.user):
        return JsonResponse({'error': 'Sem permissão.'}, status=403)
    texto = request.POST.get('mensagem', '').strip()
    if texto:
        AgenteIAMensagem.objects.create(conversa=conversa, role='human', conteudo=texto)
    return redirect(request.META.get('HTTP_REFERER', 'agente_historico'))


@login_required
def agente_conversa_liberar_view(request, pk):
    """POST — devolve a conversa para a IA."""
    conversa = get_object_or_404(AgenteIAConversa, pk=pk)
    if not (request.user.is_owner() or conversa.agente.franqueado == request.user):
        return JsonResponse({'error': 'Sem permissão.'}, status=403)
    conversa.modo = 'ia'
    conversa.assumido_por = None
    conversa.assumido_em = None
    conversa.save(update_fields=['modo', 'assumido_por', 'assumido_em'])
    AgenteIAMensagem.objects.create(
        conversa=conversa,
        role='human',
        conteudo='🤖 Atendimento devolvido ao agente de IA.',
    )
    messages.success(request, 'Conversa devolvida para a IA.')
    return redirect(request.META.get('HTTP_REFERER', 'agente_historico'))
