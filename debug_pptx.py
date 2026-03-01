import zipfile, re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree

PPTX = 'ds.pptx'

# --- 1. Show raw media files and rId mapping ---
with zipfile.ZipFile(PPTX) as z:
    rels_xml = z.read('ppt/slides/_rels/slide1.xml.rels').decode('utf-8')
    rids = re.findall(r'Id="(rId\d+)"[^>]*Target="([^"]+)"', rels_xml)
    rid_map = {k: v for k, v in rids}
    print("=== rId -> media ===")
    for k, v in rid_map.items():
        print(f"  {k} -> {v}")

# --- 2. Inspect each shape for blipFill via lxml ---
prs = Presentation(PPTX)
slide = prs.slides[0]
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

print("\n=== Shapes with blipFill ===")
for shape in slide.shapes:
    el = shape.element
    blips = el.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
    if blips:
        for blip in blips:
            rid = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
            print(f"  shape={shape.name!r} type={shape.shape_type} rId={rid} -> {rid_map.get(rid, '?')}")
    # Check group children
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            el2 = child.element
            blips2 = el2.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
            if blips2:
                for blip in blips2:
                    rid = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                    print(f"    child={child.name!r} type={child.shape_type} rId={rid} -> {rid_map.get(rid, '?')}")
